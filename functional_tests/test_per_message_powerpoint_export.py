#!/usr/bin/env python3
# test_per_message_powerpoint_export.py
"""
Functional test for per-message PowerPoint export.
Version: 0.241.146
Implemented in: 0.241.033

This test ensures the message export flow exposes a PowerPoint route,
uses the frontend PowerPoint action hook, prefers the message model
deployment for AI slide planning, and produces a valid .pptx deck with
appendix slides for visuals, tables, code, and references. It also
ensures already structured markdown slide decks are exported without
AI summarization or slide-count compression. It also ensures generated
Markdown artifacts can be used as the PowerPoint export source and that
slide-local charts/images render as PNGs without leaking authoring labels.
It also verifies edited assistant markdown can be sent with direct message exports.
"""

import ast
import base64
import io
import json
import os
import re
import traceback
import zipfile
from html import escape as _escape_html
from pathlib import Path
from types import SimpleNamespace
from typing import Any, Dict, List, Optional, Tuple

import markdown2
from bs4 import BeautifulSoup, NavigableString, Tag
from PIL import Image
try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    PPTX_IMPORT_ERROR = None
except ModuleNotFoundError as exc:
    Presentation = None
    RGBColor = None
    MSO_AUTO_SHAPE_TYPE = None
    PP_ALIGN = None
    PptxInches = None
    PptxPt = None
    PPTX_IMPORT_ERROR = exc


REPO_ROOT = Path(__file__).resolve().parents[1]
ROUTE_FILE = REPO_ROOT / "application" / "single_app" / "route_backend_conversation_export.py"
FRONTEND_FILE = REPO_ROOT / "application" / "single_app" / "static" / "js" / "chat" / "chat-message-export.js"
MENU_FILE = REPO_ROOT / "application" / "single_app" / "static" / "js" / "chat" / "chat-messages.js"


def _normalize_content(content: Any) -> str:
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        parts = []
        for item in content:
            if isinstance(item, dict):
                if item.get('type') == 'text':
                    parts.append(item.get('text', ''))
                elif item.get('type') == 'image_url':
                    parts.append('[Image]')
                else:
                    parts.append(str(item))
            else:
                parts.append(str(item))
        return '\n'.join(parts)
    if isinstance(content, dict):
        if content.get('type') == 'text':
            return content.get('text', '')
        return str(content)
    return str(content) if content else ''


def _decode_base64_image_data_uri(data_uri: Optional[str]) -> Optional[bytes]:
    if not data_uri or not isinstance(data_uri, str):
        return None

    match = re.match(r'^data:image\/[a-zA-Z0-9.+-]+;base64,(.+)$', data_uri.strip())
    if not match:
        return None

    try:
        return base64.b64decode(match.group(1))
    except Exception:
        return None


def _build_test_image_data_uri(color: Tuple[int, int, int] = (37, 99, 235)) -> str:
    image = Image.new('RGB', (24, 16), color=color)
    buffer = io.BytesIO()
    image.save(buffer, format='PNG')
    encoded = base64.b64encode(buffer.getvalue()).decode('ascii')
    return f'data:image/png;base64,{encoded}'


def _collect_slide_titles(presentation: Presentation) -> List[str]:
    titles = []
    for slide in presentation.slides:
        title_shape = slide.shapes.title
        if title_shape and title_shape.text:
            titles.append(title_shape.text.strip())
    return titles


def _collect_slide_text(presentation: Presentation) -> str:
    text_parts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, 'has_text_frame', False):
                text_parts.append(shape.text)
            if getattr(shape, 'has_table', False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        text_parts.append(cell.text)
    return '\n'.join(text_parts)


def _powerpoint_dependencies_available(test_name: str) -> bool:
    if PPTX_IMPORT_ERROR is None:
        return True

    print(f"SKIP: {test_name} requires python-pptx: {PPTX_IMPORT_ERROR}")
    return False


def _load_powerpoint_helpers():
    if PPTX_IMPORT_ERROR is not None:
        raise RuntimeError(f'python-pptx is required for PowerPoint helper tests: {PPTX_IMPORT_ERROR}')

    helper_names = {
        '_message_to_pptx_bytes',
        '_attach_generated_image_proposal_assets',
        '_load_generated_image_proposal_assets',
        '_build_export_image_asset_from_message',
        '_resolve_image_message_export_data_uri',
        '_image_bytes_to_png_data_uri',
        '_render_message_export_content',
        '_get_message_export_image_assets',
        '_normalize_export_image_asset',
        '_replace_inline_image_proposal_blocks_with_export_html',
        '_parse_inline_image_proposal_payload',
        '_find_export_image_asset_for_proposal',
        '_build_export_inline_image_html',
        '_build_missing_export_inline_image_html',
        '_clean_export_visual_text',
        '_normalize_export_visual_id',
        '_normalize_export_prompt',
        '_sanitize_powerpoint_source_content',
        '_load_powerpoint_export_message_for_user',
        '_load_generated_markdown_artifact_for_user',
        '_parse_powerpoint_requested_slide_count',
        '_resolve_powerpoint_slide_count',
        '_build_structured_markdown_powerpoint_plan',
        '_extract_powerpoint_structured_slide_sections',
        '_extract_powerpoint_numbered_slide_sections',
        '_extract_powerpoint_separator_slide_sections',
        '_match_powerpoint_slide_marker',
        '_resolve_powerpoint_slide_title_and_content',
        '_extract_powerpoint_labeled_title',
        '_parse_powerpoint_structured_label_line',
        '_should_prefer_labeled_powerpoint_title',
        '_extract_powerpoint_preamble_title',
        '_is_powerpoint_slide_separator',
        '_looks_like_powerpoint_front_matter_block',
        '_is_powerpoint_title_section',
        '_build_powerpoint_title_metadata_from_section',
        '_extract_powerpoint_title_section_lines',
        '_line_contains_powerpoint_inline_visual',
        '_build_powerpoint_slide_footer_label',
        '_trim_powerpoint_structured_section_content',
        '_is_powerpoint_non_slide_tail_marker',
        '_extract_structured_powerpoint_bullets',
        '_extract_structured_powerpoint_tables',
        '_extract_structured_powerpoint_images',
        '_parse_powerpoint_markdown_table_block',
        '_split_powerpoint_markdown_table_line',
        '_is_powerpoint_markdown_table_separator_row',
        '_clean_powerpoint_markdown_table_cell',
        '_build_message_powerpoint_plan',
        '_build_fallback_powerpoint_plan',
        '_extract_message_powerpoint_model',
        '_normalize_powerpoint_model_candidate',
        '_generate_powerpoint_slide_plan_with_model',
        '_extract_json_object',
        '_sanitize_powerpoint_plan',
        '_extract_powerpoint_sections',
        '_extract_powerpoint_bullets',
        '_sentence_bullets',
        '_looks_like_markdown_table_row',
        '_looks_like_markdown_table_divider',
        '_markdown_to_plain_text',
        '_derive_powerpoint_title',
        '_build_powerpoint_subtitle',
        '_clean_slide_text',
        '_extract_powerpoint_appendix_assets',
        '_extract_powerpoint_images',
        '_extract_powerpoint_tables',
        '_extract_powerpoint_code_blocks',
        '_add_powerpoint_title_slide',
        '_add_powerpoint_content_slide',
        '_add_powerpoint_inline_image_to_slide',
        '_fit_powerpoint_image_within_bounds',
        '_add_powerpoint_inline_table',
        '_append_powerpoint_appendix_slides',
        '_add_powerpoint_image_slide',
        '_fit_powerpoint_image',
        '_add_powerpoint_table_slide',
        '_add_powerpoint_code_slide',
        '_apply_powerpoint_background',
        '_chunk_items',
    }

    source = ROUTE_FILE.read_text(encoding='utf-8')
    tree = ast.parse(source)
    selected_nodes = [
        node for node in tree.body
        if isinstance(node, ast.FunctionDef) and node.name in helper_names
    ]

    loaded_names = {node.name for node in selected_nodes}
    missing_names = helper_names - loaded_names
    assert not missing_names, f"Missing PowerPoint helpers in route file: {sorted(missing_names)}"

    requested_models: List[str] = []
    artifact_lookup_messages: Dict[str, Dict[str, Any]] = {}
    artifact_download_requests: List[Tuple[str, str]] = []

    class _FakeCompletions:
        def create(self, model, messages):
            slide_plan = {
                'presentation_title': 'Quarterly Review',
                'presentation_subtitle': 'Assistant | Generated from chat message',
                'slides': [
                    {
                        'title': 'Overview',
                        'bullets': [
                            'Revenue grew 18 percent.',
                            'Support backlog dropped 12 percent.',
                        ],
                    },
                    {
                        'title': 'Presenter Notes',
                        'bullets': [
                            'Use the visual, table, and code appendix as backup.',
                        ],
                    },
                ],
            }
            return SimpleNamespace(
                choices=[
                    SimpleNamespace(
                        message=SimpleNamespace(content=json.dumps(slide_plan))
                    )
                ]
            )

    class _FakeClient:
        def __init__(self):
            self.chat = SimpleNamespace(completions=_FakeCompletions())

    def _fake_initialize_gpt_client(settings, requested_model=''):
        requested_models.append(requested_model or '')
        return _FakeClient(), requested_model or 'fallback-model'

    def _fake_load_export_message_for_user(user_id, conversation_id, message_id):
        message = artifact_lookup_messages.get(message_id)
        if not message:
            raise LookupError('Message not found')
        return dict(message)

    def _fake_download_blob_content(container_name, blob_path):
        artifact_download_requests.append((container_name, blob_path))
        return b'# Artifact Deck\n\n## Slide 1 - Title\nArtifact title'

    module = ast.Module(body=selected_nodes, type_ignores=[])
    ast.fix_missing_locations(module)

    namespace = {
        'Any': Any,
        'base64': base64,
        'Dict': Dict,
        'List': List,
        'Optional': Optional,
        'Tuple': Tuple,
        'io': io,
        'json': json,
        're': re,
        '_escape_html': _escape_html,
        'markdown2': markdown2,
        'BeautifulSoup': BeautifulSoup,
        'NavigableString': NavigableString,
        'Tag': Tag,
        'Image': Image,
        'Presentation': Presentation,
        'RGBColor': RGBColor,
        'MSO_AUTO_SHAPE_TYPE': MSO_AUTO_SHAPE_TYPE,
        'PP_ALIGN': PP_ALIGN,
        'PptxInches': PptxInches,
        'PptxPt': PptxPt,
        'DOCX_MARKDOWN_EXTRAS': ['fenced-code-blocks', 'tables', 'break-on-newline', 'cuddled-lists', 'strike'],
        'POWERPOINT_PLAN_SOURCE_CHAR_LIMIT': 24000,
        'POWERPOINT_DEFAULT_SLIDES': 7,
        'POWERPOINT_MAX_SLIDES': 30,
        'POWERPOINT_MAX_STRUCTURED_SLIDES': 60,
        'POWERPOINT_MAX_BULLETS_PER_SLIDE': 5,
        'POWERPOINT_MAX_STRUCTURED_BULLETS_PER_SLIDE': 12,
        'POWERPOINT_BULLET_CHAR_LIMIT': 120,
        'POWERPOINT_STRUCTURED_BULLET_CHAR_LIMIT': 180,
        'POWERPOINT_MAX_APPENDIX_IMAGES': 4,
        'POWERPOINT_MAX_APPENDIX_TABLES': 3,
        'POWERPOINT_MAX_APPENDIX_CODE_BLOCKS': 2,
        'POWERPOINT_MAX_INLINE_IMAGES_PER_SLIDE': 2,
        'POWERPOINT_MAX_TABLE_ROWS': 8,
        'POWERPOINT_MAX_TABLE_COLS': 5,
        'POWERPOINT_TITLE_BG': RGBColor(22, 37, 66),
        'POWERPOINT_ACCENT': RGBColor(37, 99, 235),
        'POWERPOINT_BG': RGBColor(248, 250, 252),
        'POWERPOINT_PANEL': RGBColor(255, 255, 255),
        'POWERPOINT_TEXT': RGBColor(31, 41, 55),
        'POWERPOINT_MUTED': RGBColor(100, 116, 139),
        'POWERPOINT_TITLE_TEXT': RGBColor(255, 255, 255),
        'POWERPOINT_DATA_URI_PATTERN': re.compile(
            r"data:image\/[a-zA-Z0-9.+-]+;base64,[^\"'\s)]+",
            re.IGNORECASE,
        ),
        'INLINE_IMAGE_PROPOSAL_BLOCK_LANGUAGE': 'simpleimage',
        'INLINE_IMAGE_PROPOSAL_EXPORT_REGEX': re.compile(
            r"```simpleimage\s*([\s\S]*?)```",
            re.IGNORECASE,
        ),
        '_normalize_content': _normalize_content,
        '_load_export_message_for_user': _fake_load_export_message_for_user,
        'download_blob_content': _fake_download_blob_content,
        '_role_to_label': lambda role: {
            'assistant': 'Assistant',
            'user': 'User',
            'system': 'System',
        }.get(role, str(role).capitalize() or 'Message'),
        '_build_message_citation_labels': lambda message: [
            citation.get('title') or citation.get('label') or citation.get('url') or str(citation)
            for citation in message.get('citations', [])
        ],
        'replace_inline_chart_blocks_with_export_html': lambda content: content,
        'decode_base64_image_data_uri': _decode_base64_image_data_uri,
        'decode_image_content': lambda image_content: ('image/png', _decode_base64_image_data_uri(image_content) or b''),
        'get_complete_image_content': lambda *_args, **_kwargs: ({}, ''),
        'is_blob_backed_image_message': lambda *_args, **_kwargs: False,
        'is_external_image_url': lambda image_content: str(image_content or '').startswith(('http://', 'https://')),
        '_initialize_gpt_client': _fake_initialize_gpt_client,
        'debug_print': lambda *args, **kwargs: None,
        'log_event': lambda *args, **kwargs: None,
        '_artifact_lookup_messages': artifact_lookup_messages,
        '_artifact_download_requests': artifact_download_requests,
    }

    exec(compile(module, str(ROUTE_FILE), 'exec'), namespace)
    return namespace, requested_models


def test_export_powerpoint_route_definition_present() -> bool:
    """Route regression: the backend must define POST /api/message/export-powerpoint."""
    print("Testing backend route definition for PowerPoint export...")

    source = ROUTE_FILE.read_text(encoding='utf-8')
    tree = ast.parse(source)
    register_func = next(
        (
            node for node in tree.body
            if isinstance(node, ast.FunctionDef)
            and node.name == 'register_route_backend_conversation_export'
        ),
        None,
    )

    assert register_func is not None, 'register_route_backend_conversation_export should exist'

    route_found = False
    for node in register_func.body:
        if not isinstance(node, ast.FunctionDef):
            continue

        for decorator in node.decorator_list:
            if not isinstance(decorator, ast.Call):
                continue
            if not isinstance(decorator.func, ast.Attribute) or decorator.func.attr != 'route':
                continue
            if not decorator.args:
                continue

            route_arg = decorator.args[0]
            if not isinstance(route_arg, ast.Constant) or route_arg.value != '/api/message/export-powerpoint':
                continue

            methods_kw = next((keyword for keyword in decorator.keywords if keyword.arg == 'methods'), None)
            assert methods_kw is not None, 'PowerPoint export route should declare allowed methods'
            methods = [
                item.value for item in methods_kw.value.elts
                if isinstance(item, ast.Constant)
            ]
            assert 'POST' in methods, f'Expected POST method, found {methods}'
            assert node.name == 'api_export_message_powerpoint', f'Unexpected route handler name: {node.name}'
            route_found = True
            break

        if route_found:
            break

    assert route_found, 'Expected POST /api/message/export-powerpoint to be defined'
    print("PASS: backend route definition present")
    return True


def test_powerpoint_frontend_hooks_present() -> bool:
    """Frontend regression: menu and fetch path should expose PowerPoint export."""
    print("Testing frontend PowerPoint export hooks...")

    frontend_source = FRONTEND_FILE.read_text(encoding='utf-8')
    menu_source = MENU_FILE.read_text(encoding='utf-8')

    assert "fetch('/api/message/export-powerpoint'" in frontend_source, 'Expected frontend fetch for the PowerPoint export endpoint'
    assert 'exportMessageAsPowerPoint' in frontend_source, 'Expected frontend PowerPoint export helper'
    assert 'message_content_override' in frontend_source, 'Expected PowerPoint export requests to support edited assistant markdown overrides'
    assert 'delete requestBody.message_content_override;' in frontend_source, 'Expected artifact PowerPoint export to ignore visible-message overrides'
    assert 'dropdown-export-ppt-btn' in menu_source, 'Expected chat message menu PowerPoint action'
    assert 'artifact_message_id' in frontend_source, 'Expected PowerPoint export request to support artifact_message_id'
    assert 'generated-artifact-export-ppt-btn' in menu_source, 'Expected generated Markdown artifact PowerPoint action'

    print("PASS: frontend PowerPoint hooks present")
    return True


def test_powerpoint_export_can_use_generated_markdown_artifact_source() -> bool:
    """Generated Markdown artifact exports should load artifact blob content."""
    print("Testing generated Markdown artifact PowerPoint source loading...")

    if not _powerpoint_dependencies_available('generated Markdown artifact PowerPoint source loading'):
        return True

    helpers, _ = _load_powerpoint_helpers()
    lookup_messages = helpers['_artifact_lookup_messages']
    download_requests = helpers['_artifact_download_requests']
    lookup_messages.update({
        'assistant-message': {
            'id': 'assistant-message',
            'conversation_id': 'conversation-1',
            'role': 'assistant',
            'content': 'Short assistant preview only.',
            'timestamp': '2026-05-16T12:00:00Z',
            'metadata': {},
        },
        'artifact-message': {
            'id': 'artifact-message',
            'conversation_id': 'conversation-1',
            'role': 'file',
            'filename': 'generated-deck.md',
            'file_content_source': 'blob',
            'blob_container': 'chat-container',
            'blob_path': 'user/conversation/generated/generated-deck.md',
            'timestamp': '2026-05-16T12:01:00Z',
            'metadata': {
                'is_generated_chat_artifact': True,
                'generated_artifact_output_format': 'md',
            },
        },
    })

    export_message = helpers['_load_powerpoint_export_message_for_user'](
        user_id='user-1',
        conversation_id='conversation-1',
        message_id='assistant-message',
        artifact_message_id='artifact-message',
    )

    assert export_message['role'] == 'assistant', 'Artifact-backed exports should render as assistant content'
    assert export_message['content'].startswith('# Artifact Deck'), 'Expected artifact blob content as export source'
    assert export_message['timestamp'] == '2026-05-16T12:01:00Z', 'Expected artifact timestamp on export source'
    assert export_message['metadata']['powerpoint_export_source'] == 'generated_markdown_artifact'
    assert export_message['metadata']['powerpoint_export_artifact_message_id'] == 'artifact-message'
    assert download_requests == [('chat-container', 'user/conversation/generated/generated-deck.md')]

    print("PASS: PowerPoint export can use generated Markdown artifact source")
    return True


def test_powerpoint_export_prefers_message_model_and_renders_appendix() -> bool:
    """PowerPoint export should use the message model hint and render appendix slides."""
    print("Testing PowerPoint slide generation...")

    if not _powerpoint_dependencies_available('PowerPoint slide generation'):
        return True

    helpers, requested_models = _load_powerpoint_helpers()
    image_data_uri = _build_test_image_data_uri()
    message = {
        'role': 'assistant',
        'timestamp': '2026-05-04T12:00:00Z',
        'model_deployment_name': 'gpt-4o-mini',
        'content': '\n'.join([
            '# Quarterly Review',
            '',
            '- Revenue grew 18 percent.',
            '- Support backlog dropped 12 percent.',
            '',
            '| Metric | Value |',
            '| --- | --- |',
            '| Revenue | +18% |',
            '| Backlog | -12% |',
            '',
            '```python',
            'print("hello slides")',
            '```',
            '',
            f'<div class="export-inline-chart"><img src="{image_data_uri}" alt="Trend chart" /><div class="export-inline-chart-caption">Trend chart</div></div>',
        ]),
        'citations': [
            {'title': 'Quarterly workbook'},
            {'title': 'Operations dashboard'},
        ],
    }

    pptx_bytes = helpers['_message_to_pptx_bytes'](message, {'gpt_model': {'selected': [{'deploymentName': 'fallback-model'}]}})

    assert requested_models == ['gpt-4o-mini'], f'Expected the message deployment to be reused, found {requested_models}'
    assert pptx_bytes[:2] == b'PK', 'PowerPoint export should return a zipped OOXML payload'

    presentation = Presentation(io.BytesIO(pptx_bytes))
    slide_titles = _collect_slide_titles(presentation)

    assert len(presentation.slides) >= 6, f'Expected multiple slides including appendix content, found {len(presentation.slides)}'
    assert slide_titles[0] == 'Quarterly Review', f'Unexpected title slide text: {slide_titles[0]}'
    assert 'Overview' in slide_titles, f'Missing AI outline slide in {slide_titles}'
    assert 'Presenter Notes' in slide_titles, f'Missing second AI outline slide in {slide_titles}'
    assert 'Visual 1' in slide_titles, f'Missing visual appendix slide in {slide_titles}'
    assert 'Table 1' in slide_titles, f'Missing table appendix slide in {slide_titles}'
    assert 'Code Example 1' in slide_titles, f'Missing code appendix slide in {slide_titles}'
    assert 'References' in slide_titles, f'Missing references slide in {slide_titles}'

    print("PASS: PowerPoint export renders slide deck with appendix content")
    return True


def test_structured_markdown_powerpoint_export_preserves_slide_count() -> bool:
    """Structured markdown decks should bypass AI planning and preserve slide count."""
    print("Testing structured markdown PowerPoint export preservation...")

    if not _powerpoint_dependencies_available('structured markdown PowerPoint export preservation'):
        return True

    helpers, requested_models = _load_powerpoint_helpers()
    content_lines = [
        'Here is a polished PowerPoint-ready version of your draft.',
        '',
        '---',
        '',
        "# FAA FY 2026 President's Budget Submission",
        '**Source:** *FAA_FY_2026_Budget_Estimates_CJ.pdf*',
        '**Basis:** Consolidated review of provided pages **1-385**',
        '',
        '---',
        '',
        '## Slide 1 \u2014 Title',
        "**FAA FY 2026 President's Budget Submission**",
        '*FAA_FY_2026_Budget_Estimates_CJ.pdf*',
        'Consolidated review of pages **1-385**',
        '',
        '---',
        '',
    ]

    for index in range(2, 33):
        if index == 3:
            content_lines.extend([
                '## Slide 3 \u2014 Top-Line Budget by Account',
                '### FY 2026 Request',
                '',
                '| Account | FY 2026 Request |',
                '|---|---:|',
                '| Operations | **$13.842B** |',
                '| Facilities & Equipment | **$4.000B** |',
                '| Research, Engineering & Development | **$165.0M** |',
                '| Grants-in-Aid for Airports | **$4.000B** |',
                '',
                '### Total with IIJA / All Appropriations',
                '- **$27.005B**',
                '',
                '---',
                '',
            ])
            continue

        content_lines.extend([
            f'## Slide {index} \u2014 Topic {index}',
            '',
            f'- Key detail for slide {index}',
            f'- Supporting evidence for slide {index}',
            '',
            '---',
            '',
        ])

    content_lines.extend([
        '# Optional Speaker Notes Summary',
        'This should not become slide content.',
        '',
        '# Coverage Caveat',
        'This should also stay out of the exported deck.',
        '',
        'If you want, I can also turn this into a shorter executive briefing.',
    ])

    message = {
        'role': 'assistant',
        'timestamp': '2026-05-16T12:00:00Z',
        'model_deployment_name': 'gpt-4o-mini',
        'content': '\n'.join(content_lines),
        'citations': [
            {'title': 'Large source presentation'},
            {'title': 'Supporting notes'},
        ],
    }

    pptx_bytes = helpers['_message_to_pptx_bytes'](
        message,
        {'gpt_model': {'selected': [{'deploymentName': 'fallback-model'}]}},
    )

    assert requested_models == [], f'Structured markdown should not call AI planning, found {requested_models}'

    presentation = Presentation(io.BytesIO(pptx_bytes))
    slide_titles = _collect_slide_titles(presentation)
    slide_text = _collect_slide_text(presentation)

    assert len(presentation.slides) == 32, f'Expected exactly 32 structured slides, found {len(presentation.slides)}'
    assert slide_titles[0] == "FAA FY 2026 President's Budget Submission", f'Unexpected title slide text: {slide_titles[0]}'
    assert slide_titles[1] == 'Topic 2', f'Unexpected first content slide title: {slide_titles[1]}'
    assert slide_titles[2] == 'Top-Line Budget by Account', f'Unexpected table slide title: {slide_titles[2]}'
    assert slide_titles[-1] == 'Topic 32', f'Unexpected final slide title: {slide_titles[-1]}'
    assert any(shape.has_table for shape in presentation.slides[2].shapes), 'Expected native table on Slide 3'
    assert 'Operations' in slide_text and '$13.842B' in slide_text, 'Expected table content in exported deck'
    assert 'Slide 32' in slide_text, 'Expected authored slide number in footer text'
    assert 'Optional Speaker Notes Summary' not in slide_text, 'Speaker notes tail should not be exported'
    assert 'If you want' not in slide_text, 'Follow-up offer text should not be exported'

    print("PASS: structured markdown deck exported without slide-count compression")
    return True


def test_structured_powerpoint_export_embeds_slide_visuals_and_strips_labels() -> bool:
    """Structured slide exports should place PNG visuals on their source slides."""
    print("Testing structured PowerPoint visual placement and label cleanup...")

    if not _powerpoint_dependencies_available('structured PowerPoint visual placement and label cleanup'):
        return True

    helpers, requested_models = _load_powerpoint_helpers()
    generated_image_data_uri = _build_test_image_data_uri((37, 99, 235))
    chart_data_uri = _build_test_image_data_uri((215, 91, 53))
    image_proposal = {
        'version': 1,
        'visualId': 'colonial_map_1700',
        'title': 'Map of British North American Colonies',
        'description': 'A classroom map of British North America around 1700.',
        'prompt': 'Create a labeled classroom map of British North America around 1700.',
        'visualType': 'map',
        'slideNumber': 1,
        'context': 'Introductory overview of colonial North America.',
    }
    image_block = '```simpleimage\n' + json.dumps(image_proposal) + '\n```'
    chart_html = (
        '<div class="export-inline-chart">'
        f'<p><img src="{chart_data_uri}" alt="Regional trade chart" /></p>'
        '<p class="export-inline-chart-caption"><em>Regional trade comparison</em></p>'
        '</div>'
    )
    message = {
        'id': 'assistant-message',
        'role': 'assistant',
        'timestamp': '2026-06-04T12:00:00Z',
        'model_deployment_name': 'gpt-4o-mini',
        'content': '\n'.join([
            '## Slide 1: Introduction',
            'Title: What Was Early America Like in 1700?',
            'Bullet Points:',
            '- Thirteen British colonies stretched along the Atlantic coast',
            '- Indigenous nations and European empires shaped daily life',
            'Speaker Note:',
            'By 1700, the British colonies were growing quickly.',
            image_block,
            '',
            '---',
            '',
            '## Slide 2: The Thirteen Colonies',
            'Title: Regional Differences in the Colonies',
            'Bullet Points:',
            '- New England Colonies: small farms, fishing, shipbuilding, trade',
            '- Middle Colonies: diverse population, farming, trade, growing cities',
            chart_html,
        ]),
        '_export_generated_image_assets': [
            {
                'data_uri': generated_image_data_uri,
                'proposal': image_proposal,
                'title': image_proposal['title'],
                'caption': image_proposal['description'],
            }
        ],
        'citations': [],
    }

    pptx_bytes = helpers['_message_to_pptx_bytes'](
        message,
        {'gpt_model': {'selected': [{'deploymentName': 'fallback-model'}]}},
    )

    assert requested_models == [], f'Structured markdown should not call AI planning, found {requested_models}'
    presentation = Presentation(io.BytesIO(pptx_bytes))
    slide_titles = _collect_slide_titles(presentation)
    slide_text = _collect_slide_text(presentation)

    with zipfile.ZipFile(io.BytesIO(pptx_bytes), 'r') as archive:
        media_names = [name for name in archive.namelist() if name.startswith('ppt/media/')]

    assert len(presentation.slides) == 2, f'Expected two source slides, found {len(presentation.slides)}'
    assert slide_titles[0] == 'What Was Early America Like in 1700?', slide_titles
    assert slide_titles[1] == 'The Thirteen Colonies', slide_titles
    assert len(media_names) >= 2, f'Expected generated image and chart PNG media, found {media_names}'
    assert 'Title:' not in slide_text, slide_text
    assert 'Bullet Points:' not in slide_text, slide_text
    assert 'simpleimage' not in slide_text, slide_text
    assert 'visualId' not in slide_text, slide_text
    assert 'Regional Differences in the Colonies' not in slide_text, slide_text

    print("PASS: structured PowerPoint export embeds slide visuals and strips labels")
    return True


def test_powerpoint_slide_count_request_validation() -> bool:
    """Optional slide_count should accept bounded integers and reject invalid values."""
    print("Testing PowerPoint slide count request validation...")

    if not _powerpoint_dependencies_available('PowerPoint slide count request validation'):
        return True

    helpers, _ = _load_powerpoint_helpers()
    parse_slide_count = helpers['_parse_powerpoint_requested_slide_count']

    assert parse_slide_count(None) is None, 'Missing slide_count should keep default behavior'
    assert parse_slide_count('15') == 15, 'String slide_count should parse to an integer'
    assert parse_slide_count(30) == 30, 'Maximum supported slide_count should be accepted'

    invalid_values = ['0', '31', '2.5', True, {'slides': 15}]
    for invalid_value in invalid_values:
        try:
            parse_slide_count(invalid_value)
        except ValueError:
            continue
        raise AssertionError(f'Expected slide_count={invalid_value!r} to be rejected')

    print("PASS: PowerPoint slide count validation is bounded and explicit")
    return True


if __name__ == '__main__':
    tests = [
        test_export_powerpoint_route_definition_present,
        test_powerpoint_frontend_hooks_present,
        test_powerpoint_export_can_use_generated_markdown_artifact_source,
        test_powerpoint_export_prefers_message_model_and_renders_appendix,
        test_structured_markdown_powerpoint_export_preserves_slide_count,
        test_structured_powerpoint_export_embeds_slide_visuals_and_strips_labels,
        test_powerpoint_slide_count_request_validation,
    ]

    results = []
    for test in tests:
        print(f"\nRunning {test.__name__}...")
        try:
            results.append(bool(test()))
        except Exception as exc:
            print(f"FAIL: {test.__name__}: {exc}")
            traceback.print_exc()
            results.append(False)

    passed = sum(1 for result in results if result)
    print(f"\nResults: {passed}/{len(results)} tests passed")
    raise SystemExit(0 if all(results) else 1)