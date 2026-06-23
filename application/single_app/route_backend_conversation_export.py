# route_backend_conversation_export.py

import base64
import io
import json
import markdown2
import os
import re
import tempfile
import zipfile
from collections import Counter, defaultdict
from datetime import datetime
from html import escape as _escape_html
from typing import Any, Dict, List, Optional, Tuple

from bs4 import BeautifulSoup, NavigableString, Tag
from config import *
from flask import jsonify, make_response, request
from functions_appinsights import log_event
from functions_authentication import *
from functions_chat import sort_messages_by_thread
from functions_chart_export import (
    decode_base64_image_data_uri,
    replace_inline_chart_blocks_with_export_html,
)
from functions_collaboration import (
    assert_user_can_view_collaboration_conversation,
    get_accessible_collaboration_message_thoughts,
    get_collaboration_conversation,
    is_collaboration_conversation,
    list_collaboration_messages,
)
from functions_conversation_metadata import update_conversation_with_metadata
from functions_debug import debug_print
from functions_group import get_group_model_endpoints, get_user_groups
from functions_image_generation import INLINE_IMAGE_PROPOSAL_BLOCK_LANGUAGE
from functions_image_messages import (
    decode_image_content,
    get_complete_image_content,
    is_blob_backed_image_message,
    is_external_image_url,
)
from functions_message_artifacts import (
    build_message_artifact_payload_map,
    hydrate_agent_citations_from_artifacts,
    is_assistant_artifact_role,
)
from functions_settings import *
from functions_keyvault import SecretReturnType, keyvault_model_endpoint_get_helper
from functions_simplechat_operations import download_blob_content
from functions_thoughts import get_thoughts_for_conversation
from foundry_agent_runtime import resolve_authority
from model_endpoint_clients import (
    MODEL_ENDPOINT_PROTOCOL_ANTHROPIC,
    MODEL_ENDPOINT_PROTOCOL_AZURE_OPENAI,
    MODEL_ENDPOINT_PROTOCOL_OPENAI_STYLE,
    build_anthropic_chat_client,
    build_openai_style_chat_client,
    infer_model_endpoint_protocol,
)
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptxInches, Pt as PptxPt
from swagger_wrapper import swagger_route, get_auth_security
from docx import Document as DocxDocument
from docx.shared import Inches, Pt


TRANSCRIPT_ROLES = {'user', 'assistant'}
SUMMARY_SOURCE_CHAR_LIMIT = 60000
DOCX_MARKDOWN_EXTRAS = ['fenced-code-blocks', 'tables', 'break-on-newline', 'cuddled-lists', 'strike']
EMAIL_SUBJECT_CHAR_LIMIT = 120
EMAIL_SUBJECT_SOURCE_CHAR_LIMIT = 12000
EMAIL_CHART_ATTACHMENT_FILENAME_PREFIX = 'message_chart'
EMAIL_IMAGE_ATTACHMENT_FILENAME_PREFIX = 'message_image'
POWERPOINT_PLAN_SOURCE_CHAR_LIMIT = 24000
POWERPOINT_DEFAULT_SLIDES = 7
POWERPOINT_MAX_SLIDES = 30
POWERPOINT_MAX_STRUCTURED_SLIDES = 60
POWERPOINT_MAX_BULLETS_PER_SLIDE = 5
POWERPOINT_MAX_STRUCTURED_BULLETS_PER_SLIDE = 12
POWERPOINT_BULLET_CHAR_LIMIT = 120
POWERPOINT_STRUCTURED_BULLET_CHAR_LIMIT = 180
POWERPOINT_MAX_APPENDIX_IMAGES = 4
POWERPOINT_MAX_APPENDIX_TABLES = 3
POWERPOINT_MAX_APPENDIX_CODE_BLOCKS = 2
POWERPOINT_MAX_INLINE_IMAGES_PER_SLIDE = 2
POWERPOINT_MAX_TABLE_ROWS = 8
MESSAGE_EXPORT_CONTENT_OVERRIDE_MAX_LENGTH = 1000000
POWERPOINT_MAX_TABLE_COLS = 5

POWERPOINT_TITLE_BG = RGBColor(22, 37, 66)
POWERPOINT_ACCENT = RGBColor(37, 99, 235)
POWERPOINT_BG = RGBColor(248, 250, 252)
POWERPOINT_PANEL = RGBColor(255, 255, 255)
POWERPOINT_TEXT = RGBColor(31, 41, 55)
POWERPOINT_MUTED = RGBColor(100, 116, 139)
POWERPOINT_TITLE_TEXT = RGBColor(255, 255, 255)

POWERPOINT_DATA_URI_PATTERN = re.compile(
    r"data:image\/[a-zA-Z0-9.+-]+;base64,[^\"'\s)]+",
    re.IGNORECASE,
)
INLINE_IMAGE_PROPOSAL_EXPORT_REGEX = re.compile(
    rf"```{re.escape(INLINE_IMAGE_PROPOSAL_BLOCK_LANGUAGE)}\s*([\s\S]*?)```",
    re.IGNORECASE,
)


def register_route_backend_conversation_export(app):
    """Register conversation export API routes."""

    @app.route('/api/conversations/export', methods=['POST'])
    @swagger_route(security=get_auth_security())
    @login_required
    @user_required
    def api_export_conversations():
        """
        Export one or more conversations in JSON or Markdown format.
        Supports single-file or ZIP packaging.

        Request body:
            conversation_ids (list): List of conversation IDs to export.
            format (str): Export format — "json" or "markdown".
            packaging (str): Output packaging — "single" or "zip".
            include_summary_intro (bool): Whether to generate a per-conversation intro.
            summary_model_deployment (str): Optional model deployment for summary generation.
            summary_model_endpoint_id (str): Optional configured endpoint id for summary generation.
            summary_model_id (str): Optional configured endpoint model id for summary generation.
            summary_model_provider (str): Optional configured endpoint provider for summary generation.
        """
        user_id = get_current_user_id()
        if not user_id:
            return jsonify({'error': 'User not authenticated'}), 401

        data = request.get_json(silent=True)
        if not data:
            return jsonify({'error': 'Request body is required'}), 400

        conversation_ids = data.get('conversation_ids', [])
        export_format = str(data.get('format', 'json')).lower()
        packaging = str(data.get('packaging', 'single')).lower()
        include_summary_intro = bool(data.get('include_summary_intro', False))
        summary_model_deployment = str(data.get('summary_model_deployment', '') or '').strip()
        summary_model_endpoint_id = str(data.get('summary_model_endpoint_id', '') or '').strip()
        summary_model_id = str(data.get('summary_model_id', '') or '').strip()
        summary_model_provider = str(data.get('summary_model_provider', '') or '').strip()

        if not conversation_ids or not isinstance(conversation_ids, list):
            return jsonify({'error': 'At least one conversation_id is required'}), 400

        if export_format not in ('json', 'markdown', 'pdf'):
            return jsonify({'error': 'Format must be "json", "markdown", or "pdf"'}), 400

        if packaging not in ('single', 'zip'):
            return jsonify({'error': 'Packaging must be "single" or "zip"'}), 400

        try:
            settings = get_settings()
            exported = []
            for conv_id in conversation_ids:
                conversation = None
                messages = []
                try:
                    conversation = cosmos_conversations_container.read_item(
                        item=conv_id,
                        partition_key=conv_id
                    )
                    if conversation.get('user_id') != user_id:
                        debug_print(f"Export: user {user_id} does not own conversation {conv_id}")
                        continue

                    message_query = """
                        SELECT * FROM c
                        WHERE c.conversation_id = @conversation_id
                        ORDER BY c.timestamp ASC
                    """
                    messages = list(cosmos_messages_container.query_items(
                        query=message_query,
                        parameters=[{'name': '@conversation_id', 'value': conv_id}],
                        partition_key=conv_id
                    ))
                except Exception:
                    try:
                        conversation = get_collaboration_conversation(conv_id)
                        access_context = assert_user_can_view_collaboration_conversation(
                            user_id,
                            conversation,
                            allow_pending=True,
                        )
                        user_state = access_context.get('user_state') or {}
                        conversation = dict(conversation)
                        conversation['is_pinned'] = bool(user_state.get('is_pinned', False))
                        conversation['is_hidden'] = bool(user_state.get('is_hidden', False))
                        messages = list_collaboration_messages(conv_id)
                    except Exception:
                        debug_print(f"Export: conversation {conv_id} not found or access denied")
                        continue

                exported.append(
                    _build_export_entry(
                        conversation=conversation,
                        raw_messages=messages,
                        user_id=user_id,
                        settings=settings,
                        include_summary_intro=include_summary_intro,
                        summary_model_deployment=summary_model_deployment,
                        summary_model_endpoint_id=summary_model_endpoint_id,
                        summary_model_id=summary_model_id,
                        summary_model_provider=summary_model_provider,
                    )
                )

            if not exported:
                return jsonify({'error': 'No accessible conversations found'}), 404

            timestamp_str = datetime.utcnow().strftime('%Y%m%d_%H%M%S')

            if packaging == 'zip':
                return _build_zip_response(exported, export_format, timestamp_str)

            return _build_single_file_response(exported, export_format, timestamp_str)

        except Exception as exc:
            debug_print(f"Export error: {str(exc)}")
            log_event(f"Conversation export failed: {exc}", level="WARNING")
            return jsonify({'error': f'Export failed: {str(exc)}'}), 500

    @app.route('/api/message/export-word', methods=['POST'])
    @swagger_route(security=get_auth_security())
    @login_required
    @user_required
    def api_export_message_word():
        """
        Export a single message as a Word (.docx) document.

        Request body:
            message_id (str): ID of the message to export.
            conversation_id (str): ID of the conversation the message belongs to.
        """
        user_id = get_current_user_id()
        if not user_id:
            return jsonify({'error': 'User not authenticated'}), 401

        data = request.get_json(silent=True)
        if not data:
            return jsonify({'error': 'Request body is required'}), 400

        message_id = str(data.get('message_id', '') or '').strip()
        conversation_id = str(data.get('conversation_id', '') or '').strip()

        if not message_id or not conversation_id:
            return jsonify({'error': 'message_id and conversation_id are required'}), 400

        try:
            message = _load_export_message_for_user(
                user_id=user_id,
                conversation_id=conversation_id,
                message_id=message_id
            )
            message_content_override = _get_message_export_content_override(data)
            message = _apply_message_export_content_override(message, message_content_override)

            document_bytes = _message_to_docx_bytes(message)
            timestamp_str = datetime.utcnow().strftime('%Y%m%d_%H%M%S')
            filename = f"message_export_{timestamp_str}.docx"

            response = make_response(document_bytes)
            response.headers['Content-Type'] = (
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            )
            response.headers['Content-Disposition'] = f'attachment; filename="{filename}"'
            return response

        except LookupError as exc:
            return jsonify({'error': str(exc)}), 404
        except PermissionError as exc:
            return jsonify({'error': str(exc)}), 403
        except ValueError as exc:
            return jsonify({'error': str(exc)}), 400

        except Exception as exc:
            debug_print(f"Message export error: {str(exc)}")
            log_event(f"Message export failed: {exc}", level="WARNING")
            return jsonify({'error': 'Export failed due to a server error. Please try again later.'}), 500

    @app.route('/api/message/export-powerpoint', methods=['POST'])
    @swagger_route(security=get_auth_security())
    @login_required
    @user_required
    def api_export_message_powerpoint():
        """
        Export a single message as a PowerPoint (.pptx) presentation.

        Request body:
            message_id (str): ID of the message to export.
            conversation_id (str): ID of the conversation the message belongs to.
        """
        user_id = get_current_user_id()
        if not user_id:
            return jsonify({'error': 'User not authenticated'}), 401

        data = request.get_json(silent=True)
        if not data:
            return jsonify({'error': 'Request body is required'}), 400

        message_id = str(data.get('message_id', '') or '').strip()
        conversation_id = str(data.get('conversation_id', '') or '').strip()
        artifact_message_id = str(data.get('artifact_message_id', '') or '').strip()
        slide_count_value = data.get('slide_count')
        if slide_count_value in (None, ''):
            slide_count_value = data.get('target_slide_count')

        if not message_id or not conversation_id:
            return jsonify({'error': 'message_id and conversation_id are required'}), 400

        try:
            requested_slide_count = _parse_powerpoint_requested_slide_count(slide_count_value)
        except ValueError as exc:
            return jsonify({'error': str(exc)}), 400

        try:
            settings = get_settings()
            message = _load_powerpoint_export_message_for_user(
                user_id=user_id,
                conversation_id=conversation_id,
                message_id=message_id,
                artifact_message_id=artifact_message_id,
            )
            if not artifact_message_id:
                message_content_override = _get_message_export_content_override(data)
                message = _apply_message_export_content_override(message, message_content_override)

            presentation_bytes = _message_to_pptx_bytes(
                message,
                settings,
                requested_slide_count=requested_slide_count,
            )
            timestamp_str = datetime.utcnow().strftime('%Y%m%d_%H%M%S')
            filename = f"message_export_{timestamp_str}.pptx"

            response = make_response(presentation_bytes)
            response.headers['Content-Type'] = (
                'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
            response.headers['Content-Disposition'] = f'attachment; filename="{filename}"'
            return response

        except LookupError as exc:
            return jsonify({'error': str(exc)}), 404
        except PermissionError as exc:
            return jsonify({'error': str(exc)}), 403
        except ValueError as exc:
            return jsonify({'error': str(exc)}), 400
        except Exception as exc:
            debug_print(f"Message PowerPoint export error: {str(exc)}")
            log_event(f"Message PowerPoint export failed: {exc}", level="WARNING")
            return jsonify({'error': 'PowerPoint export failed due to a server error. Please try again later.'}), 500

    @app.route('/api/message/export-email-draft', methods=['POST'])
    @swagger_route(security=get_auth_security())
    @login_required
    @user_required
    def api_export_message_email_draft():
        """
        Build a mailto-ready email draft for a single message.

        Request body:
            message_id (str): ID of the message to export.
            conversation_id (str): ID of the conversation the message belongs to.
            summary_model_deployment (str): Optional model deployment for subject generation.
        """
        user_id = get_current_user_id()
        if not user_id:
            return jsonify({'error': 'User not authenticated'}), 401

        data = request.get_json(silent=True)
        if not data:
            return jsonify({'error': 'Request body is required'}), 400

        message_id = str(data.get('message_id', '') or '').strip()
        conversation_id = str(data.get('conversation_id', '') or '').strip()
        summary_model_deployment = str(data.get('summary_model_deployment', '') or '').strip()

        if not message_id or not conversation_id:
            return jsonify({'error': 'message_id and conversation_id are required'}), 400

        try:
            settings = get_settings()
            message = _load_export_message_for_user(
                user_id=user_id,
                conversation_id=conversation_id,
                message_id=message_id
            )
            message_content_override = _get_message_export_content_override(data)
            message = _apply_message_export_content_override(message, message_content_override)
            draft_payload = _message_to_email_draft_payload(
                message=message,
                settings=settings,
                summary_model_deployment=summary_model_deployment
            )
            return jsonify(draft_payload), 200

        except LookupError as exc:
            return jsonify({'error': str(exc)}), 404
        except PermissionError as exc:
            return jsonify({'error': str(exc)}), 403
        except ValueError as exc:
            return jsonify({'error': str(exc)}), 400

        except Exception as exc:
            debug_print(f"Message email draft export error: {str(exc)}")
            log_event(f"Message email draft export failed: {exc}", level="WARNING")
            return jsonify({'error': 'Email draft export failed due to a server error. Please try again later.'}), 500


def _build_export_entry(
    conversation: Dict[str, Any],
    raw_messages: List[Dict[str, Any]],
    user_id: str,
    settings: Dict[str, Any],
    include_summary_intro: bool = False,
    summary_model_deployment: str = '',
    summary_model_endpoint_id: str = '',
    summary_model_id: str = '',
    summary_model_provider: str = ''
) -> Dict[str, Any]:
    artifact_payload_map = build_message_artifact_payload_map(raw_messages)
    filtered_messages = _filter_messages_for_export(raw_messages)
    filtered_messages = hydrate_agent_citations_from_artifacts(filtered_messages, artifact_payload_map)
    ordered_messages = sort_messages_by_thread(filtered_messages)

    raw_thoughts = [] if is_collaboration_conversation(conversation) else get_thoughts_for_conversation(conversation.get('id'), user_id)
    thoughts_by_message = defaultdict(list)
    for thought in raw_thoughts:
        thoughts_by_message[thought.get('message_id')].append(_sanitize_thought(thought))

    exported_messages = []
    role_counts = Counter()
    total_citation_counts = Counter({'document': 0, 'web': 0, 'agent_tool': 0, 'legacy': 0, 'total': 0})
    transcript_index = 0
    total_thoughts = 0

    for sequence_index, message in enumerate(ordered_messages, start=1):
        role = message.get('role', 'unknown')
        role_counts[role] += 1

        message_transcript_index = None
        if role in TRANSCRIPT_ROLES:
            transcript_index += 1
            message_transcript_index = transcript_index

        thoughts = thoughts_by_message.get(message.get('id'), [])
        if not thoughts and is_collaboration_conversation(conversation):
            collaboration_thoughts = get_accessible_collaboration_message_thoughts(
                conversation,
                message,
                user_id,
            )
            thoughts = [_sanitize_thought(thought) for thought in collaboration_thoughts]
        exported_message = _sanitize_message(
            message,
            sequence_index=sequence_index,
            transcript_index=message_transcript_index,
            thoughts=thoughts
        )
        exported_messages.append(exported_message)

        counts = exported_message.get('citation_counts', {})
        for key in total_citation_counts:
            total_citation_counts[key] += counts.get(key, 0)
        total_thoughts += len(thoughts)

    # Compute message time range for summary caching
    message_time_start = None
    message_time_end = None
    if ordered_messages:
        message_time_start = ordered_messages[0].get('timestamp')
        message_time_end = ordered_messages[-1].get('timestamp')

    sanitized_conversation = _sanitize_conversation(
        conversation,
        messages=exported_messages,
        role_counts=role_counts,
        citation_counts=total_citation_counts,
        thought_count=total_thoughts
    )
    summary_intro = _build_summary_intro(
        messages=exported_messages,
        conversation=conversation,
        sanitized_conversation=sanitized_conversation,
        settings=settings,
        enabled=include_summary_intro,
        summary_model_deployment=summary_model_deployment,
        summary_model_endpoint_id=summary_model_endpoint_id,
        summary_model_id=summary_model_id,
        summary_model_provider=summary_model_provider,
        user_id=user_id,
        message_time_start=message_time_start,
        message_time_end=message_time_end
    )

    return {
        'conversation': sanitized_conversation,
        'summary_intro': summary_intro,
        'messages': exported_messages
    }


def _filter_messages_for_export(messages: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    filtered_messages = []
    for message in messages:
        if is_assistant_artifact_role(message.get('role')):
            continue

        metadata = message.get('metadata', {}) or {}
        if metadata.get('is_deleted') is True:
            continue

        thread_info = metadata.get('thread_info', {}) or {}
        active = thread_info.get('active_thread')
        if active is True or active is None or 'active_thread' not in thread_info:
            filtered_messages.append(message)

    return filtered_messages


def _sanitize_conversation(
    conversation: Dict[str, Any],
    messages: List[Dict[str, Any]],
    role_counts: Counter,
    citation_counts: Counter,
    thought_count: int
) -> Dict[str, Any]:
    transcript_count = sum(1 for message in messages if message.get('is_transcript_message'))
    return {
        'id': conversation.get('id'),
        'title': conversation.get('title', 'Untitled'),
        'last_updated': conversation.get('last_updated') or conversation.get('updated_at', ''),
        'chat_type': conversation.get('chat_type', 'personal'),
        'tags': conversation.get('tags', []),
        'context': conversation.get('context', []),
        'classification': conversation.get('classification', []),
        'strict': conversation.get('strict', False),
        'is_pinned': conversation.get('is_pinned', False),
        'scope_locked': conversation.get('scope_locked'),
        'locked_contexts': conversation.get('locked_contexts', []),
        'message_count': len(messages),
        'transcript_message_count': transcript_count,
        'message_counts_by_role': dict(role_counts),
        'citation_counts': dict(citation_counts),
        'thought_count': thought_count
    }


def _sanitize_message(
    message: Dict[str, Any],
    sequence_index: int,
    transcript_index: Optional[int],
    thoughts: List[Dict[str, Any]]
) -> Dict[str, Any]:
    role = message.get('role', '')
    content = message.get('content', '')
    raw_citation_buckets = _collect_raw_citation_buckets(message)
    normalized_citations = _normalize_citations(raw_citation_buckets)
    citation_counts = _build_citation_counts(normalized_citations)
    details = _curate_message_details(message, citation_counts, len(thoughts))

    return {
        'id': message.get('id'),
        'role': role,
        'speaker_label': _role_to_label(role),
        'sequence_index': sequence_index,
        'transcript_index': transcript_index,
        'label': f"Turn {transcript_index}" if transcript_index else f"Message {sequence_index}",
        'is_transcript_message': role in TRANSCRIPT_ROLES,
        'timestamp': message.get('timestamp', ''),
        'content': content,
        'content_text': _normalize_content(content),
        'details': details,
        'citations': normalized_citations,
        'citation_counts': citation_counts,
        'thoughts': thoughts,
        'legacy_citations': raw_citation_buckets['legacy'],
        'hybrid_citations': raw_citation_buckets['hybrid'],
        'web_search_citations': raw_citation_buckets['web'],
        'agent_citations': raw_citation_buckets['agent']
    }


def _sanitize_thought(thought: Dict[str, Any]) -> Dict[str, Any]:
    return {
        'step_index': thought.get('step_index'),
        'step_type': thought.get('step_type'),
        'content': thought.get('content'),
        'detail': thought.get('detail'),
        'duration_ms': thought.get('duration_ms'),
        'timestamp': thought.get('timestamp')
    }


def _collect_raw_citation_buckets(message: Dict[str, Any]) -> Dict[str, List[Any]]:
    def ensure_list(value: Any) -> List[Any]:
        if not value:
            return []
        return value if isinstance(value, list) else [value]

    return {
        'legacy': ensure_list(message.get('citations')),
        'hybrid': ensure_list(message.get('hybrid_citations')),
        'web': ensure_list(message.get('web_search_citations')),
        'agent': ensure_list(message.get('agent_citations'))
    }


def _normalize_citations(raw_citation_buckets: Dict[str, List[Any]]) -> List[Dict[str, Any]]:
    normalized = []

    for citation in raw_citation_buckets.get('hybrid', []):
        if isinstance(citation, dict):
            normalized.append({
                'citation_type': 'document',
                'label': _build_document_citation_label(citation),
                'file_name': citation.get('file_name'),
                'title': citation.get('title') or citation.get('file_name'),
                'page_number': citation.get('page_number'),
                'citation_id': citation.get('citation_id'),
                'chunk_id': citation.get('chunk_id'),
                'metadata_type': citation.get('metadata_type'),
                'metadata_content': citation.get('metadata_content'),
                'score': citation.get('score'),
                'classification': citation.get('classification'),
                'url': citation.get('url')
            })
        else:
            normalized.append({
                'citation_type': 'document',
                'label': str(citation),
                'value': str(citation)
            })

    for citation in raw_citation_buckets.get('web', []):
        if isinstance(citation, dict):
            title = citation.get('title') or citation.get('url') or 'Web source'
            normalized.append({
                'citation_type': 'web',
                'label': title,
                'title': title,
                'url': citation.get('url')
            })
        else:
            normalized.append({
                'citation_type': 'web',
                'label': str(citation),
                'value': str(citation)
            })

    for citation in raw_citation_buckets.get('agent', []):
        if isinstance(citation, dict):
            tool_name = citation.get('tool_name') or citation.get('function_name') or 'Tool invocation'
            normalized.append({
                'citation_type': 'agent_tool',
                'label': tool_name,
                'tool_name': citation.get('tool_name'),
                'function_name': citation.get('function_name'),
                'plugin_name': citation.get('plugin_name'),
                'success': citation.get('success'),
                'timestamp': citation.get('timestamp')
            })
        else:
            normalized.append({
                'citation_type': 'agent_tool',
                'label': str(citation),
                'value': str(citation)
            })

    for citation in raw_citation_buckets.get('legacy', []):
        if isinstance(citation, dict):
            title = citation.get('title') or citation.get('filepath') or citation.get('url') or 'Legacy citation'
            normalized.append({
                'citation_type': 'legacy',
                'label': title,
                'title': title,
                'url': citation.get('url'),
                'filepath': citation.get('filepath')
            })
        else:
            normalized.append({
                'citation_type': 'legacy',
                'label': str(citation),
                'value': str(citation)
            })

    return normalized


def _build_document_citation_label(citation: Dict[str, Any]) -> str:
    file_name = citation.get('file_name') or citation.get('title') or 'Document source'
    metadata_type = citation.get('metadata_type')
    page_number = citation.get('page_number')

    if metadata_type:
        return f"{file_name} — {metadata_type.replace('_', ' ').title()}"
    if page_number not in (None, ''):
        return f"{file_name} — Page {page_number}"
    return file_name


def _build_citation_counts(citations: List[Dict[str, Any]]) -> Dict[str, int]:
    counts = {
        'document': 0,
        'web': 0,
        'agent_tool': 0,
        'legacy': 0,
        'total': len(citations)
    }
    for citation in citations:
        citation_type = citation.get('citation_type')
        if citation_type in counts:
            counts[citation_type] += 1
    return counts


def _curate_message_details(
    message: Dict[str, Any],
    citation_counts: Dict[str, int],
    thought_count: int
) -> Dict[str, Any]:
    role = message.get('role', '')
    metadata = message.get('metadata', {}) or {}
    details: Dict[str, Any] = {}

    if role == 'user':
        details['interaction_mode'] = _remove_empty_values({
            'button_states': metadata.get('button_states'),
            'workspace_search': _curate_workspace_search(metadata.get('workspace_search')),
            'prompt_selection': _curate_prompt_selection(metadata.get('prompt_selection')),
            'agent_selection': _curate_agent_selection(metadata.get('agent_selection')),
            'model_selection': _curate_model_selection(metadata.get('model_selection'))
        })
    elif role == 'assistant':
        details['generation'] = _remove_empty_values({
            'augmented': message.get('augmented'),
            'model_deployment': message.get('model_deployment_name'),
            'agent_name': message.get('agent_name'),
            'agent_display_name': message.get('agent_display_name'),
            'reasoning_effort': metadata.get('reasoning_effort'),
            'hybrid_search_query': message.get('hybridsearch_query'),
            'token_usage': _curate_token_usage(metadata.get('token_usage')),
            'citation_counts': citation_counts,
            'thought_count': thought_count
        })
    else:
        details['message_context'] = _remove_empty_values({
            'filename': message.get('filename'),
            'prompt': message.get('prompt'),
            'is_table': message.get('is_table'),
            'model_deployment': message.get('model_deployment_name')
        })

    return _remove_empty_values(details)


def _curate_workspace_search(workspace_search: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    if not isinstance(workspace_search, dict):
        return {}
    return _remove_empty_values({
        'search_enabled': workspace_search.get('search_enabled'),
        'document_scope': workspace_search.get('document_scope'),
        'document_name': workspace_search.get('document_name'),
        'document_filename': workspace_search.get('document_filename'),
        'group_name': workspace_search.get('group_name'),
        'classification': workspace_search.get('classification'),
        'public_workspace_id': workspace_search.get('active_public_workspace_id')
    })


def _curate_prompt_selection(prompt_selection: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    if not isinstance(prompt_selection, dict):
        return {}
    return _remove_empty_values({
        'prompt_name': prompt_selection.get('prompt_name'),
        'selected_prompt_index': prompt_selection.get('selected_prompt_index'),
        'selected_prompt_text': prompt_selection.get('selected_prompt_text')
    })


def _curate_agent_selection(agent_selection: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    if not isinstance(agent_selection, dict):
        return {}
    return _remove_empty_values({
        'selected_agent': agent_selection.get('selected_agent'),
        'agent_display_name': agent_selection.get('agent_display_name'),
        'is_global': agent_selection.get('is_global'),
        'is_group': agent_selection.get('is_group'),
        'group_name': agent_selection.get('group_name')
    })


def _curate_model_selection(model_selection: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    if not isinstance(model_selection, dict):
        return {}
    return _remove_empty_values({
        'selected_model': model_selection.get('selected_model'),
        'frontend_requested_model': model_selection.get('frontend_requested_model'),
        'reasoning_effort': model_selection.get('reasoning_effort'),
        'streaming': model_selection.get('streaming')
    })


def _curate_token_usage(token_usage: Any) -> Dict[str, Any]:
    if not isinstance(token_usage, dict):
        return {}
    return _remove_empty_values({
        'prompt_tokens': token_usage.get('prompt_tokens'),
        'completion_tokens': token_usage.get('completion_tokens'),
        'total_tokens': token_usage.get('total_tokens')
    })


def _remove_empty_values(value: Any) -> Any:
    if isinstance(value, dict):
        cleaned = {}
        for key, item in value.items():
            cleaned_item = _remove_empty_values(item)
            if cleaned_item in (None, '', [], {}):
                continue
            cleaned[key] = cleaned_item
        return cleaned

    if isinstance(value, list):
        cleaned_list = []
        for item in value:
            cleaned_item = _remove_empty_values(item)
            if cleaned_item in (None, '', [], {}):
                continue
            cleaned_list.append(cleaned_item)
        return cleaned_list

    return value


def generate_conversation_summary(
    messages: List[Dict[str, Any]],
    conversation_title: str,
    settings: Dict[str, Any],
    model_deployment: str,
    message_time_start: str = None,
    message_time_end: str = None,
    conversation_id: str = None,
    user_id: str = None,
    model_endpoint_id: str = '',
    model_id: str = '',
    model_provider: str = ''
) -> Dict[str, Any]:
    """Generate a conversation summary using the LLM and optionally persist it.

    This is the shared helper used by both the export pipeline and the
    on-demand summary API endpoint.  Returns a summary dict suitable for
    storage in conversation metadata.

    Raises ValueError when there is no content to summarise and
    RuntimeError on model errors.
    """
    transcript_lines = []
    for message in messages:
        content_text = message.get('content_text', '')
        if not content_text:
            continue
        role = message.get('role', 'unknown')
        speaker = message.get('speaker_label', role).upper()
        transcript_lines.append(f"{speaker}: {content_text}")

    transcript_text = '\n\n'.join(transcript_lines).strip()
    if not transcript_text:
        raise ValueError('No message content was available to summarize.')

    transcript_text = _truncate_for_summary(transcript_text)

    gpt_client, gpt_model = _initialize_gpt_client(
        settings,
        model_deployment,
        user_id=user_id,
        requested_endpoint_id=model_endpoint_id,
        requested_model_id=model_id,
        requested_provider=model_provider,
    )
    summary_prompt = (
        "You are summarizing a conversation for an export document. "
        "Read the full conversation below and write a concise summary. "
        "Use your judgement on length: for short conversations write one brief paragraph, "
        "for longer or more detailed conversations write two paragraphs. "
        "If you need refer to the user, use their name, but do not refer to the user too often."
        "Cover the goals, the key topics discussed, any data or tools referenced, "
        "and the main outcomes or answers provided. "
        "Be factual and neutral. Return plain text only — no headings, no bullet points, no markdown formatting."
    )

    model_lower = gpt_model.lower()
    is_reasoning_model = (
        'o1' in model_lower or 'o3' in model_lower or 'gpt-5' in model_lower
    )
    instruction_role = 'developer' if is_reasoning_model else 'system'

    debug_print(f"Summary generation: sending {len(transcript_lines)} messages "
                f"({len(transcript_text)} chars) to {gpt_model} (role={instruction_role})")

    summary_response = gpt_client.chat.completions.create(
        model=gpt_model,
        messages=[
            {
                'role': instruction_role,
                'content': summary_prompt
            },
            {
                'role': 'user',
                'content': (
                    f"Conversation Title: {conversation_title}\n\n"
                    f"{transcript_text}"
                )
            }
        ]
    )

    debug_print(f"Summary generation: response choices="
                f"{len(summary_response.choices) if summary_response.choices else 0}, "
                f"finish_reason={summary_response.choices[0].finish_reason if summary_response.choices else 'N/A'}")

    summary_text = (summary_response.choices[0].message.content or '').strip() if summary_response.choices else ''
    if not summary_text:
        debug_print('Summary generation: model returned an empty response')
        log_event('Conversation summary generation returned empty response', level='WARNING')
        raise RuntimeError('Summary model returned an empty response.')

    summary_data = {
        'content': summary_text,
        'model_deployment': gpt_model,
        'generated_at': datetime.utcnow().isoformat(),
        'message_time_start': message_time_start,
        'message_time_end': message_time_end
    }

    # Persist to Cosmos when a conversation_id is available
    if conversation_id:
        try:
            summary_persisted = update_conversation_with_metadata(conversation_id, {'summary': summary_data})
            if summary_persisted:
                debug_print(f"Summary persisted to conversation {conversation_id}")
            else:
                debug_print(f"Summary was generated but not persisted for conversation {conversation_id}")
                log_event(
                    f"Conversation summary persistence returned false for {conversation_id}",
                    level='WARNING'
                )
        except Exception as persist_exc:
            debug_print(f"Failed to persist summary to Cosmos: {persist_exc}")
            log_event(f"Failed to persist conversation summary: {persist_exc}", level="WARNING")

    return summary_data


def _build_summary_intro(
    messages: List[Dict[str, Any]],
    conversation: Dict[str, Any],
    sanitized_conversation: Dict[str, Any],
    settings: Dict[str, Any],
    enabled: bool,
    summary_model_deployment: str,
    summary_model_endpoint_id: str = '',
    summary_model_id: str = '',
    summary_model_provider: str = '',
    user_id: str = None,
    message_time_start: str = None,
    message_time_end: str = None
) -> Dict[str, Any]:
    """Build the summary_intro block for the export payload.

    Uses cached summary from conversation metadata when present and
    still current (no newer messages).  Otherwise generates a fresh
    summary via ``generate_conversation_summary`` and persists it.
    """
    summary_intro = {
        'enabled': enabled,
        'generated': False,
        'model_deployment': summary_model_deployment or None,
        'generated_at': None,
        'content': '',
        'error': None
    }

    if not enabled:
        return summary_intro

    # Check for a cached summary stored in the conversation document
    existing_summary = conversation.get('summary')
    if existing_summary and isinstance(existing_summary, dict):
        cached_end = existing_summary.get('message_time_end')
        if cached_end and message_time_end and cached_end >= message_time_end:
            debug_print('Export summary: using cached summary from conversation metadata')
            summary_intro.update({
                'generated': True,
                'model_deployment': existing_summary.get('model_deployment'),
                'generated_at': existing_summary.get('generated_at'),
                'content': existing_summary.get('content', ''),
                'error': None
            })
            return summary_intro
        debug_print('Export summary: cached summary is stale, regenerating')

    try:
        conversation_id = conversation.get('id')
        conversation_title = sanitized_conversation.get('title', 'Untitled')

        summary_data = generate_conversation_summary(
            messages=messages,
            conversation_title=conversation_title,
            settings=settings,
            model_deployment=summary_model_deployment,
            message_time_start=message_time_start,
            message_time_end=message_time_end,
            conversation_id=conversation_id,
            user_id=user_id,
            model_endpoint_id=summary_model_endpoint_id,
            model_id=summary_model_id,
            model_provider=summary_model_provider,
        )

        summary_intro.update({
            'generated': True,
            'model_deployment': summary_data.get('model_deployment'),
            'generated_at': summary_data.get('generated_at'),
            'content': summary_data.get('content', ''),
            'error': None
        })
        return summary_intro

    except (ValueError, RuntimeError) as known_exc:
        debug_print(f"Export summary generation issue: {known_exc}")
        summary_intro['error'] = str(known_exc)
        if hasattr(known_exc, 'model_deployment'):
            summary_intro['model_deployment'] = known_exc.model_deployment
        return summary_intro

    except Exception as exc:
        debug_print(f"Export summary generation failed: {exc}")
        log_event(f"Conversation export summary generation failed: {exc}", level="WARNING")
        summary_intro['error'] = str(exc)
        return summary_intro


def _truncate_for_summary(transcript_text: str) -> str:
    if len(transcript_text) <= SUMMARY_SOURCE_CHAR_LIMIT:
        return transcript_text

    head_chars = SUMMARY_SOURCE_CHAR_LIMIT // 2
    tail_chars = SUMMARY_SOURCE_CHAR_LIMIT - head_chars
    return (
        transcript_text[:head_chars]
        + "\n\n[... transcript truncated for export summary generation ...]\n\n"
        + transcript_text[-tail_chars:]
    )


def _normalize_summary_model_value(value: Any) -> str:
    return str(value or '').strip()


def _append_summary_endpoint_candidates(
    candidates: List[Dict[str, Any]],
    endpoints: List[Dict[str, Any]],
    endpoint_scope: str,
) -> None:
    normalized_endpoints, _ = normalize_model_endpoints(endpoints or [])
    for endpoint in normalized_endpoints:
        if isinstance(endpoint, dict):
            candidates.append({**endpoint, '_endpoint_scope': endpoint_scope})


def _get_summary_model_endpoint_candidates(settings: Dict[str, Any], user_id: str = None) -> List[Dict[str, Any]]:
    candidates: List[Dict[str, Any]] = []
    settings = settings or {}

    _append_summary_endpoint_candidates(candidates, settings.get('model_endpoints', []) or [], 'global')

    if not user_id:
        return candidates

    if settings.get('allow_user_custom_endpoints', False):
        try:
            user_settings_doc = get_user_settings(user_id)
            user_settings = user_settings_doc.get('settings', {}) if isinstance(user_settings_doc, dict) else {}
            _append_summary_endpoint_candidates(
                candidates,
                user_settings.get('personal_model_endpoints', []) or [],
                'user',
            )
        except Exception as exc:
            debug_print(f"[Summary][Model Resolution] Failed to load personal endpoints: {exc}")

    if settings.get('enable_group_workspaces', False) and settings.get('allow_group_custom_endpoints', False):
        try:
            user_groups = get_user_groups(user_id)
        except Exception as exc:
            user_groups = []
            debug_print(f"[Summary][Model Resolution] Failed to load user groups: {exc}")

        for group_doc in user_groups:
            group_id = _normalize_summary_model_value(group_doc.get('id') if isinstance(group_doc, dict) else '')
            if not group_id:
                continue
            try:
                _append_summary_endpoint_candidates(
                    candidates,
                    get_group_model_endpoints(group_id) or [],
                    'group',
                )
            except Exception as exc:
                debug_print(
                    f"[Summary][Model Resolution] Failed to load group endpoints for group_id={group_id}: {exc}"
                )

    return candidates


def _summary_model_matches(model_cfg: Dict[str, Any], requested_model: str, requested_model_id: str) -> bool:
    model_values = {
        _normalize_summary_model_value(model_cfg.get('id')),
        _normalize_summary_model_value(model_cfg.get('deploymentName')),
        _normalize_summary_model_value(model_cfg.get('deployment')),
        _normalize_summary_model_value(model_cfg.get('modelName')),
        _normalize_summary_model_value(model_cfg.get('name')),
    }
    model_values.discard('')

    if requested_model_id and requested_model_id in model_values:
        return True
    return bool(requested_model and requested_model in model_values)


def _find_summary_endpoint_model(
    endpoint_cfg: Dict[str, Any],
    requested_model: str,
    requested_model_id: str,
) -> Optional[Dict[str, Any]]:
    models = endpoint_cfg.get('models', []) or []
    for model_cfg in models:
        if not isinstance(model_cfg, dict) or not model_cfg.get('enabled', True):
            continue
        if _summary_model_matches(model_cfg, requested_model, requested_model_id):
            return model_cfg
    return None


def _resolve_summary_foundry_scope_for_auth(auth_settings: Dict[str, Any], endpoint: str = None) -> str:
    auth_settings = auth_settings or {}
    custom_scope = _normalize_summary_model_value(auth_settings.get('foundry_scope'))
    if custom_scope:
        return custom_scope

    management_cloud = _normalize_summary_model_value(auth_settings.get('management_cloud') or 'public').lower()
    if management_cloud in ('government', 'usgovernment', 'usgov'):
        return 'https://ai.azure.us/.default'
    if management_cloud == 'china':
        return 'https://ai.azure.cn/.default'
    if management_cloud == 'germany':
        return 'https://ai.azure.de/.default'

    endpoint_value = _normalize_summary_model_value(endpoint).lower()
    if 'azure.us' in endpoint_value:
        return 'https://ai.azure.us/.default'
    if 'azure.cn' in endpoint_value:
        return 'https://ai.azure.cn/.default'
    if 'azure.de' in endpoint_value:
        return 'https://ai.azure.de/.default'

    return 'https://ai.azure.com/.default'


def _build_summary_model_endpoint_client(
    auth_settings: Dict[str, Any],
    provider: str,
    endpoint: str,
    api_version: str,
    deployment_name: str,
):
    auth_settings = auth_settings or {}
    auth_type = _normalize_summary_model_value(auth_settings.get('type') or 'managed_identity').lower()
    normalized_provider = _normalize_summary_model_value(provider or 'aoai').lower()
    runtime_protocol = infer_model_endpoint_protocol(normalized_provider, endpoint, deployment_name)

    if auth_type in ('api_key', 'key'):
        api_key = auth_settings.get('api_key')
        if not api_key:
            raise ValueError('Selected summary model endpoint is missing an API key.')
        if runtime_protocol == MODEL_ENDPOINT_PROTOCOL_ANTHROPIC:
            return build_anthropic_chat_client(endpoint=endpoint, api_key=api_key)
        if runtime_protocol == MODEL_ENDPOINT_PROTOCOL_OPENAI_STYLE:
            return build_openai_style_chat_client(api_key, endpoint, api_version)
        return AzureOpenAI(
            api_version=api_version,
            azure_endpoint=endpoint,
            api_key=api_key,
        )

    if auth_type == 'service_principal':
        credential = ClientSecretCredential(
            tenant_id=auth_settings.get('tenant_id'),
            client_id=auth_settings.get('client_id'),
            client_secret=auth_settings.get('client_secret'),
            authority=resolve_authority(auth_settings),
        )
    else:
        managed_identity_client_id = auth_settings.get('managed_identity_client_id') or None
        credential = DefaultAzureCredential(managed_identity_client_id=managed_identity_client_id)

    scope = cognitive_services_scope
    if normalized_provider in ('aifoundry', 'new_foundry') or runtime_protocol != MODEL_ENDPOINT_PROTOCOL_AZURE_OPENAI:
        scope = _resolve_summary_foundry_scope_for_auth(auth_settings, endpoint=endpoint)

    if runtime_protocol == MODEL_ENDPOINT_PROTOCOL_ANTHROPIC:
        token = credential.get_token(scope).token
        return build_anthropic_chat_client(endpoint=endpoint, bearer_token=token)

    if runtime_protocol == MODEL_ENDPOINT_PROTOCOL_OPENAI_STYLE:
        token = credential.get_token(scope).token
        return build_openai_style_chat_client(token, endpoint, api_version)

    token_provider = get_bearer_token_provider(credential, scope)
    return AzureOpenAI(
        api_version=api_version,
        azure_endpoint=endpoint,
        azure_ad_token_provider=token_provider,
    )


def _resolve_summary_multi_endpoint_client(
    settings: Dict[str, Any],
    requested_model: str = '',
    user_id: str = None,
    requested_endpoint_id: str = '',
    requested_model_id: str = '',
    requested_provider: str = '',
):
    settings = settings or {}
    if not settings.get('enable_multi_model_endpoints', False):
        return None

    requested_model = _normalize_summary_model_value(requested_model)
    requested_endpoint_id = _normalize_summary_model_value(requested_endpoint_id)
    requested_model_id = _normalize_summary_model_value(requested_model_id)
    requested_provider = _normalize_summary_model_value(requested_provider).lower()
    selection_source = 'request' if (requested_model or requested_endpoint_id or requested_model_id) else ''

    if not selection_source:
        default_selection = settings.get('default_model_selection', {}) or {}
        requested_endpoint_id = _normalize_summary_model_value(default_selection.get('endpoint_id'))
        requested_model_id = _normalize_summary_model_value(default_selection.get('model_id'))
        requested_provider = requested_provider or _normalize_summary_model_value(default_selection.get('provider')).lower()
        selection_source = 'default' if (requested_endpoint_id or requested_model_id) else ''

    if not selection_source:
        return None

    endpoint_candidates = _get_summary_model_endpoint_candidates(settings, user_id=user_id)
    if requested_endpoint_id:
        endpoint_candidates = [
            endpoint for endpoint in endpoint_candidates
            if _normalize_summary_model_value(endpoint.get('id')) == requested_endpoint_id
        ]
        if not endpoint_candidates:
            if selection_source == 'request':
                raise ValueError('Selected summary model endpoint could not be found.')
            return None

    for endpoint_cfg in endpoint_candidates:
        if not isinstance(endpoint_cfg, dict) or not endpoint_cfg.get('enabled', True):
            continue

        model_cfg = _find_summary_endpoint_model(endpoint_cfg, requested_model, requested_model_id)
        if not model_cfg:
            continue

        endpoint_scope = endpoint_cfg.get('_endpoint_scope', 'global')
        resolved_endpoint_cfg = dict(endpoint_cfg)
        resolved_endpoint_cfg.pop('_endpoint_scope', None)
        endpoint_id = _normalize_summary_model_value(resolved_endpoint_cfg.get('id'))
        resolved_endpoint_cfg = keyvault_model_endpoint_get_helper(
            resolved_endpoint_cfg,
            endpoint_id,
            scope=endpoint_scope,
            return_type=SecretReturnType.VALUE,
        )

        provider = _normalize_summary_model_value(resolved_endpoint_cfg.get('provider') or requested_provider or 'aoai').lower()
        connection = resolved_endpoint_cfg.get('connection', {}) or {}
        auth_settings = resolved_endpoint_cfg.get('auth', {}) or {}
        deployment = _normalize_summary_model_value(
            model_cfg.get('deploymentName') or model_cfg.get('deployment') or model_cfg.get('id')
        )
        endpoint = _normalize_summary_model_value(connection.get('endpoint'))
        api_version = _normalize_summary_model_value(connection.get('openai_api_version') or connection.get('api_version'))
        runtime_protocol = infer_model_endpoint_protocol(provider, endpoint, deployment)

        missing_required_config = not endpoint or not deployment or (
            runtime_protocol == MODEL_ENDPOINT_PROTOCOL_AZURE_OPENAI and not api_version
        )
        if missing_required_config:
            if selection_source == 'request' and requested_endpoint_id:
                raise ValueError('Selected summary model endpoint is missing endpoint or deployment configuration.')
            continue

        gpt_client = _build_summary_model_endpoint_client(
            auth_settings,
            provider,
            endpoint,
            api_version,
            deployment,
        )
        debug_print(
            f"[Summary][Model Resolution] Resolved {selection_source} multi-endpoint model | "
            f"provider={provider} | endpoint_id={endpoint_id} | model_id={model_cfg.get('id')} | "
            f"deployment={deployment} | api_version={api_version} | protocol={runtime_protocol}"
        )
        return gpt_client, deployment

    if selection_source == 'request' and requested_endpoint_id:
        raise ValueError('Selected summary model could not be found on the configured endpoint.')

    return None


def _initialize_gpt_client(
    settings: Dict[str, Any],
    requested_model: str = '',
    user_id: str = None,
    requested_endpoint_id: str = '',
    requested_model_id: str = '',
    requested_provider: str = '',
):
    settings = settings or {}
    multi_endpoint_client = _resolve_summary_multi_endpoint_client(
        settings,
        requested_model=requested_model,
        user_id=user_id,
        requested_endpoint_id=requested_endpoint_id,
        requested_model_id=requested_model_id,
        requested_provider=requested_provider,
    )
    if multi_endpoint_client:
        return multi_endpoint_client

    enable_gpt_apim = settings.get('enable_gpt_apim', False)

    if enable_gpt_apim:
        raw_models = settings.get('azure_apim_gpt_deployment', '') or ''
        apim_models = [model.strip() for model in raw_models.split(',') if model.strip()]
        if not apim_models:
            raise ValueError('APIM GPT deployment name is not configured.')

        if requested_model and requested_model not in apim_models:
            raise ValueError(f"Requested summary model '{requested_model}' is not configured for APIM.")

        gpt_model = requested_model or apim_models[0]
        gpt_client = AzureOpenAI(
            api_version=settings.get('azure_apim_gpt_api_version'),
            azure_endpoint=settings.get('azure_apim_gpt_endpoint'),
            api_key=settings.get('azure_apim_gpt_subscription_key')
        )
        return gpt_client, gpt_model

    auth_type = settings.get('azure_openai_gpt_authentication_type')
    endpoint = settings.get('azure_openai_gpt_endpoint')
    api_version = settings.get('azure_openai_gpt_api_version')
    gpt_model_obj = settings.get('gpt_model', {}) or {}

    if requested_model:
        gpt_model = requested_model
    elif gpt_model_obj.get('selected'):
        gpt_model = gpt_model_obj['selected'][0]['deploymentName']
    else:
        raise ValueError('No GPT model selected or configured for export summary generation.')

    if auth_type == 'managed_identity':
        token_provider = get_bearer_token_provider(DefaultAzureCredential(), cognitive_services_scope)
        gpt_client = AzureOpenAI(
            api_version=api_version,
            azure_endpoint=endpoint,
            azure_ad_token_provider=token_provider
        )
    else:
        api_key = settings.get('azure_openai_gpt_key')
        if not api_key:
            raise ValueError('Azure OpenAI API Key not configured.')
        gpt_client = AzureOpenAI(
            api_version=api_version,
            azure_endpoint=endpoint,
            api_key=api_key
        )

    return gpt_client, gpt_model


def _build_single_file_response(exported: List[Dict[str, Any]], export_format: str, timestamp_str: str):
    """Build a single-file download response."""
    if export_format == 'json':
        content = json.dumps(exported, indent=2, ensure_ascii=False, default=str)
        filename = f"conversations_export_{timestamp_str}.json"
        content_type = 'application/json; charset=utf-8'
    elif export_format == 'pdf':
        if len(exported) == 1:
            content = _conversation_to_pdf_bytes(exported[0])
        else:
            combined_parts = []
            for idx, entry in enumerate(exported):
                if idx > 0:
                    combined_parts.append(
                        '<div style="margin-top: 24pt; border-top: 2px solid #999; '
                        'padding-top: 12pt;"></div>'
                    )
                combined_parts.append(_build_pdf_html_body(entry))
            content = _html_body_to_pdf_bytes('\n'.join(combined_parts))
        filename = f"conversations_export_{timestamp_str}.pdf"
        content_type = 'application/pdf'
    else:
        parts = []
        for entry in exported:
            parts.append(_conversation_to_markdown(entry))
        content = '\n\n---\n\n'.join(parts)
        filename = f"conversations_export_{timestamp_str}.md"
        content_type = 'text/markdown; charset=utf-8'

    response = make_response(content)
    response.headers['Content-Type'] = content_type
    response.headers['Content-Disposition'] = f'attachment; filename="{filename}"'
    return response


def _build_zip_response(exported: List[Dict[str, Any]], export_format: str, timestamp_str: str):
    """Build a ZIP archive containing one file per conversation."""
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
        for entry in exported:
            conversation = entry['conversation']
            safe_title = _safe_filename(conversation.get('title', 'Untitled'))
            conversation_id_short = conversation.get('id', 'unknown')[:8]

            if export_format == 'json':
                file_content = json.dumps(entry, indent=2, ensure_ascii=False, default=str)
                ext = 'json'
            elif export_format == 'pdf':
                file_content = _conversation_to_pdf_bytes(entry)
                ext = 'pdf'
            else:
                file_content = _conversation_to_markdown(entry)
                ext = 'md'

            file_name = f"{safe_title}_{conversation_id_short}.{ext}"
            zf.writestr(file_name, file_content)

    buffer.seek(0)
    filename = f"conversations_export_{timestamp_str}.zip"

    response = make_response(buffer.read())
    response.headers['Content-Type'] = 'application/zip'
    response.headers['Content-Disposition'] = f'attachment; filename="{filename}"'
    return response


def _conversation_to_markdown(entry: Dict[str, Any]) -> str:
    """Convert a conversation + messages entry to Markdown format."""
    conversation = entry['conversation']
    messages = entry['messages']
    summary_intro = entry.get('summary_intro', {}) or {}

    transcript_messages = [message for message in messages if message.get('is_transcript_message')]
    detail_messages = [message for message in messages if message.get('details')]
    reference_messages = [message for message in messages if message.get('citations')]
    thought_messages = [message for message in messages if message.get('thoughts')]
    supplemental_messages = [message for message in messages if not message.get('is_transcript_message')]

    lines: List[str] = []
    lines.append(f"# {conversation.get('title', 'Untitled')}")
    lines.append('')
    lines.append(f"**Last Updated:** {conversation.get('last_updated', '')}  ")
    lines.append(f"**Chat Type:** {conversation.get('chat_type', 'personal')}  ")
    lines.append(f"**Messages:** {conversation.get('message_count', len(messages))}  ")
    if conversation.get('tags'):
        lines.append(f"**Tags:** {', '.join(_format_tag(tag) for tag in conversation.get('tags', []))}  ")
    if conversation.get('classification'):
        lines.append(f"**Classification:** {', '.join(_format_tag(item) for item in conversation.get('classification', []))}  ")
    lines.append('')

    if summary_intro.get('enabled') and summary_intro.get('generated') and summary_intro.get('content'):
        lines.append('## Abstract')
        lines.append('')
        lines.append(summary_intro.get('content', ''))
        lines.append('')
        lines.append(f"_Generated with {summary_intro.get('model_deployment') or 'configured model'} on {summary_intro.get('generated_at')}_")
        lines.append('')
    elif summary_intro.get('enabled') and summary_intro.get('error'):
        lines.append('> _A summary intro was requested, but it could not be generated for this export._')
        lines.append(f"> _Error: {summary_intro.get('error')}_")
        lines.append('')

    lines.append('## Transcript')
    lines.append('')
    if not transcript_messages:
        lines.append('_No user or assistant transcript messages were available for export._')
        lines.append('')
    else:
        for message in transcript_messages:
            lines.append(f"### {message.get('label')} — {message.get('speaker_label')}")
            if message.get('timestamp'):
                lines.append(f"*{message.get('timestamp')}*")
            lines.append('')
            lines.append(
                replace_inline_chart_blocks_with_export_html(
                    message.get('content_text') or '_No content recorded._'
                )
            )
            lines.append('')

    lines.append('## Appendix A — Conversation Metadata')
    lines.append('')
    metadata_to_render = _remove_empty_values({
        'context': conversation.get('context'),
        'classification': conversation.get('classification'),
        'strict': conversation.get('strict'),
        'is_pinned': conversation.get('is_pinned'),
        'scope_locked': conversation.get('scope_locked'),
        'locked_contexts': conversation.get('locked_contexts'),
        'message_counts_by_role': conversation.get('message_counts_by_role'),
        'citation_counts': conversation.get('citation_counts'),
        'thought_count': conversation.get('thought_count')
    })
    _append_markdown_mapping(lines, metadata_to_render)
    lines.append('')

    if detail_messages:
        lines.append('## Appendix B — Message Details')
        lines.append('')
        for message in detail_messages:
            lines.append(f"### {message.get('label')} — {message.get('speaker_label')}")
            if message.get('timestamp'):
                lines.append(f"*{message.get('timestamp')}*")
            lines.append('')
            _append_markdown_mapping(lines, message.get('details', {}))
            lines.append('')

    if reference_messages:
        lines.append('## Appendix C — References')
        lines.append('')
        for message in reference_messages:
            lines.append(f"### {message.get('label')} — {message.get('speaker_label')}")
            if message.get('timestamp'):
                lines.append(f"*{message.get('timestamp')}*")
            lines.append('')
            _append_citations_markdown(lines, message)
            lines.append('')

    if thought_messages:
        lines.append('## Appendix D — Processing Thoughts')
        lines.append('')
        for message in thought_messages:
            lines.append(f"### {message.get('label')} — {message.get('speaker_label')}")
            if message.get('timestamp'):
                lines.append(f"*{message.get('timestamp')}*")
            lines.append('')
            for thought in message.get('thoughts', []):
                thought_label = thought.get('step_type', 'step').replace('_', ' ').title()
                lines.append(f"1. **{thought_label}:** {thought.get('content') or 'No content recorded.'}")
                if thought.get('duration_ms') is not None:
                    lines.append(f"   - **Duration:** {thought.get('duration_ms')} ms")
                if thought.get('timestamp'):
                    lines.append(f"   - **Timestamp:** {thought.get('timestamp')}")
                if thought.get('detail'):
                    lines.append('   - **Detail:**')
                    _append_code_block(lines, thought.get('detail'), indent='     ')
            lines.append('')

    if supplemental_messages:
        lines.append('## Appendix E — Supplemental Messages')
        lines.append('')
        for message in supplemental_messages:
            lines.append(f"### {message.get('label')} — {message.get('speaker_label')}")
            if message.get('timestamp'):
                lines.append(f"*{message.get('timestamp')}*")
            lines.append('')
            lines.append(
                replace_inline_chart_blocks_with_export_html(
                    message.get('content_text') or '_No content recorded._'
                )
            )
            lines.append('')

    return '\n'.join(lines).strip()


def _append_citations_markdown(lines: List[str], message: Dict[str, Any]):
    document_citations = [citation for citation in message.get('citations', []) if citation.get('citation_type') == 'document']
    web_citations = [citation for citation in message.get('citations', []) if citation.get('citation_type') == 'web']
    agent_citations = message.get('agent_citations', []) or []
    legacy_citations = [citation for citation in message.get('citations', []) if citation.get('citation_type') == 'legacy']

    if not any([document_citations, web_citations, agent_citations, legacy_citations]):
        lines.append('_No citations were recorded for this message._')
        return

    if document_citations:
        lines.append('#### Document Sources')
        lines.append('')
        for index, citation in enumerate(document_citations, start=1):
            lines.append(f"{index}. **{citation.get('label', 'Document source')}**")
            detail_mapping = _remove_empty_values({
                'citation_id': citation.get('citation_id'),
                'page_number': citation.get('page_number'),
                'classification': citation.get('classification'),
                'score': citation.get('score'),
                'metadata_type': citation.get('metadata_type')
            })
            _append_markdown_mapping(lines, detail_mapping, indent=1)
            if citation.get('metadata_content'):
                lines.append('   - **Metadata Content:**')
                _append_code_block(lines, citation.get('metadata_content'), indent='     ')
        lines.append('')

    if web_citations:
        lines.append('#### Web Sources')
        lines.append('')
        for index, citation in enumerate(web_citations, start=1):
            title = citation.get('title') or citation.get('label') or 'Web source'
            url = citation.get('url')
            if url:
                lines.append(f"{index}. [{title}]({url})")
            else:
                lines.append(f"{index}. {title}")
        lines.append('')

    if agent_citations:
        lines.append('#### Tool Invocations')
        lines.append('')
        for index, citation in enumerate(agent_citations, start=1):
            label = citation.get('tool_name') or citation.get('function_name') or f"Tool {index}"
            lines.append(f"{index}. **{label}**")
            detail_mapping = _remove_empty_values({
                'function_name': citation.get('function_name'),
                'plugin_name': citation.get('plugin_name'),
                'success': citation.get('success'),
                'timestamp': citation.get('timestamp')
            })
            _append_markdown_mapping(lines, detail_mapping, indent=1)
            if citation.get('function_arguments') not in (None, '', [], {}):
                lines.append('   - **Arguments:**')
                _append_code_block(lines, citation.get('function_arguments'), indent='     ')
            if citation.get('function_result') not in (None, '', [], {}):
                lines.append('   - **Result:**')
                _append_code_block(lines, citation.get('function_result'), indent='     ')
        lines.append('')

    if legacy_citations:
        lines.append('#### Legacy Citation Records')
        lines.append('')
        for index, citation in enumerate(legacy_citations, start=1):
            lines.append(f"{index}. {citation.get('label', 'Legacy citation')}")
        lines.append('')


def _append_markdown_mapping(lines: List[str], mapping: Dict[str, Any], indent: int = 0):
    if not isinstance(mapping, dict) or not mapping:
        return

    prefix = '  ' * indent
    for key, value in mapping.items():
        label = _format_markdown_key(key)
        if isinstance(value, dict):
            lines.append(f"{prefix}- **{label}:**")
            _append_markdown_mapping(lines, value, indent + 1)
        elif isinstance(value, list):
            if not value:
                continue
            if all(not isinstance(item, (dict, list)) for item in value):
                lines.append(f"{prefix}- **{label}:** {', '.join(_stringify_markdown_value(item) for item in value)}")
            else:
                lines.append(f"{prefix}- **{label}:**")
                for item in value:
                    if isinstance(item, dict):
                        lines.append(f"{prefix}  -")
                        _append_markdown_mapping(lines, item, indent + 2)
                    else:
                        lines.append(f"{prefix}  - {_stringify_markdown_value(item)}")
        else:
            lines.append(f"{prefix}- **{label}:** {_stringify_markdown_value(value)}")


def _append_code_block(lines: List[str], value: Any, indent: str = ''):
    if isinstance(value, (dict, list)):
        code_block = json.dumps(value, indent=2, ensure_ascii=False, default=str)
        language = 'json'
    else:
        code_block = str(value)
        language = 'text'

    lines.append(f"{indent}```{language}")
    for line in code_block.splitlines() or ['']:
        lines.append(f"{indent}{line}")
    lines.append(f"{indent}```")


def _format_markdown_key(key: str) -> str:
    return str(key).replace('_', ' ').title()


def _stringify_markdown_value(value: Any) -> str:
    if isinstance(value, bool):
        return 'Yes' if value else 'No'
    return str(value)


def _format_tag(tag: Any) -> str:
    """Format a tag or classification entry for display.

    Tags in Cosmos are stored as dicts such as
    ``{'category': 'model', 'value': 'gpt-5'}`` or
    ``{'category': 'participant', 'name': 'Alice', 'user_id': '...'}``
    but they can also be plain strings in older data.
    """
    if isinstance(tag, dict):
        category = tag.get('category', '')
        # Participant tags carry a readable name / email
        name = tag.get('name') or tag.get('email') or tag.get('display_name')
        if name:
            return f"{category}: {name}" if category else str(name)
        # Document tags carry a title
        title = tag.get('title') or tag.get('document_id')
        if title:
            return f"{category}: {title}" if category else str(title)
        # Generic category/value tags
        value = tag.get('value')
        if value:
            return f"{category}: {value}" if category else str(value)
        return category or str(tag)
    return str(tag)


def _role_to_label(role: str) -> str:
    role_map = {
        'assistant': 'Assistant',
        'user': 'User',
        'system': 'System',
        'tool': 'Tool',
        'file': 'File',
        'image': 'Image',
        'safety': 'Safety',
        'blocked': 'Blocked'
    }
    return role_map.get(role, str(role).capitalize() or 'Message')


def _normalize_content(content: Any) -> str:
    """Normalize message content to a plain string."""
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        parts = []
        for item in content:
            if isinstance(item, dict):
                if item.get('type') == 'text':
                    parts.append(item.get('text', ''))
                elif item.get('type') == 'image_url':
                    parts.append('[Image]')
                else:
                    parts.append(str(item))
            else:
                parts.append(str(item))
        return '\n'.join(parts)
    if isinstance(content, dict):
        if content.get('type') == 'text':
            return content.get('text', '')
        return str(content)
    return str(content) if content else ''


def _safe_filename(title: str) -> str:
    """Create a filesystem-safe filename from a conversation title."""
    safe = re.sub(r'[<>:"/\\|?*]', '_', title)
    safe = re.sub(r'\s+', '_', safe)
    safe = safe.strip('_. ')
    if len(safe) > 50:
        safe = safe[:50]
    return safe or 'Untitled'


def _get_message_export_content_override(data: Dict[str, Any]) -> Optional[str]:
    if not isinstance(data, dict):
        return None

    if 'message_content_override' in data:
        raw_override = data.get('message_content_override')
    elif 'content_override' in data:
        raw_override = data.get('content_override')
    else:
        return None

    if raw_override in (None, ''):
        return None
    if not isinstance(raw_override, str):
        raise ValueError('message_content_override must be a string')

    normalized_override = raw_override.replace('\r\n', '\n').replace('\r', '\n')
    if not normalized_override.strip():
        return None
    if len(normalized_override) > MESSAGE_EXPORT_CONTENT_OVERRIDE_MAX_LENGTH:
        raise ValueError('message_content_override is too large')
    return normalized_override


def _apply_message_export_content_override(
    message: Dict[str, Any],
    message_content_override: Optional[str],
) -> Dict[str, Any]:
    if message_content_override is None:
        return message

    export_message = dict(message)
    export_message['content'] = message_content_override
    metadata = dict(export_message.get('metadata') or {})
    metadata['export_content_override_applied'] = True
    export_message['metadata'] = metadata
    return export_message


def _load_export_message_for_user(user_id: str, conversation_id: str, message_id: str) -> Dict[str, Any]:
    try:
        conversation = cosmos_conversations_container.read_item(
            item=conversation_id,
            partition_key=conversation_id
        )
    except Exception as exc:
        raise LookupError('Conversation not found') from exc

    if conversation.get('user_id') != user_id:
        raise PermissionError('Access denied')

    try:
        message = cosmos_messages_container.read_item(
            item=message_id,
            partition_key=conversation_id
        )
    except Exception:
        message_query = """
            SELECT * FROM c
            WHERE c.id = @message_id AND c.conversation_id = @conversation_id
        """
        message_results = list(cosmos_messages_container.query_items(
            query=message_query,
            parameters=[
                {'name': '@message_id', 'value': message_id},
                {'name': '@conversation_id', 'value': conversation_id}
            ],
            enable_cross_partition_query=True
        ))
        if not message_results:
            raise LookupError('Message not found')
        message = message_results[0]

    if message.get('conversation_id') != conversation_id:
        raise LookupError('Message not found')

    if isinstance(message.get('agent_citations'), list) and any(
        isinstance(citation, dict) and citation.get('artifact_id')
        for citation in message.get('agent_citations', [])
    ):
        conversation_messages = list(cosmos_messages_container.query_items(
            query="SELECT * FROM c WHERE c.conversation_id = @conversation_id",
            parameters=[{'name': '@conversation_id', 'value': conversation_id}],
            partition_key=conversation_id,
        ))
        artifact_payload_map = build_message_artifact_payload_map(conversation_messages)
        hydrated_messages = hydrate_agent_citations_from_artifacts([message], artifact_payload_map)
        if hydrated_messages:
            message = hydrated_messages[0]

    if message.get('role') == 'assistant':
        message = _attach_generated_image_proposal_assets(
            message,
            conversation_id=conversation_id,
        )

    return message


def _load_powerpoint_export_message_for_user(
    user_id: str,
    conversation_id: str,
    message_id: str,
    artifact_message_id: str = '',
) -> Dict[str, Any]:
    message = _load_export_message_for_user(user_id, conversation_id, message_id)
    normalized_artifact_message_id = str(artifact_message_id or '').strip()
    if not normalized_artifact_message_id:
        return message

    artifact_message = _load_generated_markdown_artifact_for_user(
        user_id=user_id,
        conversation_id=conversation_id,
        artifact_message_id=normalized_artifact_message_id,
    )
    artifact_bytes = download_blob_content(
        artifact_message.get('blob_container'),
        artifact_message.get('blob_path'),
    )
    try:
        artifact_content = artifact_bytes.decode('utf-8-sig')
    except UnicodeDecodeError:
        artifact_content = artifact_bytes.decode('utf-8', errors='replace')

    artifact_filename = str(artifact_message.get('filename') or 'generated-artifact.md').strip()
    metadata = dict(message.get('metadata') or {})
    metadata.update({
        'powerpoint_export_source': 'generated_markdown_artifact',
        'powerpoint_export_artifact_message_id': normalized_artifact_message_id,
        'powerpoint_export_artifact_filename': artifact_filename,
    })

    artifact_backed_message = dict(message)
    artifact_backed_message.update({
        'role': 'assistant',
        'content': artifact_content,
        'timestamp': artifact_message.get('timestamp') or message.get('timestamp'),
        'metadata': metadata,
    })
    return artifact_backed_message


def _load_generated_markdown_artifact_for_user(
    user_id: str,
    conversation_id: str,
    artifact_message_id: str,
) -> Dict[str, Any]:
    artifact_message = _load_export_message_for_user(
        user_id=user_id,
        conversation_id=conversation_id,
        message_id=artifact_message_id,
    )
    artifact_metadata = artifact_message.get('metadata') if isinstance(artifact_message.get('metadata'), dict) else {}
    artifact_filename = str(artifact_message.get('filename') or '').strip().lower()
    artifact_format = str(artifact_metadata.get('generated_artifact_output_format') or '').strip().lower()

    if artifact_message.get('role') != 'file' or not artifact_metadata.get('is_generated_chat_artifact'):
        raise LookupError('Generated Markdown artifact not found')

    if artifact_format not in {'md', 'markdown'} and not artifact_filename.endswith(('.md', '.markdown')):
        raise ValueError('Only generated Markdown artifacts can be exported as PowerPoint')

    if str(artifact_message.get('file_content_source') or '').strip().lower() != 'blob':
        raise LookupError('Generated Markdown artifact content is unavailable')

    if (
        not str(artifact_message.get('blob_container') or '').strip()
        or not str(artifact_message.get('blob_path') or '').strip()
    ):
        raise LookupError('Generated Markdown artifact content is unavailable')

    return artifact_message


def _attach_generated_image_proposal_assets(
    message: Dict[str, Any],
    conversation_id: str,
) -> Dict[str, Any]:
    message_id = str(message.get('id') or '').strip()
    if not message_id:
        return message

    image_assets = _load_generated_image_proposal_assets(
        conversation_id=conversation_id,
        source_assistant_message_id=message_id,
    )
    if not image_assets:
        return message

    export_message = dict(message)
    export_message['_export_generated_image_assets'] = image_assets
    return export_message


def _load_generated_image_proposal_assets(
    conversation_id: str,
    source_assistant_message_id: str,
) -> List[Dict[str, Any]]:
    normalized_source_id = str(source_assistant_message_id or '').strip()
    if not normalized_source_id:
        return []

    try:
        image_messages = list(cosmos_messages_container.query_items(
            query=(
                'SELECT * FROM c '
                'WHERE c.conversation_id = @conversation_id AND c.role = @role'
            ),
            parameters=[
                {'name': '@conversation_id', 'value': conversation_id},
                {'name': '@role', 'value': 'image'},
            ],
            partition_key=conversation_id,
        ))
    except Exception as exc:
        debug_print(f'Image proposal export lookup failed: {exc}')
        return []

    image_assets: List[Dict[str, Any]] = []
    for image_message in image_messages:
        metadata = image_message.get('metadata') if isinstance(image_message.get('metadata'), dict) else {}
        proposal = metadata.get('image_proposal') if isinstance(metadata.get('image_proposal'), dict) else {}
        proposal_source_id = str(proposal.get('source_assistant_message_id') or '').strip()
        if proposal_source_id != normalized_source_id:
            continue

        image_asset = _build_export_image_asset_from_message(
            conversation_id=conversation_id,
            image_message=image_message,
            proposal=proposal,
        )
        if image_asset:
            image_assets.append(image_asset)

    return image_assets


def _build_export_image_asset_from_message(
    conversation_id: str,
    image_message: Dict[str, Any],
    proposal: Dict[str, Any],
) -> Optional[Dict[str, Any]]:
    data_uri = _resolve_image_message_export_data_uri(conversation_id, image_message)
    if not data_uri:
        return None

    title = _clean_export_visual_text(
        proposal.get('title') or image_message.get('filename') or 'Generated image',
        160,
    )
    caption = _clean_export_visual_text(
        proposal.get('description') or proposal.get('context') or title,
        240,
    )

    return {
        'message_id': str(image_message.get('id') or '').strip(),
        'data_uri': data_uri,
        'content_type': 'image/png',
        'proposal': proposal,
        'title': title or 'Generated image',
        'caption': caption,
        'visual_id': _normalize_export_visual_id(proposal.get('visualId') or proposal.get('visual_id')),
        'prompt': _normalize_export_prompt(proposal.get('prompt')),
    }


def _resolve_image_message_export_data_uri(
    conversation_id: str,
    image_message: Dict[str, Any],
) -> str:
    try:
        if is_blob_backed_image_message(image_message):
            image_bytes = download_blob_content(
                image_message.get('blob_container'),
                image_message.get('blob_path'),
            )
            return _image_bytes_to_png_data_uri(image_bytes)

        message_id = str(image_message.get('id') or '').strip()
        if message_id:
            try:
                _, complete_content = get_complete_image_content(
                    cosmos_messages_container,
                    conversation_id,
                    message_id,
                )
            except Exception:
                complete_content = str(image_message.get('content') or '')
        else:
            complete_content = str(image_message.get('content') or '')

        if is_external_image_url(complete_content):
            return ''

        _, image_bytes = decode_image_content(complete_content)
        return _image_bytes_to_png_data_uri(image_bytes)
    except Exception as exc:
        debug_print(f'Image proposal export data URI resolution failed: {exc}')
        return ''


def _image_bytes_to_png_data_uri(image_bytes: bytes) -> str:
    if not image_bytes:
        return ''

    png_bytes = b''
    try:
        with Image.open(io.BytesIO(image_bytes)) as image:
            image_to_save = image
            if image.mode not in {'RGB', 'RGBA'}:
                image_to_save = image.convert('RGBA' if 'A' in image.getbands() else 'RGB')
            png_buffer = io.BytesIO()
            image_to_save.save(png_buffer, format='PNG')
            png_bytes = png_buffer.getvalue()
    except Exception:
        if bytes(image_bytes).startswith(b'\x89PNG'):
            png_bytes = bytes(image_bytes)

    if not png_bytes:
        return ''

    encoded_payload = base64.b64encode(png_bytes).decode('ascii')
    return f'data:image/png;base64,{encoded_payload}'


def _render_message_export_content(
    message: Dict[str, Any],
    source_content: Optional[Any] = None,
) -> str:
    raw_content = message.get('content', '') if source_content is None else source_content
    rendered_content = replace_inline_chart_blocks_with_export_html(
        _normalize_content(raw_content)
    )
    return _replace_inline_image_proposal_blocks_with_export_html(
        rendered_content,
        _get_message_export_image_assets(message),
    )


def _get_message_export_image_assets(message: Dict[str, Any]) -> List[Dict[str, Any]]:
    raw_assets = message.get('_export_generated_image_assets')
    if not isinstance(raw_assets, list):
        raw_assets = message.get('generated_image_proposals')
    if not isinstance(raw_assets, list):
        return []

    image_assets: List[Dict[str, Any]] = []
    for raw_asset in raw_assets:
        image_asset = _normalize_export_image_asset(raw_asset)
        if image_asset:
            image_assets.append(image_asset)
    return image_assets


def _normalize_export_image_asset(raw_asset: Any) -> Optional[Dict[str, Any]]:
    if not isinstance(raw_asset, dict):
        return None

    data_uri = str(
        raw_asset.get('data_uri')
        or raw_asset.get('image_url')
        or raw_asset.get('content')
        or ''
    ).strip()
    image_bytes = decode_base64_image_data_uri(data_uri)
    if not image_bytes:
        return None

    png_data_uri = _image_bytes_to_png_data_uri(image_bytes)
    if not png_data_uri:
        return None

    raw_proposal = raw_asset.get('proposal') or raw_asset.get('image_proposal') or {}
    if not isinstance(raw_proposal, dict):
        raw_proposal = {}

    title = _clean_export_visual_text(
        raw_asset.get('title') or raw_proposal.get('title') or 'Generated image',
        160,
    )
    caption = _clean_export_visual_text(
        raw_asset.get('caption')
        or raw_proposal.get('description')
        or raw_proposal.get('context')
        or title,
        240,
    )

    return {
        'message_id': str(raw_asset.get('message_id') or raw_asset.get('id') or '').strip(),
        'data_uri': png_data_uri,
        'content_type': 'image/png',
        'proposal': raw_proposal,
        'title': title or 'Generated image',
        'caption': caption,
        'visual_id': _normalize_export_visual_id(raw_proposal.get('visualId') or raw_proposal.get('visual_id')),
        'prompt': _normalize_export_prompt(raw_proposal.get('prompt')),
    }


def _replace_inline_image_proposal_blocks_with_export_html(
    content: str,
    image_assets: List[Dict[str, Any]],
) -> str:
    rendered_content = str(content or '')
    if INLINE_IMAGE_PROPOSAL_EXPORT_REGEX.search(rendered_content) is None:
        return rendered_content

    used_asset_indexes = set()

    def replace_match(match: re.Match[str]) -> str:
        proposal = _parse_inline_image_proposal_payload(match.group(1) or '')
        image_asset, asset_index = _find_export_image_asset_for_proposal(
            proposal,
            image_assets,
            used_asset_indexes,
        )
        if image_asset is not None and asset_index is not None:
            used_asset_indexes.add(asset_index)
            return _build_export_inline_image_html(image_asset, proposal)

        return _build_missing_export_inline_image_html(proposal)

    return INLINE_IMAGE_PROPOSAL_EXPORT_REGEX.sub(replace_match, rendered_content)


def _parse_inline_image_proposal_payload(payload_text: str) -> Dict[str, Any]:
    payload_json = str(payload_text or '').strip()
    if not payload_json:
        return {}

    try:
        parsed_payload = json.loads(payload_json)
        return parsed_payload if isinstance(parsed_payload, dict) else {}
    except (TypeError, ValueError):
        return {}


def _find_export_image_asset_for_proposal(
    proposal: Dict[str, Any],
    image_assets: List[Dict[str, Any]],
    used_asset_indexes: set,
) -> Tuple[Optional[Dict[str, Any]], Optional[int]]:
    if not image_assets:
        return None, None

    visual_id = _normalize_export_visual_id(proposal.get('visualId') or proposal.get('visual_id'))
    title = _clean_export_visual_text(proposal.get('title'), 160).lower()
    prompt = _normalize_export_prompt(proposal.get('prompt'))

    for index, image_asset in enumerate(image_assets):
        if index in used_asset_indexes:
            continue

        if visual_id and visual_id == image_asset.get('visual_id'):
            return image_asset, index

        asset_proposal = image_asset.get('proposal') if isinstance(image_asset.get('proposal'), dict) else {}
        asset_title = _clean_export_visual_text(
            image_asset.get('title') or asset_proposal.get('title'),
            160,
        ).lower()
        if title and asset_title and title == asset_title:
            return image_asset, index

        asset_prompt = image_asset.get('prompt') or _normalize_export_prompt(asset_proposal.get('prompt'))
        if prompt and asset_prompt and prompt == asset_prompt:
            return image_asset, index

    for index, image_asset in enumerate(image_assets):
        if index not in used_asset_indexes:
            return image_asset, index

    return None, None


def _build_export_inline_image_html(
    image_asset: Dict[str, Any],
    proposal: Dict[str, Any],
) -> str:
    title = _clean_export_visual_text(
        image_asset.get('title') or proposal.get('title') or 'Generated image',
        160,
    ) or 'Generated image'
    caption = _clean_export_visual_text(
        image_asset.get('caption')
        or proposal.get('description')
        or proposal.get('context')
        or title,
        240,
    )
    if title and caption and title != caption:
        caption = _clean_export_visual_text(f'{title}: {caption}', 260)
    elif title and not caption:
        caption = title
    caption_html = ''
    if caption:
        caption_html = (
            '<p class="export-inline-image-caption">'
            f'<em>{_escape_html(caption)}</em>'
            '</p>'
        )

    return (
        '\n\n'
        '<div class="export-inline-image">'
        f'<p><img src="{_escape_html(image_asset.get("data_uri") or "")}" alt="{_escape_html(title)}" /></p>'
        f'{caption_html}'
        '</div>'
        '\n\n'
    )


def _build_missing_export_inline_image_html(proposal: Dict[str, Any]) -> str:
    title = _clean_export_visual_text(proposal.get('title') or 'Image proposal', 160)
    return (
        '\n\n'
        '<p><em>'
        f'Image proposal not generated: {_escape_html(title)}'
        '</em></p>'
        '\n\n'
    )


def _clean_export_visual_text(value: Any, max_chars: int) -> str:
    text = re.sub(r'\s+', ' ', str(value or '')).strip()
    if max_chars and len(text) > max_chars:
        text = text[:max_chars - 3].rstrip()
        if ' ' in text:
            text = text.rsplit(' ', 1)[0]
        text = f'{text}...'
    return text


def _normalize_export_visual_id(value: Any) -> str:
    normalized_value = re.sub(r'[^a-zA-Z0-9_.-]+', '_', str(value or '').strip())
    normalized_value = normalized_value.strip('._-')
    return normalized_value[:120]


def _normalize_export_prompt(value: Any) -> str:
    return str(value or '').replace('\r\n', '\n').replace('\r', '\n').strip()[:4000]


def _message_to_docx_bytes(message: Dict[str, Any]) -> bytes:
    doc = DocxDocument()
    doc.add_heading('Message Export', level=1)

    role_label = _role_to_label(message.get('role', 'unknown'))
    timestamp = message.get('timestamp', '')

    meta_paragraph = doc.add_paragraph()
    meta_run = meta_paragraph.add_run(f"Role: {role_label}")
    meta_run.bold = True
    if timestamp:
        meta_paragraph.add_run(f"    {timestamp}")

    doc.add_paragraph('')

    content = _render_message_export_content(message)
    if content:
        _add_markdown_content_to_doc(doc, content)
    else:
        doc.add_paragraph('No content recorded.')

    citation_labels = _build_message_citation_labels(message)
    if citation_labels:
        doc.add_heading('Citations', level=2)
        for citation_label in citation_labels:
            doc.add_paragraph(citation_label, style='List Bullet')

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.read()


def _message_to_pptx_bytes(
    message: Dict[str, Any],
    settings: Dict[str, Any],
    requested_slide_count: Optional[int] = None,
) -> bytes:
    role_label = _role_to_label(message.get('role', 'unknown'))
    timestamp = str(message.get('timestamp', '') or '')

    render_content = _render_message_export_content(message)
    slide_plan = _build_message_powerpoint_plan(
        content=render_content,
        message=message,
        settings=settings,
        requested_model=_extract_message_powerpoint_model(message),
        requested_slide_count=requested_slide_count,
    )
    if slide_plan.get('include_appendix_slides', True):
        appendix_assets = _extract_powerpoint_appendix_assets(render_content)
        citation_labels = _build_message_citation_labels(message)
    else:
        appendix_assets = {'images': [], 'tables': [], 'code_blocks': []}
        citation_labels = []

    presentation = Presentation()
    presentation.slide_width = PptxInches(13.333)
    presentation.slide_height = PptxInches(7.5)

    presentation_title = slide_plan.get('presentation_title') or f'{role_label} Message'
    presentation_subtitle = slide_plan.get('presentation_subtitle') or _build_powerpoint_subtitle(
        role_label,
        timestamp,
    )

    if slide_plan.get('include_title_slide', True):
        _add_powerpoint_title_slide(
            presentation,
            title=presentation_title,
            subtitle=presentation_subtitle,
            role_label=role_label,
            timestamp=timestamp,
        )

    rendered_slide = False
    for slide_spec in slide_plan.get('slides', []):
        _add_powerpoint_content_slide(presentation, slide_spec, role_label, timestamp)
        rendered_slide = True

    if not rendered_slide:
        _add_powerpoint_content_slide(
            presentation,
            {'title': 'Overview', 'bullets': ['No content recorded.']},
            role_label,
            timestamp,
        )

    _append_powerpoint_appendix_slides(presentation, appendix_assets, citation_labels)

    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer.read()


def _sanitize_powerpoint_source_content(content: str) -> str:
    if not content:
        return ''

    sanitized = POWERPOINT_DATA_URI_PATTERN.sub('[inline-image]', content)
    sanitized = re.sub(
        r'<img[^>]*alt="([^"]*)"[^>]*>',
        lambda match: f"[Image: {match.group(1).strip() or 'Inline visual'}]",
        sanitized,
        flags=re.IGNORECASE,
    )
    sanitized = re.sub(r'<img[^>]*>', '[Image]', sanitized, flags=re.IGNORECASE)
    sanitized = sanitized.replace('</div>', '\n')
    sanitized = sanitized.replace('</section>', '\n')
    sanitized = sanitized.replace('<br>', '\n').replace('<br/>', '\n').replace('<br />', '\n')
    sanitized = re.sub(r'<[^>]+>', ' ', sanitized)
    sanitized = re.sub(r'\n{3,}', '\n\n', sanitized)
    return sanitized.strip()


def _parse_powerpoint_requested_slide_count(raw_slide_count: Any) -> Optional[int]:
    if raw_slide_count is None:
        return None

    if isinstance(raw_slide_count, str):
        stripped_slide_count = raw_slide_count.strip()
        if not stripped_slide_count:
            return None
        if not re.fullmatch(r'\d+', stripped_slide_count):
            raise ValueError(
                f'slide_count must be a whole number between 1 and {POWERPOINT_MAX_SLIDES}.'
            )
        slide_count = int(stripped_slide_count)
    elif isinstance(raw_slide_count, bool):
        raise ValueError(
            f'slide_count must be a whole number between 1 and {POWERPOINT_MAX_SLIDES}.'
        )
    elif isinstance(raw_slide_count, float):
        if not raw_slide_count.is_integer():
            raise ValueError(
                f'slide_count must be a whole number between 1 and {POWERPOINT_MAX_SLIDES}.'
            )
        slide_count = int(raw_slide_count)
    elif isinstance(raw_slide_count, int):
        slide_count = raw_slide_count
    else:
        raise ValueError(
            f'slide_count must be a whole number between 1 and {POWERPOINT_MAX_SLIDES}.'
        )

    if slide_count < 1 or slide_count > POWERPOINT_MAX_SLIDES:
        raise ValueError(
            f'slide_count must be between 1 and {POWERPOINT_MAX_SLIDES}.'
        )

    return slide_count


def _resolve_powerpoint_slide_count(requested_slide_count: Optional[int] = None) -> int:
    if requested_slide_count is None:
        return POWERPOINT_DEFAULT_SLIDES

    return max(1, min(int(requested_slide_count), POWERPOINT_MAX_SLIDES))


def _build_structured_markdown_powerpoint_plan(
    content: str,
    message: Dict[str, Any],
) -> Optional[Dict[str, Any]]:
    slide_sections, preamble_title = _extract_powerpoint_structured_slide_sections(content)
    if len(slide_sections) < 2:
        return None

    role_label = _role_to_label(message.get('role', 'unknown'))
    timestamp = str(message.get('timestamp', '') or '')
    slides: List[Dict[str, Any]] = []
    title_section = slide_sections[0] if _is_powerpoint_title_section(slide_sections[0]) else None
    content_sections = slide_sections[1:] if title_section else slide_sections
    title_metadata = _build_powerpoint_title_metadata_from_section(
        title_section,
        preamble_title,
        role_label,
        timestamp,
    )

    for index, section in enumerate(content_sections[:POWERPOINT_MAX_STRUCTURED_SLIDES], start=1):
        section_content = _trim_powerpoint_structured_section_content(section.get('content', ''))
        title = _clean_slide_text(section.get('title') or f'Slide {index}', 100)
        bullets = _extract_structured_powerpoint_bullets(
            section_content,
            max_bullets=POWERPOINT_MAX_STRUCTURED_BULLETS_PER_SLIDE,
        )
        tables = _extract_structured_powerpoint_tables(section_content)
        images = _extract_structured_powerpoint_images(section_content)
        slides.append({
            'title': title or f'Slide {index}',
            'bullets': bullets,
            'tables': tables,
            'images': images,
            'allow_empty_body': True,
            'bullet_char_limit': POWERPOINT_STRUCTURED_BULLET_CHAR_LIMIT,
            'footer_label': _build_powerpoint_slide_footer_label(section, index),
            'source_format': 'structured_markdown',
        })

    if len(slides) < 1:
        return None

    return {
        'presentation_title': title_metadata['title'],
        'presentation_subtitle': title_metadata['subtitle'],
        'slides': slides,
        'include_title_slide': bool(title_section),
        'include_appendix_slides': False,
        'source_format': 'structured_markdown',
    }


def _extract_powerpoint_structured_slide_sections(content: str) -> Tuple[List[Dict[str, str]], str]:
    numbered_sections, numbered_title = _extract_powerpoint_numbered_slide_sections(content)
    if len(numbered_sections) >= 2:
        return numbered_sections, numbered_title

    separator_sections, separator_title = _extract_powerpoint_separator_slide_sections(content)
    if len(separator_sections) >= 2:
        return separator_sections, separator_title

    return [], ''


def _extract_powerpoint_numbered_slide_sections(content: str) -> Tuple[List[Dict[str, str]], str]:
    sections: List[Dict[str, Any]] = []
    preamble_lines: List[str] = []
    current_section: Optional[Dict[str, Any]] = None
    in_code_block = False

    for line in str(content or '').splitlines():
        stripped_line = line.strip()
        if stripped_line.startswith('```'):
            in_code_block = not in_code_block

        if not in_code_block and _is_powerpoint_slide_separator(line):
            continue

        slide_marker = None if in_code_block else _match_powerpoint_slide_marker(line)
        if slide_marker:
            if current_section:
                sections.append(current_section)

            current_section = {
                'slide_number': slide_marker.get('number', ''),
                'title': slide_marker.get('title', ''),
                'lines': [],
            }
            continue

        if current_section is None:
            preamble_lines.append(line)
        else:
            current_section['lines'].append(line)

    if current_section:
        sections.append(current_section)

    resolved_sections: List[Dict[str, str]] = []
    for index, section in enumerate(sections, start=1):
        title, section_content = _resolve_powerpoint_slide_title_and_content(
            section.get('title', ''),
            '\n'.join(section.get('lines', [])).strip(),
            fallback_title=f'Slide {index}',
        )
        resolved_sections.append({
            'slide_number': str(section.get('slide_number', '') or ''),
            'title': title,
            'content': section_content,
        })

    return resolved_sections, _extract_powerpoint_preamble_title(preamble_lines)


def _extract_powerpoint_separator_slide_sections(content: str) -> Tuple[List[Dict[str, str]], str]:
    blocks: List[List[str]] = []
    current_block: List[str] = []
    in_code_block = False

    for line in str(content or '').splitlines():
        stripped_line = line.strip()
        if stripped_line.startswith('```'):
            in_code_block = not in_code_block

        if not in_code_block and _is_powerpoint_slide_separator(line):
            blocks.append(current_block)
            current_block = []
            continue

        current_block.append(line)

    blocks.append(current_block)

    nonempty_blocks = [
        block for block in blocks
        if any(line.strip() for line in block)
    ]
    if nonempty_blocks and _looks_like_powerpoint_front_matter_block(nonempty_blocks[0]):
        nonempty_blocks = nonempty_blocks[1:]

    sections: List[Dict[str, str]] = []
    for index, block in enumerate(nonempty_blocks, start=1):
        title, section_content = _resolve_powerpoint_slide_title_and_content(
            '',
            '\n'.join(block).strip(),
            fallback_title=f'Slide {index}',
        )
        if title == f'Slide {index}':
            return [], ''
        sections.append({
            'slide_number': str(index),
            'title': title,
            'content': section_content,
        })

    return sections, sections[0]['title'] if sections else ''


def _match_powerpoint_slide_marker(line: str) -> Optional[Dict[str, str]]:
    candidate = str(line or '').strip()
    if not candidate:
        return None

    heading_match = re.match(r'^#{1,6}\s+(.+?)\s*$', candidate)
    if heading_match:
        candidate = heading_match.group(1).strip()

    candidate = candidate.strip('*_` ')
    marker_match = re.match(
        r'^(?:slide|page)\s+(\d{1,3})(?:\s*(?:[:.)-]|\u2013|\u2014)\s*(.+?))?$',
        candidate,
        flags=re.IGNORECASE,
    )
    if not marker_match:
        return None

    return {
        'number': marker_match.group(1),
        'title': _clean_slide_text(marker_match.group(2) or '', 100),
    }


def _resolve_powerpoint_slide_title_and_content(
    marker_title: str,
    content: str,
    fallback_title: str,
) -> Tuple[str, str]:
    labeled_title, content_without_labeled_title = _extract_powerpoint_labeled_title(content)
    if marker_title:
        resolved_marker_title = _clean_slide_text(marker_title, 100)
        if labeled_title and _should_prefer_labeled_powerpoint_title(resolved_marker_title):
            return _clean_slide_text(labeled_title, 100), content_without_labeled_title
        return resolved_marker_title, content_without_labeled_title

    if labeled_title:
        return _clean_slide_text(labeled_title, 100) or fallback_title, content_without_labeled_title

    content_lines = content_without_labeled_title.splitlines()
    for index, line in enumerate(content_lines):
        stripped_line = line.strip()
        if not stripped_line:
            continue

        heading_match = re.match(r'^#{1,6}\s+(.+?)\s*$', stripped_line)
        if heading_match and not _match_powerpoint_slide_marker(stripped_line):
            remaining_lines = content_lines[:index] + content_lines[index + 1:]
            return (
                _clean_slide_text(heading_match.group(1), 100) or fallback_title,
                '\n'.join(remaining_lines).strip(),
            )
        break

    return fallback_title, content_without_labeled_title


def _extract_powerpoint_labeled_title(content: str) -> Tuple[str, str]:
    content_lines = str(content or '').splitlines()
    labeled_title = ''
    retained_lines: List[str] = []
    removed_title = False

    for line in content_lines:
        label_name, label_value = _parse_powerpoint_structured_label_line(line)
        if label_name == 'title' and not removed_title:
            labeled_title = _clean_slide_text(label_value, 120)
            removed_title = True
            continue
        retained_lines.append(line)

    return labeled_title, '\n'.join(retained_lines).strip()


def _parse_powerpoint_structured_label_line(line: str) -> Tuple[str, str]:
    match = re.match(
        r'^\s*(title|subtitle|bullet\s*points?|speaker\s*notes?|visual|image|chart)\s*:\s*(.*?)\s*$',
        str(line or '').strip(),
        flags=re.IGNORECASE,
    )
    if not match:
        return '', ''

    normalized_label = re.sub(r'\s+', ' ', match.group(1).strip().lower())
    return normalized_label, match.group(2).strip()


def _should_prefer_labeled_powerpoint_title(marker_title: str) -> bool:
    normalized_title = _clean_slide_text(marker_title, 100).lower()
    if not normalized_title:
        return True

    return normalized_title in {
        'agenda',
        'intro',
        'introduction',
        'overview',
        'section',
        'title',
        'title slide',
    }


def _extract_powerpoint_preamble_title(lines: List[str]) -> str:
    for line in lines:
        stripped_line = line.strip()
        if not stripped_line:
            continue

        heading_match = re.match(r'^#{1,6}\s+(.+?)\s*$', stripped_line)
        if heading_match and not _match_powerpoint_slide_marker(stripped_line):
            return _clean_slide_text(heading_match.group(1), 100)

        cleaned_line = _clean_slide_text(stripped_line, 100)
        if cleaned_line:
            return cleaned_line

    return ''


def _is_powerpoint_slide_separator(line: str) -> bool:
    return bool(re.match(r'^\s*(?:---+|\*\*\*+|___+)\s*$', str(line or '')))


def _looks_like_powerpoint_front_matter_block(block: List[str]) -> bool:
    nonempty_lines = [line.strip() for line in block if line.strip()]
    if not nonempty_lines or len(nonempty_lines) > 12:
        return False

    return all(re.match(r'^[A-Za-z0-9_-]+\s*:', line) for line in nonempty_lines)


def _is_powerpoint_title_section(section: Dict[str, str]) -> bool:
    title = _clean_slide_text(section.get('title', ''), 80).lower()
    slide_number = str(section.get('slide_number', '') or '').strip()
    return slide_number == '1' and title in {'title', 'title slide', 'intro', 'introduction'}


def _build_powerpoint_title_metadata_from_section(
    title_section: Optional[Dict[str, str]],
    preamble_title: str,
    role_label: str,
    timestamp: str,
) -> Dict[str, str]:
    default_subtitle = _build_powerpoint_subtitle(role_label, timestamp)
    if not title_section:
        return {
            'title': _clean_slide_text(preamble_title, 100) or f'{role_label} Message',
            'subtitle': default_subtitle,
        }

    title_lines = _extract_powerpoint_title_section_lines(title_section.get('content', ''))
    title = _clean_slide_text(title_lines[0], 120) if title_lines else ''
    if not title:
        title = _clean_slide_text(preamble_title, 120) or f'{role_label} Message'

    subtitle_lines = [
        _clean_slide_text(line, 120)
        for line in title_lines[1:4]
        if _clean_slide_text(line, 120)
    ]
    subtitle = ' | '.join(subtitle_lines) if subtitle_lines else default_subtitle

    return {
        'title': _clean_slide_text(title, 120),
        'subtitle': _clean_slide_text(subtitle, 180),
    }


def _extract_powerpoint_title_section_lines(content: str) -> List[str]:
    title_lines: List[str] = []
    for line in _trim_powerpoint_structured_section_content(content).splitlines():
        stripped_line = line.strip()
        if not stripped_line:
            continue
        if _is_powerpoint_slide_separator(stripped_line):
            continue
        if _looks_like_markdown_table_row(stripped_line) or _looks_like_markdown_table_divider(stripped_line):
            continue
        if _line_contains_powerpoint_inline_visual(stripped_line):
            continue

        label_name, label_value = _parse_powerpoint_structured_label_line(stripped_line)
        if label_name in {'title', 'bullet points', 'bullet point', 'visual', 'image', 'chart'}:
            if label_name == 'title' and label_value:
                title_lines.append(label_value)
            continue
        if label_name == 'subtitle':
            if label_value:
                title_lines.append(label_value)
            continue
        if label_name in {'speaker note', 'speaker notes'}:
            continue

        title_lines.append(stripped_line)

    return title_lines


def _line_contains_powerpoint_inline_visual(line: str) -> bool:
    normalized_line = str(line or '').strip().lower()
    if not normalized_line:
        return False

    return (
        '<img' in normalized_line
        or 'export-inline-chart' in normalized_line
        or 'export-inline-image' in normalized_line
        or normalized_line.startswith('![')
    )


def _build_powerpoint_slide_footer_label(section: Dict[str, str], fallback_index: int) -> str:
    slide_number = str(section.get('slide_number', '') or '').strip()
    if slide_number:
        return f'Slide {slide_number}'

    return f'Slide {fallback_index}'


def _trim_powerpoint_structured_section_content(content: str) -> str:
    retained_lines: List[str] = []
    for line in str(content or '').splitlines():
        stripped_line = line.strip()
        if _is_powerpoint_non_slide_tail_marker(stripped_line):
            break
        retained_lines.append(line)

    return '\n'.join(retained_lines).strip()


def _is_powerpoint_non_slide_tail_marker(line: str) -> bool:
    if not line:
        return False

    normalized_line = _clean_slide_text(line, 140).lower()
    if normalized_line.startswith('if you want') or normalized_line.startswith('would you like'):
        return True

    heading_match = re.match(r'^#{1,6}\s+(.+?)\s*$', line)
    if not heading_match:
        return False

    heading_text = _clean_slide_text(heading_match.group(1), 140).lower()
    return heading_text.startswith((
        'optional speaker notes',
        'speaker notes',
        'coverage caveat',
        'coverage notes',
    ))


def _extract_structured_powerpoint_bullets(content: str, max_bullets: int) -> List[str]:
    bullets: List[str] = []
    paragraph_lines: List[str] = []
    in_code_block = False

    def flush_paragraph_lines():
        nonlocal paragraph_lines
        if not paragraph_lines or len(bullets) >= max_bullets:
            paragraph_lines = []
            return

        paragraph_text = ' '.join(paragraph_lines)
        for paragraph_bullet in _sentence_bullets(paragraph_text, max_bullets - len(bullets)):
            cleaned_bullet = _clean_slide_text(paragraph_bullet, POWERPOINT_STRUCTURED_BULLET_CHAR_LIMIT)
            if cleaned_bullet and cleaned_bullet not in bullets:
                bullets.append(cleaned_bullet)
            if len(bullets) >= max_bullets:
                break

        paragraph_lines = []

    for line in str(content or '').splitlines():
        stripped_line = line.strip()
        if stripped_line.startswith('```'):
            flush_paragraph_lines()
            in_code_block = not in_code_block
            continue

        if in_code_block:
            continue

        if not stripped_line:
            flush_paragraph_lines()
            continue

        if _is_powerpoint_slide_separator(stripped_line):
            flush_paragraph_lines()
            continue

        if _looks_like_markdown_table_row(stripped_line) or _looks_like_markdown_table_divider(stripped_line):
            flush_paragraph_lines()
            continue

        if _line_contains_powerpoint_inline_visual(stripped_line):
            flush_paragraph_lines()
            continue

        label_name, label_value = _parse_powerpoint_structured_label_line(stripped_line)
        if label_name:
            flush_paragraph_lines()
            if label_name in {'speaker note', 'speaker notes'} and label_value:
                cleaned_label_value = _clean_slide_text(
                    label_value,
                    POWERPOINT_STRUCTURED_BULLET_CHAR_LIMIT,
                )
                if cleaned_label_value and cleaned_label_value not in bullets:
                    bullets.append(cleaned_label_value)
            continue

        heading_match = re.match(r'^#{1,6}\s+(.+?)\s*$', stripped_line)
        if heading_match:
            flush_paragraph_lines()
            cleaned_heading = _clean_slide_text(heading_match.group(1), POWERPOINT_STRUCTURED_BULLET_CHAR_LIMIT)
            if cleaned_heading and cleaned_heading not in bullets:
                bullets.append(cleaned_heading)
            if len(bullets) >= max_bullets:
                break
            continue

        bullet_match = re.match(r'^(?:[-*+]\s+|\d+[.)]\s+)(.+)$', stripped_line)
        if bullet_match:
            flush_paragraph_lines()
            cleaned_bullet = _clean_slide_text(bullet_match.group(1), POWERPOINT_STRUCTURED_BULLET_CHAR_LIMIT)
            if cleaned_bullet and cleaned_bullet not in bullets:
                bullets.append(cleaned_bullet)
            if len(bullets) >= max_bullets:
                break
            continue

        if re.match(r'^!\[[^\]]*\]\([^)]+\)$', stripped_line):
            flush_paragraph_lines()
            continue

        paragraph_lines.append(stripped_line)
        if len(' '.join(paragraph_lines)) >= POWERPOINT_STRUCTURED_BULLET_CHAR_LIMIT:
            flush_paragraph_lines()

    flush_paragraph_lines()
    return bullets[:max_bullets]


def _extract_structured_powerpoint_tables(content: str) -> List[Dict[str, Any]]:
    tables: List[Dict[str, Any]] = []
    table_block: List[str] = []
    in_code_block = False

    def flush_table_block():
        nonlocal table_block
        if table_block:
            parsed_table = _parse_powerpoint_markdown_table_block(table_block)
            if parsed_table:
                tables.append(parsed_table)
            table_block = []

    for line in str(content or '').splitlines():
        stripped_line = line.strip()
        if stripped_line.startswith('```'):
            in_code_block = not in_code_block
            flush_table_block()
            continue

        if in_code_block:
            continue

        if _looks_like_markdown_table_row(stripped_line) or _looks_like_markdown_table_divider(stripped_line):
            table_block.append(stripped_line)
            continue

        flush_table_block()

    flush_table_block()
    return tables[:2]


def _extract_structured_powerpoint_images(content: str) -> List[Dict[str, Any]]:
    if not str(content or '').strip():
        return []

    html = markdown2.markdown(content, extras=DOCX_MARKDOWN_EXTRAS)
    soup = BeautifulSoup(f'<div>{html}</div>', 'html.parser')
    root = soup.div if soup.div else soup
    images: List[Dict[str, Any]] = []
    seen_keys = set()

    for image_node in root.find_all('img'):
        image_bytes = decode_base64_image_data_uri(image_node.get('src'))
        if not image_bytes:
            continue

        image_key = (len(image_bytes), image_bytes[:24])
        if image_key in seen_keys:
            continue
        seen_keys.add(image_key)

        image_wrapper = image_node.find_parent(class_='export-inline-image')
        chart_wrapper = image_node.find_parent(class_='export-inline-chart')
        wrapper = image_wrapper or chart_wrapper
        caption_node = None
        if wrapper:
            caption_node = wrapper.find(class_='export-inline-image-caption')
            if caption_node is None:
                caption_node = wrapper.find(class_='export-inline-chart-caption')
        caption = (
            caption_node.get_text(' ', strip=True)
            if caption_node and caption_node.get_text(' ', strip=True)
            else image_node.get('alt') or 'Inline visual'
        )

        images.append({
            'title': _clean_slide_text(image_node.get('alt') or f'Visual {len(images) + 1}', 100),
            'caption': _clean_slide_text(caption, 160),
            'image_bytes': image_bytes,
        })
        if len(images) >= POWERPOINT_MAX_INLINE_IMAGES_PER_SLIDE:
            break

    return images


def _parse_powerpoint_markdown_table_block(table_block: List[str]) -> Optional[Dict[str, Any]]:
    split_rows = [
        _split_powerpoint_markdown_table_line(line)
        for line in table_block
        if _looks_like_markdown_table_row(line) or _looks_like_markdown_table_divider(line)
    ]
    split_rows = [row for row in split_rows if row]
    if len(split_rows) < 2:
        return None

    separator_index = next(
        (
            index for index, row in enumerate(split_rows)
            if _is_powerpoint_markdown_table_separator_row(row)
        ),
        None,
    )
    has_header = separator_index is not None and separator_index > 0

    if has_header:
        rows = [split_rows[separator_index - 1]] + split_rows[separator_index + 1:]
    else:
        rows = [row for row in split_rows if not _is_powerpoint_markdown_table_separator_row(row)]

    rows = [
        [_clean_powerpoint_markdown_table_cell(cell) for cell in row[:POWERPOINT_MAX_TABLE_COLS]]
        for row in rows[:POWERPOINT_MAX_TABLE_ROWS]
    ]
    rows = [row for row in rows if any(cell for cell in row)]
    if len(rows) < 2:
        return None

    column_count = max(len(row) for row in rows)
    normalized_rows = [row + [''] * (column_count - len(row)) for row in rows]

    return {
        'rows': normalized_rows,
        'has_header': has_header,
    }


def _split_powerpoint_markdown_table_line(line: str) -> List[str]:
    candidate = str(line or '').strip()
    if candidate.startswith('|'):
        candidate = candidate[1:]
    if candidate.endswith('|'):
        candidate = candidate[:-1]

    return [cell.strip() for cell in candidate.split('|')]


def _is_powerpoint_markdown_table_separator_row(row: List[str]) -> bool:
    return bool(row) and all(re.match(r'^\s*:?-{3,}:?\s*$', cell or '') for cell in row)


def _clean_powerpoint_markdown_table_cell(cell: Any) -> str:
    return _clean_slide_text(str(cell or '').replace('<br>', ' '), 80)


def _build_message_powerpoint_plan(
    content: str,
    message: Dict[str, Any],
    settings: Dict[str, Any],
    requested_model: str = '',
    requested_slide_count: Optional[int] = None,
) -> Dict[str, Any]:
    structured_plan = _build_structured_markdown_powerpoint_plan(content, message)
    if structured_plan and requested_slide_count is None:
        return structured_plan

    planning_content = _sanitize_powerpoint_source_content(content)
    target_slide_count = _resolve_powerpoint_slide_count(requested_slide_count)
    fallback_plan = _build_fallback_powerpoint_plan(
        planning_content,
        message,
        slide_count=target_slide_count,
    )
    if not planning_content.strip():
        return fallback_plan

    ai_plan = _generate_powerpoint_slide_plan_with_model(
        content=planning_content,
        message=message,
        settings=settings,
        requested_model=requested_model,
        fallback_plan=fallback_plan,
        target_slide_count=target_slide_count,
    )
    return ai_plan or fallback_plan


def _build_fallback_powerpoint_plan(
    content: str,
    message: Dict[str, Any],
    slide_count: Optional[int] = None,
) -> Dict[str, Any]:
    role_label = _role_to_label(message.get('role', 'unknown'))
    timestamp = str(message.get('timestamp', '') or '')
    target_slide_count = _resolve_powerpoint_slide_count(slide_count)
    sections = _extract_powerpoint_sections(content)

    if len(sections) == 1 and not sections[0].get('title'):
        chunks = [chunk.strip() for chunk in re.split(r'\n\s*\n', sections[0].get('content', '')) if chunk.strip()]
        if len(chunks) > 1:
            sections = [
                {
                    'title': 'Overview' if index == 0 else f'Detail {index + 1}',
                    'content': chunk,
                }
                for index, chunk in enumerate(chunks)
            ]

    slides: List[Dict[str, Any]] = []
    for index, section in enumerate(sections, start=1):
        if len(slides) >= target_slide_count:
            break

        title = _clean_slide_text(
            section.get('title') or ('Overview' if index == 1 else f'Key Point {index}'),
            80,
        )
        bullets = _extract_powerpoint_bullets(
            section.get('content', ''),
            max_bullets=POWERPOINT_MAX_BULLETS_PER_SLIDE,
        )
        if not bullets:
            plain_text = _markdown_to_plain_text(section.get('content', ''))
            bullets = _sentence_bullets(plain_text, POWERPOINT_MAX_BULLETS_PER_SLIDE)

        if not bullets:
            continue

        slides.append({
            'title': title,
            'bullets': bullets,
        })

    if not slides:
        slides.append({
            'title': 'Overview',
            'bullets': _sentence_bullets(
                _markdown_to_plain_text(content),
                POWERPOINT_MAX_BULLETS_PER_SLIDE,
            ) or ['No content recorded.'],
        })

    return {
        'presentation_title': _derive_powerpoint_title(content, sections, role_label),
        'presentation_subtitle': _build_powerpoint_subtitle(role_label, timestamp),
        'slides': slides,
    }


def _extract_message_powerpoint_model(message: Dict[str, Any]) -> str:
    metadata = message.get('metadata') if isinstance(message.get('metadata'), dict) else {}
    candidates = [
        message.get('model_deployment_name'),
        metadata.get('selected_model'),
        metadata.get('model_deployment_name'),
        metadata.get('model'),
    ]

    for candidate in candidates:
        normalized_candidate = _normalize_powerpoint_model_candidate(candidate)
        if normalized_candidate:
            return normalized_candidate

    return ''


def _normalize_powerpoint_model_candidate(candidate: Any) -> str:
    if isinstance(candidate, str):
        return candidate.strip()

    if isinstance(candidate, dict):
        for key in ('deploymentName', 'deployment', 'value', 'name'):
            value = str(candidate.get(key) or '').strip()
            if value:
                return value
        return ''

    if isinstance(candidate, list):
        for item in candidate:
            normalized_candidate = _normalize_powerpoint_model_candidate(item)
            if normalized_candidate:
                return normalized_candidate

    return ''


def _generate_powerpoint_slide_plan_with_model(
    content: str,
    message: Dict[str, Any],
    settings: Dict[str, Any],
    requested_model: str,
    fallback_plan: Dict[str, Any],
    target_slide_count: int,
) -> Optional[Dict[str, Any]]:
    prompt_source = str(content or '').strip()[:POWERPOINT_PLAN_SOURCE_CHAR_LIMIT]
    if not prompt_source:
        return None

    role_label = _role_to_label(message.get('role', 'unknown'))
    timestamp = str(message.get('timestamp', '') or '')
    fallback_seed_lines = []
    for slide in fallback_plan.get('slides', [])[:3]:
        bullet_preview = '; '.join(slide.get('bullets', [])[:2])
        seed_text = bullet_preview or slide.get('body', '')
        if seed_text:
            fallback_seed_lines.append(f"- {slide.get('title', 'Slide')}: {seed_text}")

    try:
        gpt_client, gpt_model = _initialize_gpt_client(settings, requested_model)
        model_lower = gpt_model.lower()
        is_reasoning_model = (
            'o1' in model_lower or 'o3' in model_lower or 'gpt-5' in model_lower
        )
        instruction_role = 'developer' if is_reasoning_model else 'system'
        slide_prompt = (
            'You are turning a single chat message into a presentation-ready PowerPoint outline. '
            'Preserve factual content, keep numbers and sequence accurate, and do not invent new claims. '
            'Optimize for concise slide titles and visually scannable bullets. '
            'Return valid JSON only with the keys presentation_title, presentation_subtitle, and slides. '
            f'The slides value must be an array of 1 to {target_slide_count} objects. '
            f'Use {target_slide_count} slides when the source has enough distinct material; use fewer only to avoid invented content. '
            'Each slide object must contain title and bullets. It may also contain body. '
            'Use no more than 5 bullets per slide. Keep each bullet under 16 words. '
            'If body is present, keep it under 280 characters. '
            'Do not wrap the JSON in markdown fences.'
        )

        user_prompt = '\n'.join([
            f'Role: {role_label}',
            f'Timestamp: {timestamp or "Unknown"}',
            'Fallback outline seed:',
            '\n'.join(fallback_seed_lines) if fallback_seed_lines else '- Overview: Summarize the message clearly.',
            '',
            'Source message:',
            prompt_source,
        ])

        slide_response = gpt_client.chat.completions.create(
            model=gpt_model,
            messages=[
                {
                    'role': instruction_role,
                    'content': slide_prompt,
                },
                {
                    'role': 'user',
                    'content': user_prompt,
                },
            ]
        )

        raw_plan = (
            (slide_response.choices[0].message.content or '').strip()
            if slide_response.choices else ''
        )
        json_payload = _extract_json_object(raw_plan)
        if not json_payload:
            return None

        parsed_plan = json.loads(json_payload)
        return _sanitize_powerpoint_plan(
            parsed_plan,
            fallback_plan,
            max_slides=target_slide_count,
        )
    except Exception as exc:
        debug_print(f'Message PowerPoint plan generation failed: {exc}')
        log_event(
            'Message PowerPoint plan generation failed',
            extra={
                'requested_model': requested_model or None,
                'content_length': len(prompt_source),
            },
            level='WARNING'
        )
        return None


def _extract_json_object(raw_content: str) -> Optional[str]:
    if not raw_content:
        return None

    start_index = raw_content.find('{')
    end_index = raw_content.rfind('}')
    if start_index == -1 or end_index == -1 or end_index <= start_index:
        return None

    return raw_content[start_index:end_index + 1]


def _sanitize_powerpoint_plan(
    plan: Any,
    fallback_plan: Dict[str, Any],
    max_slides: Optional[int] = None,
) -> Optional[Dict[str, Any]]:
    if not isinstance(plan, dict):
        return None

    slide_limit = _resolve_powerpoint_slide_count(max_slides)
    slides: List[Dict[str, Any]] = []
    raw_slides = plan.get('slides') if isinstance(plan.get('slides'), list) else []
    for index, raw_slide in enumerate(raw_slides, start=1):
        if len(slides) >= slide_limit or not isinstance(raw_slide, dict):
            continue

        title = _clean_slide_text(raw_slide.get('title') or f'Slide {index}', 80)
        raw_bullets = raw_slide.get('bullets', [])
        if isinstance(raw_bullets, str):
            raw_bullets = [raw_bullets]

        bullets = []
        if isinstance(raw_bullets, list):
            for raw_bullet in raw_bullets:
                cleaned_bullet = _clean_slide_text(raw_bullet, 120)
                if cleaned_bullet and cleaned_bullet not in bullets:
                    bullets.append(cleaned_bullet)
                if len(bullets) >= POWERPOINT_MAX_BULLETS_PER_SLIDE:
                    break

        body = _clean_slide_text(raw_slide.get('body', ''), 280)
        if not bullets and body:
            bullets = _sentence_bullets(body, POWERPOINT_MAX_BULLETS_PER_SLIDE) or [body]

        if not bullets:
            continue

        slides.append({
            'title': title,
            'bullets': bullets[:POWERPOINT_MAX_BULLETS_PER_SLIDE],
            'bullet_char_limit': POWERPOINT_BULLET_CHAR_LIMIT,
        })

    if not slides:
        return None

    return {
        'presentation_title': _clean_slide_text(
            plan.get('presentation_title') or fallback_plan.get('presentation_title') or 'Message Export',
            100,
        ),
        'presentation_subtitle': _clean_slide_text(
            plan.get('presentation_subtitle') or fallback_plan.get('presentation_subtitle') or '',
            140,
        ),
        'slides': slides,
    }


def _extract_powerpoint_sections(content: str) -> List[Dict[str, str]]:
    if not content.strip():
        return [{'title': '', 'content': ''}]

    sections: List[Dict[str, str]] = []
    current_title = ''
    current_lines: List[str] = []
    heading_pattern = re.compile(r'^\s*#{1,6}\s+(.+?)\s*$')

    for line in content.splitlines():
        heading_match = heading_pattern.match(line)
        if heading_match:
            if current_title or current_lines:
                sections.append({
                    'title': current_title,
                    'content': '\n'.join(current_lines).strip(),
                })
            current_title = _clean_slide_text(heading_match.group(1), 80)
            current_lines = []
            continue

        current_lines.append(line)

    if current_title or current_lines:
        sections.append({
            'title': current_title,
            'content': '\n'.join(current_lines).strip(),
        })

    return sections or [{'title': '', 'content': content.strip()}]


def _extract_powerpoint_bullets(content: str, max_bullets: int) -> List[str]:
    bullets: List[str] = []
    in_code_block = False

    for line in str(content or '').splitlines():
        stripped_line = line.strip()
        if not stripped_line:
            continue

        if stripped_line.startswith('```'):
            in_code_block = not in_code_block
            continue
        if in_code_block:
            continue

        if _looks_like_markdown_table_row(stripped_line) or _looks_like_markdown_table_divider(stripped_line):
            continue

        bullet_match = re.match(r'^(?:[-*+]\s+|\d+[.)]\s+)(.+)$', stripped_line)
        if bullet_match:
            cleaned_bullet = _clean_slide_text(bullet_match.group(1), 120)
            if cleaned_bullet and cleaned_bullet not in bullets:
                bullets.append(cleaned_bullet)
            if len(bullets) >= max_bullets:
                break
            continue

        line_bullets = _sentence_bullets(stripped_line, max_bullets - len(bullets))
        for line_bullet in line_bullets:
            if line_bullet and line_bullet not in bullets:
                bullets.append(line_bullet)
            if len(bullets) >= max_bullets:
                break

        if len(bullets) >= max_bullets:
            break

    return bullets[:max_bullets]


def _sentence_bullets(text: str, max_bullets: int) -> List[str]:
    normalized_text = re.sub(r'\s+', ' ', str(text or '')).strip()
    if not normalized_text:
        return []

    sentence_candidates = re.split(r'(?<=[.!?])\s+|\s*;\s+', normalized_text)
    bullets = []
    for sentence in sentence_candidates:
        cleaned_sentence = _clean_slide_text(sentence, 120)
        if not cleaned_sentence:
            continue
        bullets.append(cleaned_sentence)
        if len(bullets) >= max_bullets:
            break

    if not bullets:
        bullets.append(_clean_slide_text(normalized_text, 120))

    return bullets[:max_bullets]


def _looks_like_markdown_table_row(line: str) -> bool:
    return line.startswith('|') and line.endswith('|') and line.count('|') >= 2


def _looks_like_markdown_table_divider(line: str) -> bool:
    return bool(re.match(r'^\|?[\s:-]+\|[\s|:-]*$', line))


def _markdown_to_plain_text(content: str) -> str:
    if not content:
        return ''

    html = markdown2.markdown(content, extras=DOCX_MARKDOWN_EXTRAS)
    soup = BeautifulSoup(f'<div>{html}</div>', 'html.parser')
    plain_text = soup.get_text('\n')
    plain_text = re.sub(r'\n{3,}', '\n\n', plain_text)
    return plain_text.strip()


def _derive_powerpoint_title(content: str, sections: List[Dict[str, str]], role_label: str) -> str:
    for section in sections:
        if section.get('title'):
            return _clean_slide_text(section['title'], 100)

    for line in str(content or '').splitlines():
        cleaned_line = _clean_slide_text(line, 100)
        if cleaned_line:
            return cleaned_line

    return f'{role_label} Message'


def _build_powerpoint_subtitle(role_label: str, timestamp: str) -> str:
    subtitle_parts = [role_label, 'Generated from chat message']
    if timestamp:
        subtitle_parts.append(timestamp)
    return ' | '.join(subtitle_parts)


def _clean_slide_text(value: Any, max_chars: int) -> str:
    text = str(value or '')
    text = re.sub(r'\[([^\]]+)\]\(([^)]+)\)', r'\1', text)
    text = re.sub(r'[`*_~]+', '', text)
    text = re.sub(r'^#{1,6}\s+', '', text)
    text = re.sub(r'^>\s*', '', text)
    text = re.sub(r'\s+', ' ', text).strip()

    if max_chars and len(text) > max_chars:
        truncated_text = text[:max_chars - 3].rstrip()
        if ' ' in truncated_text:
            truncated_text = truncated_text.rsplit(' ', 1)[0]
        text = f'{truncated_text}...'

    return text


def _extract_powerpoint_appendix_assets(content: str) -> Dict[str, List[Dict[str, Any]]]:
    if not content.strip():
        return {'images': [], 'tables': [], 'code_blocks': []}

    html = markdown2.markdown(content, extras=DOCX_MARKDOWN_EXTRAS)
    soup = BeautifulSoup(f'<div>{html}</div>', 'html.parser')
    root = soup.div if soup.div else soup

    return {
        'images': _extract_powerpoint_images(root),
        'tables': _extract_powerpoint_tables(root),
        'code_blocks': _extract_powerpoint_code_blocks(root),
    }


def _extract_powerpoint_images(root: Tag) -> List[Dict[str, Any]]:
    images: List[Dict[str, Any]] = []
    seen_keys = set()

    for image_node in root.find_all('img'):
        image_bytes = decode_base64_image_data_uri(image_node.get('src'))
        if not image_bytes:
            continue

        image_key = (len(image_bytes), image_bytes[:24])
        if image_key in seen_keys:
            continue
        seen_keys.add(image_key)

        chart_wrapper = image_node.find_parent(class_='export-inline-chart')
        caption_node = chart_wrapper.find(class_='export-inline-chart-caption') if chart_wrapper else None
        caption = (
            caption_node.get_text(' ', strip=True)
            if caption_node and caption_node.get_text(' ', strip=True)
            else image_node.get('alt') or 'Inline visual'
        )

        images.append({
            'title': f'Visual {len(images) + 1}',
            'caption': _clean_slide_text(caption, 120),
            'image_bytes': image_bytes,
        })
        if len(images) >= POWERPOINT_MAX_APPENDIX_IMAGES:
            break

    return images


def _extract_powerpoint_tables(root: Tag) -> List[Dict[str, Any]]:
    tables: List[Dict[str, Any]] = []

    for index, table_node in enumerate(root.find_all('table'), start=1):
        parsed_rows: List[List[str]] = []
        header_present = False
        for row_index, row in enumerate(table_node.find_all('tr')):
            cells = row.find_all(['th', 'td'], recursive=False)
            if not cells:
                continue
            if row_index == 0 and all(cell.name.lower() == 'th' for cell in cells):
                header_present = True

            parsed_rows.append([
                _clean_slide_text(cell.get_text(' ', strip=True), 60)
                for cell in cells[:POWERPOINT_MAX_TABLE_COLS]
            ])
            if len(parsed_rows) >= POWERPOINT_MAX_TABLE_ROWS:
                break

        if not parsed_rows:
            continue

        column_count = max(len(row) for row in parsed_rows)
        normalized_rows = [row + [''] * (column_count - len(row)) for row in parsed_rows]
        tables.append({
            'title': f'Table {index}',
            'rows': normalized_rows,
            'has_header': header_present,
        })
        if len(tables) >= POWERPOINT_MAX_APPENDIX_TABLES:
            break

    return tables


def _extract_powerpoint_code_blocks(root: Tag) -> List[Dict[str, Any]]:
    code_blocks: List[Dict[str, Any]] = []

    for index, code_node in enumerate(root.find_all('pre'), start=1):
        code_text = code_node.get_text().rstrip()
        if not code_text:
            continue

        code_lines = code_text.splitlines()[:18]
        code_blocks.append({
            'title': f'Code Example {index}',
            'code': '\n'.join(code_lines),
        })
        if len(code_blocks) >= POWERPOINT_MAX_APPENDIX_CODE_BLOCKS:
            break

    return code_blocks


def _add_powerpoint_title_slide(
    presentation: Presentation,
    title: str,
    subtitle: str,
    role_label: str,
    timestamp: str,
):
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    _apply_powerpoint_background(slide, POWERPOINT_TITLE_BG)

    accent_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        presentation.slide_height - PptxInches(0.35),
        presentation.slide_width,
        PptxInches(0.35),
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = POWERPOINT_ACCENT
    accent_shape.line.fill.background()

    title_shape = slide.shapes.title
    title_shape.text = title
    title_frame = title_shape.text_frame
    title_frame.word_wrap = True
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = PptxPt(30)
            run.font.bold = True
            run.font.color.rgb = POWERPOINT_TITLE_TEXT

    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.word_wrap = True
    for paragraph in subtitle_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = PptxPt(18)
            run.font.color.rgb = RGBColor(219, 234, 254)

    footer_box = slide.shapes.add_textbox(
        PptxInches(0.75),
        presentation.slide_height - PptxInches(0.9),
        PptxInches(5.5),
        PptxInches(0.25),
    )
    footer_frame = footer_box.text_frame
    footer_frame.text = f'{role_label} export{f" | {timestamp}" if timestamp else ""}'
    footer_paragraph = footer_frame.paragraphs[0]
    footer_paragraph.alignment = PP_ALIGN.LEFT
    for run in footer_paragraph.runs:
        run.font.size = PptxPt(10)
        run.font.color.rgb = RGBColor(191, 219, 254)


def _add_powerpoint_content_slide(
    presentation: Presentation,
    slide_spec: Dict[str, Any],
    role_label: str,
    timestamp: str,
):
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    _apply_powerpoint_background(slide, POWERPOINT_BG)

    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        PptxInches(0.22),
        presentation.slide_height,
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = POWERPOINT_ACCENT
    accent_bar.line.fill.background()

    title_shape = slide.shapes.title
    title_shape.text = _clean_slide_text(slide_spec.get('title') or 'Overview', 80)
    title_frame = title_shape.text_frame
    title_frame.word_wrap = True
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = PptxPt(24)
            run.font.bold = True
            run.font.color.rgb = POWERPOINT_TEXT

    tables = slide_spec.get('tables', []) if isinstance(slide_spec.get('tables'), list) else []
    primary_table = next(
        (table for table in tables if isinstance(table, dict) and table.get('rows')),
        None,
    )
    images = slide_spec.get('images', []) if isinstance(slide_spec.get('images'), list) else []
    primary_image = next(
        (image for image in images if isinstance(image, dict) and image.get('image_bytes')),
        None,
    )

    bullets = slide_spec.get('bullets', [])
    if isinstance(bullets, str):
        bullets = [bullets]
    if not bullets:
        bullets = [] if slide_spec.get('allow_empty_body') else ['No content recorded.']

    content_placeholder = slide.placeholders[1]
    content_left = PptxInches(0.78)
    content_top = PptxInches(1.35)
    content_width = presentation.slide_width - PptxInches(1.55)
    content_height = presentation.slide_height - content_top - PptxInches(0.95)

    if primary_image and (bullets or primary_table):
        visual_left = presentation.slide_width - PptxInches(5.25)
        visual_top = PptxInches(1.35)
        visual_width = PptxInches(4.65)
        visual_height = presentation.slide_height - PptxInches(2.15)
        _add_powerpoint_inline_image_to_slide(
            slide,
            primary_image,
            left=visual_left,
            top=visual_top,
            max_width=visual_width,
            max_height=visual_height,
        )
        content_width = visual_left - content_left - PptxInches(0.3)
    elif primary_image:
        _add_powerpoint_inline_image_to_slide(
            slide,
            primary_image,
            left=PptxInches(0.85),
            top=PptxInches(1.25),
            max_width=presentation.slide_width - PptxInches(1.7),
            max_height=presentation.slide_height - PptxInches(2.15),
        )
        content_placeholder.left = PptxInches(0.1)
        content_placeholder.top = PptxInches(0.1)
        content_placeholder.width = PptxInches(0.1)
        content_placeholder.height = PptxInches(0.1)
        content_placeholder.text_frame.clear()
        bullets = []

    if primary_table:
        table_rows = primary_table.get('rows', [])
        table_height = min(
            PptxInches(2.45),
            max(PptxInches(0.95), PptxInches(0.34 * max(len(table_rows), 1))),
        )
        table_top = content_top
        _add_powerpoint_inline_table(
            slide,
            primary_table,
            left=content_left,
            top=table_top,
            width=content_width,
            height=table_height,
        )

        content_placeholder.left = content_left
        content_placeholder.top = table_top + table_height + PptxInches(0.25)
        content_placeholder.width = content_width
        content_placeholder.height = max(
            PptxInches(0.45),
            presentation.slide_height - content_placeholder.top - PptxInches(0.95),
        )
    elif bullets:
        content_placeholder.left = content_left
        content_placeholder.top = content_top
        content_placeholder.width = content_width
        content_placeholder.height = content_height

    text_frame = content_placeholder.text_frame
    text_frame.clear()
    text_frame.word_wrap = True

    try:
        bullet_char_limit = int(slide_spec.get('bullet_char_limit') or POWERPOINT_BULLET_CHAR_LIMIT)
    except (TypeError, ValueError):
        bullet_char_limit = POWERPOINT_BULLET_CHAR_LIMIT

    if len(bullets) <= 3:
        font_size = 22
    elif len(bullets) <= 5:
        font_size = 20
    elif len(bullets) <= 8:
        font_size = 18
    elif len(bullets) <= 10:
        font_size = 16
    else:
        font_size = 14

    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = _clean_slide_text(bullet, bullet_char_limit)
        paragraph.level = 0
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = PptxPt(8)
        for run in paragraph.runs:
            run.font.size = PptxPt(font_size)
            run.font.color.rgb = POWERPOINT_TEXT

    metadata_box = slide.shapes.add_textbox(
        presentation.slide_width - PptxInches(3.1),
        presentation.slide_height - PptxInches(0.45),
        PptxInches(2.6),
        PptxInches(0.22),
    )
    metadata_frame = metadata_box.text_frame
    footer_label = str(slide_spec.get('footer_label') or '').strip()
    metadata_frame.text = footer_label or f'{role_label}{f" | {timestamp}" if timestamp else ""}'
    metadata_paragraph = metadata_frame.paragraphs[0]
    metadata_paragraph.alignment = PP_ALIGN.RIGHT
    for run in metadata_paragraph.runs:
        run.font.size = PptxPt(9)
        run.font.color.rgb = POWERPOINT_MUTED


def _add_powerpoint_inline_image_to_slide(
    slide,
    image_asset: Dict[str, Any],
    left: int,
    top: int,
    max_width: int,
    max_height: int,
):
    image_bytes = image_asset.get('image_bytes')
    if not image_bytes:
        return

    caption = _clean_slide_text(image_asset.get('caption') or '', 120)
    caption_height = PptxInches(0.38) if caption else 0
    image_max_height = max(PptxInches(0.8), int(max_height) - int(caption_height) - PptxInches(0.08))
    picture_left, picture_top, picture_width, picture_height = _fit_powerpoint_image_within_bounds(
        image_bytes,
        left=left,
        top=top,
        max_width=max_width,
        max_height=image_max_height,
    )
    slide.shapes.add_picture(
        io.BytesIO(image_bytes),
        picture_left,
        picture_top,
        width=picture_width,
        height=picture_height,
    )

    if not caption:
        return

    caption_box = slide.shapes.add_textbox(
        int(left),
        int(top) + int(max_height) - int(caption_height),
        int(max_width),
        int(caption_height),
    )
    caption_frame = caption_box.text_frame
    caption_frame.text = caption
    caption_paragraph = caption_frame.paragraphs[0]
    caption_paragraph.alignment = PP_ALIGN.CENTER
    for run in caption_paragraph.runs:
        run.font.size = PptxPt(10)
        run.font.color.rgb = POWERPOINT_MUTED


def _fit_powerpoint_image_within_bounds(
    image_bytes: bytes,
    left: int,
    top: int,
    max_width: int,
    max_height: int,
) -> Tuple[int, int, int, int]:
    with Image.open(io.BytesIO(image_bytes)) as image:
        image_width, image_height = image.size

    left = int(left)
    top = int(top)
    max_width = int(max_width)
    max_height = int(max_height)

    if not image_width or not image_height:
        return left, top, max_width, max_height

    aspect_ratio = image_width / image_height
    target_width = max_width
    target_height = int(target_width / aspect_ratio)

    if target_height > max_height:
        target_height = max_height
        target_width = int(target_height * aspect_ratio)

    fitted_left = left + int((max_width - target_width) / 2)
    fitted_top = top + int((max_height - target_height) / 2)
    return fitted_left, fitted_top, target_width, target_height


def _add_powerpoint_inline_table(
    slide,
    table_asset: Dict[str, Any],
    left: int,
    top: int,
    width: int,
    height: int,
):
    rows = table_asset.get('rows', [])
    if not rows:
        return

    row_count = len(rows)
    column_count = max(len(row) for row in rows)
    if row_count < 1 or column_count < 1:
        return

    table_shape = slide.shapes.add_table(row_count, column_count, left, top, width, height)
    table = table_shape.table
    has_header = bool(table_asset.get('has_header'))
    font_size = 10 if row_count <= 5 and column_count <= 3 else 9

    for row_index, row_values in enumerate(rows):
        for column_index in range(column_count):
            cell = table.cell(row_index, column_index)
            cell.text = row_values[column_index] if column_index < len(row_values) else ''
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = POWERPOINT_ACCENT if has_header and row_index == 0 else POWERPOINT_PANEL

            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.size = PptxPt(font_size)
                run.font.bold = bool(has_header and row_index == 0)
                run.font.color.rgb = POWERPOINT_TITLE_TEXT if has_header and row_index == 0 else POWERPOINT_TEXT


def _append_powerpoint_appendix_slides(
    presentation: Presentation,
    appendix_assets: Dict[str, List[Dict[str, Any]]],
    citation_labels: List[str],
):
    for image_asset in appendix_assets.get('images', []):
        _add_powerpoint_image_slide(presentation, image_asset)

    for table_asset in appendix_assets.get('tables', []):
        _add_powerpoint_table_slide(presentation, table_asset)

    for code_asset in appendix_assets.get('code_blocks', []):
        _add_powerpoint_code_slide(presentation, code_asset)

    for index, citation_chunk in enumerate(_chunk_items(citation_labels, 8), start=1):
        title = 'References' if index == 1 else f'References {index}'
        _add_powerpoint_content_slide(
            presentation,
            {'title': title, 'bullets': citation_chunk},
            role_label='Sources',
            timestamp='',
        )


def _add_powerpoint_image_slide(presentation: Presentation, image_asset: Dict[str, Any]):
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    _apply_powerpoint_background(slide, POWERPOINT_PANEL)

    title_shape = slide.shapes.title
    title_shape.text = _clean_slide_text(image_asset.get('title') or 'Visual', 80)
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = PptxPt(24)
            run.font.bold = True
            run.font.color.rgb = POWERPOINT_TEXT

    image_bytes = image_asset.get('image_bytes')
    if image_bytes:
        left, top, width, height = _fit_powerpoint_image(
            image_bytes,
            max_width=presentation.slide_width - PptxInches(1.4),
            max_height=presentation.slide_height - PptxInches(2.3),
            top_offset=PptxInches(1.2),
        )
        slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width=width, height=height)

    caption = _clean_slide_text(image_asset.get('caption') or '', 160)
    if caption:
        caption_box = slide.shapes.add_textbox(
            PptxInches(0.8),
            presentation.slide_height - PptxInches(0.8),
            presentation.slide_width - PptxInches(1.6),
            PptxInches(0.35),
        )
        caption_frame = caption_box.text_frame
        caption_frame.text = caption
        caption_paragraph = caption_frame.paragraphs[0]
        caption_paragraph.alignment = PP_ALIGN.CENTER
        for run in caption_paragraph.runs:
            run.font.size = PptxPt(12)
            run.font.color.rgb = POWERPOINT_MUTED


def _fit_powerpoint_image(
    image_bytes: bytes,
    max_width: int,
    max_height: int,
    top_offset: int,
) -> Tuple[int, int, int, int]:
    with Image.open(io.BytesIO(image_bytes)) as image:
        image_width, image_height = image.size

    max_width = int(max_width)
    max_height = int(max_height)
    top_offset = int(top_offset)

    if not image_width or not image_height:
        return PptxInches(0.8), top_offset, max_width, max_height

    aspect_ratio = image_width / image_height
    target_width = max_width
    target_height = int(target_width / aspect_ratio)

    if target_height > max_height:
        target_height = max_height
        target_width = int(target_height * aspect_ratio)

    left = int((13333320 - target_width) / 2)
    top = top_offset + int((max_height - target_height) / 2)
    return left, top, target_width, target_height


def _add_powerpoint_table_slide(presentation: Presentation, table_asset: Dict[str, Any]):
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    _apply_powerpoint_background(slide, POWERPOINT_PANEL)

    title_shape = slide.shapes.title
    title_shape.text = _clean_slide_text(table_asset.get('title') or 'Table', 80)
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = PptxPt(24)
            run.font.bold = True
            run.font.color.rgb = POWERPOINT_TEXT

    rows = table_asset.get('rows', [])
    if not rows:
        return

    row_count = len(rows)
    column_count = max(len(row) for row in rows)
    table_shape = slide.shapes.add_table(
        row_count,
        column_count,
        PptxInches(0.6),
        PptxInches(1.4),
        presentation.slide_width - PptxInches(1.2),
        presentation.slide_height - PptxInches(2.0),
    )
    table = table_shape.table
    has_header = bool(table_asset.get('has_header'))

    for row_index, row_values in enumerate(rows):
        for column_index in range(column_count):
            cell = table.cell(row_index, column_index)
            cell.text = row_values[column_index] if column_index < len(row_values) else ''
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = POWERPOINT_ACCENT if has_header and row_index == 0 else POWERPOINT_PANEL
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.size = PptxPt(12)
                run.font.bold = bool(has_header and row_index == 0)
                run.font.color.rgb = POWERPOINT_TITLE_TEXT if has_header and row_index == 0 else POWERPOINT_TEXT


def _add_powerpoint_code_slide(presentation: Presentation, code_asset: Dict[str, Any]):
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    _apply_powerpoint_background(slide, POWERPOINT_PANEL)

    title_shape = slide.shapes.title
    title_shape.text = _clean_slide_text(code_asset.get('title') or 'Code Example', 80)
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = PptxPt(24)
            run.font.bold = True
            run.font.color.rgb = POWERPOINT_TEXT

    code_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        PptxInches(0.7),
        PptxInches(1.35),
        presentation.slide_width - PptxInches(1.4),
        presentation.slide_height - PptxInches(2.0),
    )
    code_panel.fill.solid()
    code_panel.fill.fore_color.rgb = RGBColor(241, 245, 249)
    code_panel.line.color.rgb = RGBColor(203, 213, 225)

    code_box = slide.shapes.add_textbox(
        PptxInches(0.95),
        PptxInches(1.55),
        presentation.slide_width - PptxInches(1.9),
        presentation.slide_height - PptxInches(2.4),
    )
    code_frame = code_box.text_frame
    code_frame.word_wrap = False
    code_frame.text = str(code_asset.get('code') or '')
    for paragraph in code_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.name = 'Consolas'
            run.font.size = PptxPt(12)
            run.font.color.rgb = POWERPOINT_TEXT


def _apply_powerpoint_background(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _chunk_items(items: List[str], chunk_size: int) -> List[List[str]]:
    if not items:
        return []
    return [items[index:index + chunk_size] for index in range(0, len(items), chunk_size)]


def _message_to_email_draft_payload(
    message: Dict[str, Any],
    settings: Dict[str, Any],
    summary_model_deployment: str = ''
) -> Dict[str, Any]:
    content = _normalize_content(message.get('content', ''))
    subject_payload = _build_message_email_subject(
        content=content,
        settings=settings,
        requested_model=summary_model_deployment
    )
    body_content = _strip_explicit_message_email_subject(content)
    rendered_body_content = _render_message_export_content(
        message,
        source_content=body_content,
    )
    chart_attachments = _extract_email_chart_png_attachments(rendered_body_content)
    image_labels_by_src = {
        attachment['data_uri']: _format_email_chart_attachment_reference(attachment)
        for attachment in chart_attachments
        if attachment.get('data_uri')
    }

    body_lines = []

    if rendered_body_content.strip():
        body_lines.extend(_render_markdown_to_email_lines(
            rendered_body_content,
            image_labels_by_src=image_labels_by_src,
        ))
    else:
        body_lines.append('No content recorded.')

    citation_labels = _build_message_citation_labels(message)
    if citation_labels:
        if body_lines and body_lines[-1] != '':
            body_lines.append('')
        body_lines.append('Citations')
        body_lines.append('---------')
        for citation_label in citation_labels:
            body_lines.append(f'- {citation_label}')

    body = _finalize_email_body_text(body_lines)
    return {
        'subject': subject_payload['subject'],
        'subject_source': subject_payload['source'],
        'body': body,
        'attachments': chart_attachments,
    }


def _extract_email_chart_png_attachments(rendered_content: str) -> List[Dict[str, str]]:
    if not str(rendered_content or '').strip():
        return []

    html = markdown2.markdown(rendered_content, extras=DOCX_MARKDOWN_EXTRAS)
    soup = BeautifulSoup(f'<div>{html}</div>', 'html.parser')
    root = soup.div if soup.div else soup
    attachments: List[Dict[str, str]] = []
    seen_keys = set()

    for image_node in root.find_all('img'):
        chart_wrapper = image_node.find_parent(class_='export-inline-chart')
        image_wrapper = image_node.find_parent(class_='export-inline-image')
        visual_wrapper = chart_wrapper or image_wrapper
        if not visual_wrapper:
            continue
        visual_type = 'chart' if chart_wrapper else 'image'

        data_uri = str(image_node.get('src') or '').strip()
        image_bytes = decode_base64_image_data_uri(data_uri)
        if not image_bytes:
            continue

        image_key = (len(image_bytes), image_bytes[:24])
        if image_key in seen_keys:
            continue
        seen_keys.add(image_key)

        caption_node = visual_wrapper.find(class_='export-inline-chart-caption')
        if caption_node is None:
            caption_node = visual_wrapper.find(class_='export-inline-image-caption')
        caption = (
            caption_node.get_text(' ', strip=True)
            if caption_node and caption_node.get_text(' ', strip=True)
            else ''
        )
        alt_text = str(image_node.get('alt') or '').strip()
        attachment_label = caption or alt_text or f'{visual_type.title()} {len(attachments) + 1}'
        if visual_type == 'chart':
            filename = _safe_email_chart_attachment_filename(
                attachment_label,
                len(attachments) + 1,
            )
        else:
            filename = _safe_email_image_attachment_filename(
                attachment_label,
                len(attachments) + 1,
            )
        attachments.append({
            'filename': filename,
            'content_type': 'image/png',
            'data_uri': data_uri,
            'caption': caption,
            'alt': alt_text or attachment_label,
            'visual_type': visual_type,
        })

    return attachments


def _safe_email_chart_attachment_filename(label: str, sequence_number: int) -> str:
    return _safe_email_visual_attachment_filename(
        label,
        sequence_number,
        EMAIL_CHART_ATTACHMENT_FILENAME_PREFIX,
        'chart',
    )


def _safe_email_image_attachment_filename(label: str, sequence_number: int) -> str:
    return _safe_email_visual_attachment_filename(
        label,
        sequence_number,
        EMAIL_IMAGE_ATTACHMENT_FILENAME_PREFIX,
        'image',
    )


def _safe_email_visual_attachment_filename(
    label: str,
    sequence_number: int,
    prefix: str,
    fallback_stem: str,
) -> str:
    stem = re.sub(r'[^A-Za-z0-9._-]+', '_', str(label or '').lower())
    stem = stem.strip('_.-')
    if len(stem) > 48:
        stem = stem[:48].strip('_.-')
    if not stem:
        stem = f'{fallback_stem}_{sequence_number}'
    return f'{prefix}_{sequence_number}_{stem}.png'


def _format_email_chart_attachment_reference(attachment: Dict[str, str]) -> str:
    visual_type = str(attachment.get('visual_type') or 'chart').strip().lower()
    filename = str(attachment.get('filename') or f'{visual_type}.png').strip()
    caption = str(attachment.get('caption') or attachment.get('alt') or '').strip()
    if visual_type == 'image':
        if caption and caption != filename:
            return f'Image PNG exported as {filename}: {caption}'
        return f'Image PNG exported as {filename}'

    if caption and caption != filename:
        return f'Chart image exported as {filename}: {caption}'
    return f'Chart image exported as {filename}'


def _render_markdown_to_email_lines(
    content: str,
    image_labels_by_src: Optional[Dict[str, str]] = None,
) -> List[str]:
    html = markdown2.markdown(content, extras=DOCX_MARKDOWN_EXTRAS)
    soup = BeautifulSoup(f'<div>{html}</div>', 'html.parser')
    root = soup.div if soup.div else soup
    lines: List[str] = []

    for child in root.children:
        if isinstance(child, NavigableString):
            text = str(child).strip()
            if text:
                lines.append(text)
                lines.append('')
            continue

        if isinstance(child, Tag):
            _append_html_block_to_email_lines(
                lines,
                child,
                image_labels_by_src=image_labels_by_src,
            )

    return lines


def _append_html_block_to_email_lines(
    lines: List[str],
    node: Tag,
    list_level: int = 0,
    image_labels_by_src: Optional[Dict[str, str]] = None,
):
    tag_name = node.name.lower()

    if tag_name in {'h1', 'h2', 'h3', 'h4', 'h5', 'h6'}:
        heading_text = _extract_email_inline_text(
            node,
            image_labels_by_src=image_labels_by_src,
        ).strip()
        if heading_text:
            underline_char = '=' if tag_name in {'h1', 'h2'} else '-'
            lines.append(heading_text)
            lines.append(underline_char * min(len(heading_text), 80))
            lines.append('')
        return

    if tag_name == 'p':
        if image_labels_by_src and any(
            class_name in {'export-inline-chart-caption', 'export-inline-image-caption'}
            for class_name in (node.get('class') or [])
        ):
            return

        paragraph_text = _extract_email_inline_text(
            node,
            image_labels_by_src=image_labels_by_src,
        ).strip()
        if paragraph_text:
            lines.extend(paragraph_text.splitlines())
            lines.append('')
        return

    if tag_name in {'ul', 'ol'}:
        _append_html_list_to_email_lines(
            lines,
            node,
            ordered=(tag_name == 'ol'),
            level=list_level,
            image_labels_by_src=image_labels_by_src,
        )
        lines.append('')
        return

    if tag_name == 'pre':
        code_text = node.get_text().rstrip('\n')
        if code_text:
            for code_line in code_text.splitlines():
                lines.append(f'    {code_line.rstrip()}')
            lines.append('')
        return

    if tag_name == 'blockquote':
        quote_text = _extract_email_inline_text(
            node,
            image_labels_by_src=image_labels_by_src,
        ).strip()
        if quote_text:
            for quote_line in quote_text.splitlines():
                lines.append(f'    {quote_line}')
            lines.append('')
        return

    if tag_name == 'table':
        _append_html_table_to_email_lines(
            lines,
            node,
            image_labels_by_src=image_labels_by_src,
        )
        lines.append('')
        return

    if tag_name == 'hr':
        lines.append('-' * 40)
        lines.append('')
        return

    if tag_name in {'div', 'section', 'article'}:
        for child in node.children:
            if isinstance(child, NavigableString):
                text = str(child).strip()
                if text:
                    lines.append(text)
                    lines.append('')
                continue

            if isinstance(child, Tag):
                _append_html_block_to_email_lines(
                    lines,
                    child,
                    list_level=list_level,
                    image_labels_by_src=image_labels_by_src,
                )
        return

    fallback_text = _extract_email_inline_text(
        node,
        image_labels_by_src=image_labels_by_src,
    ).strip()
    if fallback_text:
        lines.extend(fallback_text.splitlines())
        lines.append('')


def _append_html_list_to_email_lines(
    lines: List[str],
    list_node: Tag,
    ordered: bool,
    level: int = 0,
    image_labels_by_src: Optional[Dict[str, str]] = None,
):
    item_number = 1
    indent = '  ' * level

    for item in list_node.find_all('li', recursive=False):
        prefix = f'{item_number}. ' if ordered else '- '
        item_parts = []

        for child in item.children:
            if isinstance(child, Tag) and child.name.lower() in {'ul', 'ol'}:
                continue
            item_parts.append(_extract_email_inline_text(
                child,
                image_labels_by_src=image_labels_by_src,
            ))

        item_text = ''.join(item_parts).strip()
        if item_text:
            lines.append(f'{indent}{prefix}{item_text}')
        else:
            lines.append(f'{indent}{prefix}'.rstrip())

        for nested_list in item.find_all(['ul', 'ol'], recursive=False):
            _append_html_list_to_email_lines(
                lines,
                nested_list,
                ordered=(nested_list.name.lower() == 'ol'),
                level=level + 1,
                image_labels_by_src=image_labels_by_src,
            )

        if ordered:
            item_number += 1


def _append_html_table_to_email_lines(
    lines: List[str],
    table_node: Tag,
    image_labels_by_src: Optional[Dict[str, str]] = None,
):
    rows = table_node.find_all('tr')
    if not rows:
        return

    parsed_rows = []
    header_present = False
    for row_index, row in enumerate(rows):
        cells = row.find_all(['th', 'td'], recursive=False)
        if not cells:
            continue
        if row_index == 0 and all(cell.name.lower() == 'th' for cell in cells):
            header_present = True
        parsed_rows.append([
            re.sub(
                r'\s+',
                ' ',
                _extract_email_inline_text(
                    cell,
                    image_labels_by_src=image_labels_by_src,
                ),
            ).strip()
            for cell in cells
        ])

    if not parsed_rows:
        return

    column_count = max(len(row) for row in parsed_rows)
    normalized_rows = [row + [''] * (column_count - len(row)) for row in parsed_rows]
    column_widths = [
        max(len(row[column_index]) for row in normalized_rows)
        for column_index in range(column_count)
    ]

    def format_row(row_values: List[str]) -> str:
        padded_cells = [
            row_values[column_index].ljust(column_widths[column_index])
            for column_index in range(column_count)
        ]
        return '  '.join(padded_cells).rstrip()

    lines.append(format_row(normalized_rows[0]))
    if header_present:
        separator = '  '.join(
            '-' * max(column_widths[column_index], 3)
            for column_index in range(column_count)
        )
        lines.append(separator)
        data_rows = normalized_rows[1:]
    else:
        data_rows = normalized_rows[1:]

    for row_values in data_rows:
        lines.append(format_row(row_values))


def _extract_email_inline_text(
    node: Any,
    image_labels_by_src: Optional[Dict[str, str]] = None,
) -> str:
    if isinstance(node, NavigableString):
        return str(node)

    if not isinstance(node, Tag):
        return ''

    tag_name = node.name.lower()
    if tag_name == 'br':
        return '\n'
    if tag_name == 'img':
        src = str(node.get('src') or '').strip()
        image_label = (
            image_labels_by_src.get(src)
            if image_labels_by_src and src in image_labels_by_src
            else str(node.get('alt') or 'Image').strip() or 'Image'
        )
        return f'[{image_label}]'
    if tag_name == 'a':
        label = ''.join(
            _extract_email_inline_text(child, image_labels_by_src=image_labels_by_src)
            for child in node.children
        ).strip()
        href = str(node.get('href') or '').strip()
        if href and href != label:
            if label:
                return f'{label} ({href})'
            return href
        return label

    return ''.join(
        _extract_email_inline_text(child, image_labels_by_src=image_labels_by_src)
        for child in node.children
    )


def _finalize_email_body_text(lines: List[str]) -> str:
    normalized_lines: List[str] = []

    for raw_line in lines:
        line = str(raw_line or '').rstrip()
        if not line:
            if normalized_lines and normalized_lines[-1] != '':
                normalized_lines.append('')
            continue
        normalized_lines.append(line)

    while normalized_lines and normalized_lines[-1] == '':
        normalized_lines.pop()

    return '\n'.join(normalized_lines)


def _build_message_email_subject(
    content: str,
    settings: Dict[str, Any],
    requested_model: str = ''
) -> Dict[str, str]:
    explicit_subject = _extract_message_email_subject(content)
    if explicit_subject:
        return {
            'subject': explicit_subject,
            'source': 'message'
        }

    generated_subject = _generate_message_email_subject_with_model(
        content=content,
        settings=settings,
        requested_model=requested_model
    )
    if generated_subject:
        return {
            'subject': generated_subject,
            'source': 'model'
        }

    return {
        'subject': _fallback_message_email_subject(content),
        'source': 'fallback'
    }


def _extract_message_email_subject(content: str) -> Optional[str]:
    if not content:
        return None

    lines = content.splitlines()
    explicit_patterns = [
        re.compile(r'^\s*(?:\*\*|__)?(?:email\s+)?subject(?:\*\*|__)?\s*:\s*(.+?)\s*$', re.IGNORECASE),
        re.compile(r'^\s*(?:\*\*|__)?title(?:\*\*|__)?\s*:\s*(.+?)\s*$', re.IGNORECASE),
    ]

    for line in lines:
        stripped_line = line.strip()
        if not stripped_line:
            continue
        for pattern in explicit_patterns:
            match = pattern.match(stripped_line)
            if not match:
                continue
            cleaned_subject = _clean_email_subject(match.group(1))
            if cleaned_subject:
                return cleaned_subject

    for line in lines:
        stripped_line = line.strip()
        if not stripped_line:
            continue

        heading_match = re.match(r'^#{1,6}\s+(.+)$', stripped_line)
        if heading_match:
            cleaned_subject = _clean_email_subject(heading_match.group(1))
            if cleaned_subject:
                return cleaned_subject
        break

    return None


def _strip_explicit_message_email_subject(content: str) -> str:
    if not content:
        return ''

    lines = content.splitlines()
    explicit_patterns = [
        re.compile(r'^\s*(?:\*\*|__)?(?:email\s+)?subject(?:\*\*|__)?\s*:\s*(.+?)\s*$', re.IGNORECASE),
        re.compile(r'^\s*(?:\*\*|__)?title(?:\*\*|__)?\s*:\s*(.+?)\s*$', re.IGNORECASE),
    ]

    first_non_empty_index = None
    for index, line in enumerate(lines):
        if line.strip():
            first_non_empty_index = index
            break

    if first_non_empty_index is None:
        return ''

    first_line = lines[first_non_empty_index].strip()
    if not any(pattern.match(first_line) for pattern in explicit_patterns):
        return content

    remaining_lines = lines[:first_non_empty_index] + lines[first_non_empty_index + 1:]
    while remaining_lines and not remaining_lines[0].strip():
        remaining_lines.pop(0)
    return '\n'.join(remaining_lines)


def _clean_email_subject(subject: str) -> str:
    cleaned_subject = re.sub(r'[`*_~]+', '', str(subject or ''))
    cleaned_subject = re.sub(r'\s+', ' ', cleaned_subject).strip()
    cleaned_subject = cleaned_subject.strip('"\'')
    cleaned_subject = cleaned_subject.rstrip(' .:;-')
    if len(cleaned_subject) > EMAIL_SUBJECT_CHAR_LIMIT:
        cleaned_subject = cleaned_subject[:EMAIL_SUBJECT_CHAR_LIMIT].rstrip(' .:;-')
    return cleaned_subject


def _generate_message_email_subject_with_model(
    content: str,
    settings: Dict[str, Any],
    requested_model: str = ''
) -> Optional[str]:
    subject_source = str(content or '').strip()
    if not subject_source:
        return None

    truncated_source = subject_source[:EMAIL_SUBJECT_SOURCE_CHAR_LIMIT]

    try:
        gpt_client, gpt_model = _initialize_gpt_client(settings, requested_model)
        model_lower = gpt_model.lower()
        is_reasoning_model = (
            'o1' in model_lower or 'o3' in model_lower or 'gpt-5' in model_lower
        )
        instruction_role = 'developer' if is_reasoning_model else 'system'
        subject_prompt = (
            'You are generating an email subject line for a mailto draft from a single chat message. '
            'If the message already contains a subject or clear title, reuse it in cleaned form. '
            'Otherwise, write a concise and specific subject line. '
            'Return plain text only with no quotes, no markdown, and no more than 10 words.'
        )

        subject_response = gpt_client.chat.completions.create(
            model=gpt_model,
            messages=[
                {
                    'role': instruction_role,
                    'content': subject_prompt
                },
                {
                    'role': 'user',
                    'content': truncated_source
                }
            ]
        )
        raw_subject = (
            (subject_response.choices[0].message.content or '').strip()
            if subject_response.choices else ''
        )
        cleaned_subject = _clean_email_subject(raw_subject)
        if cleaned_subject:
            return cleaned_subject
    except Exception as exc:
        debug_print(f'Message email subject generation failed: {exc}')
        log_event(
            'Message email subject generation failed',
            extra={
                'requested_model': requested_model or None,
                'content_length': len(subject_source)
            },
            level='WARNING'
        )

    return None


def _fallback_message_email_subject(content: str) -> str:
    extracted_subject = _extract_message_email_subject(content)
    if extracted_subject:
        return extracted_subject

    for line in str(content or '').splitlines():
        cleaned_subject = _clean_email_subject(line)
        if cleaned_subject:
            return cleaned_subject

    return 'Shared chat message'


def _build_message_citation_labels(message: Dict[str, Any]) -> List[str]:
    normalized_citations = _normalize_citations(_collect_raw_citation_buckets(message))
    citation_labels: List[str] = []
    seen_labels = set()

    for citation in normalized_citations:
        label = str(
            citation.get('label')
            or citation.get('title')
            or citation.get('url')
            or citation.get('filepath')
            or citation.get('tool_name')
            or citation.get('function_name')
            or ''
        ).strip()
        if not label or label in seen_labels:
            continue
        seen_labels.add(label)
        citation_labels.append(label)

    return citation_labels


def _add_markdown_content_to_doc(doc: DocxDocument, content: str):
    html = markdown2.markdown(content, extras=DOCX_MARKDOWN_EXTRAS)
    soup = BeautifulSoup(f'<div>{html}</div>', 'html.parser')
    root = soup.div if soup.div else soup
    rendered_blocks = False

    for child in root.children:
        if isinstance(child, NavigableString):
            text = str(child).strip()
            if not text:
                continue
            paragraph = doc.add_paragraph()
            paragraph.add_run(text)
            rendered_blocks = True
            continue

        if not isinstance(child, Tag):
            continue

        _append_html_block_to_doc(doc, child)
        rendered_blocks = True

    if not rendered_blocks and content.strip():
        doc.add_paragraph(content.strip())


def _append_html_block_to_doc(doc: DocxDocument, node: Tag, list_level: int = 0):
    tag_name = node.name.lower()

    if tag_name in {'h1', 'h2', 'h3', 'h4', 'h5', 'h6'}:
        paragraph = doc.add_heading('', level=min(int(tag_name[1]), 4))
        _append_inline_html_runs(paragraph, node)
        return

    if tag_name == 'p':
        paragraph = doc.add_paragraph()
        _append_inline_html_runs(paragraph, node)
        return

    if tag_name in {'ul', 'ol'}:
        _append_list_items_to_doc(doc, node, ordered=(tag_name == 'ol'), level=list_level)
        return

    if tag_name == 'pre':
        _add_code_block_to_doc(doc, node)
        return

    if tag_name == 'blockquote':
        paragraph = doc.add_paragraph()
        paragraph.paragraph_format.left_indent = Inches(0.3)
        _append_inline_html_runs(paragraph, node, {'italic': True})
        return

    if tag_name == 'table':
        _add_html_table_to_doc(doc, node)
        return

    if tag_name == 'hr':
        doc.add_paragraph('')
        return

    if tag_name in {'div', 'section', 'article'}:
        for child in node.children:
            if isinstance(child, NavigableString):
                text = str(child).strip()
                if not text:
                    continue
                paragraph = doc.add_paragraph()
                paragraph.add_run(text)
                continue

            if isinstance(child, Tag):
                _append_html_block_to_doc(doc, child, list_level=list_level)
        return

    paragraph = doc.add_paragraph()
    _append_inline_html_runs(paragraph, node)


def _append_list_items_to_doc(doc: DocxDocument, list_node: Tag, ordered: bool, level: int = 0):
    style_name = 'List Number' if ordered else 'List Bullet'

    for item in list_node.find_all('li', recursive=False):
        paragraph = doc.add_paragraph(style=style_name)
        if level:
            paragraph.paragraph_format.left_indent = Inches(0.25 * level)

        rendered_inline = False
        for child in item.children:
            if isinstance(child, Tag) and child.name.lower() in {'ul', 'ol'}:
                continue
            if isinstance(child, NavigableString) and not str(child).strip():
                continue

            _append_inline_html_runs(paragraph, child)
            rendered_inline = True

        if not rendered_inline:
            text = item.get_text(' ', strip=True)
            if text:
                paragraph.add_run(text)

        for nested_list in item.find_all(['ul', 'ol'], recursive=False):
            _append_list_items_to_doc(
                doc,
                nested_list,
                ordered=(nested_list.name.lower() == 'ol'),
                level=level + 1
            )


def _add_code_block_to_doc(doc: DocxDocument, node: Tag):
    code_text = node.get_text().rstrip('\n')
    if not code_text:
        return

    paragraph = doc.add_paragraph()
    paragraph.paragraph_format.left_indent = Inches(0.25)
    paragraph.paragraph_format.space_before = Pt(6)
    paragraph.paragraph_format.space_after = Pt(6)
    run = paragraph.add_run(code_text)
    run.font.name = 'Consolas'
    run.font.size = Pt(9)


def _add_html_table_to_doc(doc: DocxDocument, table_node: Tag):
    rows = table_node.find_all('tr')
    if not rows:
        return

    column_count = max(
        len(row.find_all(['th', 'td'], recursive=False))
        for row in rows
    )
    if column_count == 0:
        return

    table = doc.add_table(rows=len(rows), cols=column_count)
    table.style = 'Table Grid'

    for row_index, row in enumerate(rows):
        cells = row.find_all(['th', 'td'], recursive=False)
        for column_index in range(column_count):
            cell = table.cell(row_index, column_index)
            cell.text = ''

            if column_index >= len(cells):
                continue

            _populate_table_cell(
                cell,
                cells[column_index],
                is_header=(cells[column_index].name.lower() == 'th')
            )


def _populate_table_cell(cell, node: Tag, is_header: bool = False):
    paragraph = cell.paragraphs[0]
    _append_inline_html_runs(paragraph, node, {'bold': is_header})


def _append_inline_html_runs(paragraph, node: Any, formatting: Optional[Dict[str, bool]] = None):
    if formatting is None:
        formatting = {}

    if isinstance(node, NavigableString):
        text = str(node)
        if not text:
            return

        run = paragraph.add_run(text)
        _apply_run_formatting(run, formatting)
        return

    if not isinstance(node, Tag):
        return

    tag_name = node.name.lower()
    if tag_name == 'br':
        paragraph.add_run().add_break()
        return

    if tag_name == 'img':
        image_bytes = decode_base64_image_data_uri(node.get('src'))
        if image_bytes:
            try:
                paragraph.add_run().add_picture(io.BytesIO(image_bytes), width=Inches(6.0))
                return
            except Exception:
                pass

        alt_text = node.get('alt') or 'Image'
        run = paragraph.add_run(f'[{alt_text}]')
        _apply_run_formatting(run, formatting)
        return

    next_formatting = dict(formatting)
    if tag_name in {'strong', 'b'}:
        next_formatting['bold'] = True
    elif tag_name in {'em', 'i'}:
        next_formatting['italic'] = True
    elif tag_name in {'s', 'strike', 'del'}:
        next_formatting['strike'] = True
    elif tag_name == 'code':
        next_formatting['code'] = True
    elif tag_name == 'a':
        next_formatting['underline'] = True

    for child in node.children:
        _append_inline_html_runs(paragraph, child, next_formatting)

    if tag_name == 'a':
        href = str(node.get('href') or '').strip()
        label = node.get_text(' ', strip=True)
        if href and href != label:
            suffix_run = paragraph.add_run(f' ({href})')
            _apply_run_formatting(suffix_run, formatting)


def _apply_run_formatting(run, formatting: Dict[str, bool]):
    if formatting.get('bold'):
        run.bold = True
    if formatting.get('italic'):
        run.italic = True
    if formatting.get('underline'):
        run.underline = True
    if formatting.get('strike'):
        run.font.strike = True
    if formatting.get('code'):
        run.font.name = 'Consolas'
        run.font.size = Pt(9)


# ---------------------------------------------------------------------------
# PDF Export — HTML generation and PyMuPDF Story rendering
# ---------------------------------------------------------------------------

_PDF_CSS = """
body {
    font-family: sans-serif;
    font-size: 10pt;
    color: #222;
    line-height: 1.4;
}
h1 {
    font-size: 16pt;
    color: #1a1a2e;
    margin-bottom: 2pt;
}
h2 {
    font-size: 13pt;
    color: #16213e;
    margin-top: 16pt;
    margin-bottom: 6pt;
    border-bottom: 1px solid #ccc;
    padding-bottom: 4pt;
}
h3 {
    font-size: 11pt;
    color: #0f3460;
    margin-top: 10pt;
    margin-bottom: 4pt;
}
h4 {
    font-size: 10pt;
    color: #333;
    margin-top: 8pt;
    margin-bottom: 4pt;
}
p {
    margin-top: 2pt;
    margin-bottom: 4pt;
}
.metadata {
    font-size: 8pt;
    color: #666;
}
.abstract {
    background-color: #f8f9fa;
    padding: 8pt;
    margin-bottom: 8pt;
}
.note {
    font-size: 9pt;
    color: #856404;
    background-color: #fff3cd;
    padding: 6pt;
}
.bubble {
    padding: 8pt 12pt;
    margin-bottom: 8pt;
}
.bubble-header {
    font-size: 8pt;
    color: #444;
    margin-bottom: 2pt;
}
.ts {
    font-weight: normal;
    color: #888;
}
.user-bubble {
    background-color: #c8e0fa;
    margin-left: 60pt;
}
.assistant-bubble {
    background-color: #f1f0f0;
    margin-right: 60pt;
}
.system-bubble {
    background-color: #fff3cd;
    margin-left: 30pt;
    margin-right: 30pt;
    font-size: 9pt;
}
.file-bubble {
    background-color: #e8f5e9;
    margin-right: 60pt;
    font-size: 9pt;
}
.other-bubble {
    background-color: #f5f5f5;
    margin-left: 30pt;
    margin-right: 30pt;
    font-size: 9pt;
}
table {
    border-collapse: collapse;
    width: 100%;
    font-size: 9pt;
    margin-bottom: 8pt;
}
th, td {
    border: 1px solid #ddd;
    padding: 4pt 6pt;
    text-align: left;
}
th {
    background-color: #f5f5f5;
    font-weight: bold;
}
pre {
    background-color: #f5f5f5;
    padding: 6pt;
    font-size: 8pt;
    font-family: monospace;
}
code {
    font-family: monospace;
    font-size: 9pt;
    background-color: #f0f0f0;
    padding: 1pt 3pt;
}
ol, ul {
    margin-top: 4pt;
    margin-bottom: 8pt;
}
li {
    margin-bottom: 4pt;
}
small {
    font-size: 8pt;
    color: #666;
}
.export-inline-chart {
    background-color: #fafafa;
    border: 1px solid #ddd;
    padding: 8pt;
    margin-top: 6pt;
    margin-bottom: 10pt;
}
.export-inline-chart img {
    max-width: 100%;
    height: auto;
    display: block;
    margin: 0 auto;
}
.export-inline-chart-caption {
    font-size: 8pt;
    color: #666;
    text-align: center;
    margin-top: 4pt;
}
a {
    color: #0066cc;
}
"""


def _pdf_bubble_class(role: str) -> str:
    """Return the CSS class for a chat bubble based on message role."""
    role_classes = {
        'user': 'user-bubble',
        'assistant': 'assistant-bubble',
        'system': 'system-bubble',
        'file': 'file-bubble',
        'image': 'file-bubble'
    }
    return role_classes.get(role, 'other-bubble')


def _build_pdf_html_body(entry: Dict[str, Any]) -> str:
    """Build the HTML body content for a single conversation PDF."""
    conversation = entry['conversation']
    messages = entry['messages']
    summary_intro = entry.get('summary_intro', {}) or {}

    transcript_messages = [m for m in messages if m.get('is_transcript_message')]
    detail_messages = [m for m in messages if m.get('details')]
    reference_messages = [m for m in messages if m.get('citations')]
    thought_messages = [m for m in messages if m.get('thoughts')]
    supplemental_messages = [m for m in messages if not m.get('is_transcript_message')]

    parts: List[str] = []

    # --- Title and metadata ---
    parts.append(f'<h1>{_escape_html(conversation.get("title", "Untitled"))}</h1>')
    meta_items = [
        f'<b>Last Updated:</b> {_escape_html(str(conversation.get("last_updated", "")))}',
        f'<b>Chat Type:</b> {_escape_html(str(conversation.get("chat_type", "personal")))}',
        f'<b>Messages:</b> {conversation.get("message_count", len(messages))}'
    ]
    tags = conversation.get('tags')
    if tags:
        meta_items.append(f'<b>Tags:</b> {_escape_html(", ".join(_format_tag(t) for t in tags))}')
    classification = conversation.get('classification')
    if classification:
        meta_items.append(
            f'<b>Classification:</b> {_escape_html(", ".join(_format_tag(c) for c in classification))}'
        )
    parts.append(f'<p class="metadata">{" &nbsp;|&nbsp; ".join(meta_items)}</p>')

    # --- Abstract ---
    if summary_intro.get('enabled') and summary_intro.get('generated') and summary_intro.get('content'):
        parts.append('<h2>Abstract</h2>')
        abstract_html = markdown2.markdown(
            replace_inline_chart_blocks_with_export_html(summary_intro.get('content', '')),
            extras=['fenced-code-blocks', 'tables']
        )
        parts.append(f'<div class="abstract">{abstract_html}</div>')
        parts.append(
            f'<p class="metadata"><i>Generated with '
            f'{_escape_html(str(summary_intro.get("model_deployment") or "configured model"))} on '
            f'{_escape_html(str(summary_intro.get("generated_at", "")))}</i></p>'
        )
    elif summary_intro.get('enabled') and summary_intro.get('error'):
        error_text = _escape_html(str(summary_intro.get('error', '')))
        parts.append(
            '<p class="note"><i>A summary intro was requested, '
            'but could not be generated for this export.</i><br/>'
            f'<small>Error: {error_text}</small></p>'
        )

    # --- Transcript with chat bubbles ---
    parts.append('<h2>Transcript</h2>')
    if not transcript_messages:
        parts.append(
            '<p><i>No user or assistant transcript messages were available for export.</i></p>'
        )
    else:
        for message in transcript_messages:
            role = message.get('role', '')
            bubble_class = _pdf_bubble_class(role)
            label = message.get('label', '')
            speaker = message.get('speaker_label', '')
            timestamp = message.get('timestamp', '')
            content = message.get('content_text', '') or 'No content recorded.'

            parts.append(f'<div class="bubble {bubble_class}">')
            ts_str = (
                f' &nbsp;|&nbsp; <span class="ts">{_escape_html(str(timestamp))}</span>'
                if timestamp else ''
            )
            parts.append(
                f'<p class="bubble-header"><b>{_escape_html(label)} — '
                f'{_escape_html(speaker)}</b>{ts_str}</p>'
            )
            content_html = markdown2.markdown(
                replace_inline_chart_blocks_with_export_html(content),
                extras=['fenced-code-blocks', 'tables', 'break-on-newline']
            )
            parts.append(content_html)
            parts.append('</div>')

    # --- Appendix A: Conversation Metadata ---
    parts.append('<h2>Appendix A — Conversation Metadata</h2>')
    metadata_to_render = _remove_empty_values({
        'context': conversation.get('context'),
        'classification': conversation.get('classification'),
        'strict': conversation.get('strict'),
        'is_pinned': conversation.get('is_pinned'),
        'scope_locked': conversation.get('scope_locked'),
        'locked_contexts': conversation.get('locked_contexts'),
        'message_counts_by_role': conversation.get('message_counts_by_role'),
        'citation_counts': conversation.get('citation_counts'),
        'thought_count': conversation.get('thought_count')
    })
    _append_html_table(parts, metadata_to_render)

    # --- Appendix B: Message Details ---
    if detail_messages:
        parts.append('<h2>Appendix B — Message Details</h2>')
        for message in detail_messages:
            parts.append(
                f'<h3>{_escape_html(message.get("label", ""))} — '
                f'{_escape_html(message.get("speaker_label", ""))}</h3>'
            )
            if message.get('timestamp'):
                parts.append(
                    f'<p class="metadata"><i>{_escape_html(str(message.get("timestamp")))}</i></p>'
                )
            _append_html_table(parts, message.get('details', {}))

    # --- Appendix C: References ---
    if reference_messages:
        parts.append('<h2>Appendix C — References</h2>')
        for message in reference_messages:
            parts.append(
                f'<h3>{_escape_html(message.get("label", ""))} — '
                f'{_escape_html(message.get("speaker_label", ""))}</h3>'
            )
            if message.get('timestamp'):
                parts.append(
                    f'<p class="metadata"><i>{_escape_html(str(message.get("timestamp")))}</i></p>'
                )
            _append_html_citations(parts, message)

    # --- Appendix D: Processing Thoughts ---
    if thought_messages:
        parts.append('<h2>Appendix D — Processing Thoughts</h2>')
        for message in thought_messages:
            parts.append(
                f'<h3>{_escape_html(message.get("label", ""))} — '
                f'{_escape_html(message.get("speaker_label", ""))}</h3>'
            )
            if message.get('timestamp'):
                parts.append(
                    f'<p class="metadata"><i>{_escape_html(str(message.get("timestamp")))}</i></p>'
                )
            parts.append('<ol>')
            for thought in message.get('thoughts', []):
                thought_label = (thought.get('step_type') or 'step').replace('_', ' ').title()
                parts.append(
                    f'<li><b>{_escape_html(thought_label)}:</b> '
                    f'{_escape_html(str(thought.get("content") or "No content recorded."))}'
                )
                if thought.get('duration_ms') is not None:
                    parts.append(
                        f'<br/><small><b>Duration:</b> {thought.get("duration_ms")} ms</small>'
                    )
                if thought.get('timestamp'):
                    parts.append(
                        f'<br/><small><b>Timestamp:</b> '
                        f'{_escape_html(str(thought.get("timestamp")))}</small>'
                    )
                if thought.get('detail'):
                    parts.append('<br/><small><b>Detail:</b></small>')
                    _append_html_code_block(parts, thought.get('detail'))
                parts.append('</li>')
            parts.append('</ol>')

    # --- Appendix E: Supplemental Messages ---
    if supplemental_messages:
        parts.append('<h2>Appendix E — Supplemental Messages</h2>')
        for message in supplemental_messages:
            parts.append(
                f'<h3>{_escape_html(message.get("label", ""))} — '
                f'{_escape_html(message.get("speaker_label", ""))}</h3>'
            )
            if message.get('timestamp'):
                parts.append(
                    f'<p class="metadata"><i>{_escape_html(str(message.get("timestamp")))}</i></p>'
                )
            content = message.get('content_text', '') or 'No content recorded.'
            content_html = markdown2.markdown(
                replace_inline_chart_blocks_with_export_html(content),
                extras=['fenced-code-blocks', 'tables', 'break-on-newline']
            )
            parts.append(content_html)

    return '\n'.join(parts)


def _render_pdf_bytes(body_html: str) -> bytes:
    """Render HTML body content to PDF bytes using PyMuPDF Story API."""
    MEDIABOX = fitz.paper_rect("letter")
    WHERE = MEDIABOX + (36, 36, -36, -36)

    story = fitz.Story(html=body_html, user_css=_PDF_CSS)

    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp_path = tmp.name

        writer = fitz.DocumentWriter(tmp_path)
        more = True
        while more:
            device = writer.begin_page(MEDIABOX)
            more, _ = story.place(WHERE)
            story.draw(device)
            writer.end_page()
        writer.close()
        del story
        del writer

        with open(tmp_path, 'rb') as f:
            return f.read()
    finally:
        if tmp_path:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass


def _conversation_to_pdf_bytes(entry: Dict[str, Any]) -> bytes:
    """Convert a conversation export entry to PDF bytes."""
    body_html = _build_pdf_html_body(entry)
    return _render_pdf_bytes(body_html)


def _html_body_to_pdf_bytes(body_html: str) -> bytes:
    """Convert raw HTML body content to PDF bytes."""
    return _render_pdf_bytes(body_html)


def _append_html_table(parts: List[str], mapping: Dict[str, Any]):
    """Append a key-value mapping as an HTML table."""
    if not isinstance(mapping, dict) or not mapping:
        parts.append('<p><i>No data available.</i></p>')
        return

    parts.append('<table>')
    parts.append('<tr><th>Property</th><th>Value</th></tr>')
    for key, value in mapping.items():
        label = _format_markdown_key(key)
        if isinstance(value, dict):
            formatted = _format_nested_html_value(value)
        elif isinstance(value, list):
            formatted = (
                ', '.join(_escape_html(str(item)) for item in value)
                if value else '<i>None</i>'
            )
        elif isinstance(value, bool):
            formatted = 'Yes' if value else 'No'
        else:
            formatted = _escape_html(str(value))
        parts.append(f'<tr><td><b>{_escape_html(label)}</b></td><td>{formatted}</td></tr>')
    parts.append('</table>')


def _format_nested_html_value(mapping: Dict[str, Any], depth: int = 0) -> str:
    """Format a nested dict as an HTML string for table cells."""
    if not mapping:
        return '<i>None</i>'

    items = []
    for key, value in mapping.items():
        label = _format_markdown_key(key)
        if isinstance(value, dict):
            nested = _format_nested_html_value(value, depth + 1)
            items.append(f'<b>{_escape_html(label)}:</b><br/>{nested}')
        elif isinstance(value, list):
            list_str = (
                ', '.join(_escape_html(str(v)) for v in value)
                if value else 'None'
            )
            items.append(f'<b>{_escape_html(label)}:</b> {list_str}')
        elif isinstance(value, bool):
            items.append(f'<b>{_escape_html(label)}:</b> {"Yes" if value else "No"}')
        else:
            items.append(f'<b>{_escape_html(label)}:</b> {_escape_html(str(value))}')
    return '<br/>'.join(items)


def _append_html_citations(parts: List[str], message: Dict[str, Any]):
    """Append citation data as HTML."""
    citations = message.get('citations', [])
    if not citations:
        parts.append('<p><i>No citations were recorded for this message.</i></p>')
        return

    doc_citations = [c for c in citations if c.get('citation_type') == 'document']
    web_citations = [c for c in citations if c.get('citation_type') == 'web']
    agent_citations = [c for c in citations if c.get('citation_type') == 'agent_tool']
    legacy_citations = [c for c in citations if c.get('citation_type') == 'legacy']

    if doc_citations:
        parts.append('<h4>Document Sources</h4>')
        parts.append('<ol>')
        for citation in doc_citations:
            parts.append(
                f'<li><b>{_escape_html(str(citation.get("label", "Document source")))}</b>'
            )
            detail_items = _remove_empty_values({
                'citation_id': citation.get('citation_id'),
                'page_number': citation.get('page_number'),
                'classification': citation.get('classification'),
                'score': citation.get('score'),
                'metadata_type': citation.get('metadata_type')
            })
            if detail_items:
                detail_str = '; '.join(
                    f'{_format_markdown_key(k)}: {_escape_html(str(v))}'
                    for k, v in detail_items.items()
                )
                parts.append(f'<br/><small>{detail_str}</small>')
            if citation.get('metadata_content'):
                parts.append('<br/><small><b>Metadata Content:</b></small>')
                _append_html_code_block(parts, citation.get('metadata_content'))
            parts.append('</li>')
        parts.append('</ol>')

    if web_citations:
        parts.append('<h4>Web Sources</h4>')
        parts.append('<ol>')
        for citation in web_citations:
            title = _escape_html(
                str(citation.get('title') or citation.get('label') or 'Web source')
            )
            url = citation.get('url')
            if url:
                parts.append(f'<li><a href="{_escape_html(url)}">{title}</a></li>')
            else:
                parts.append(f'<li>{title}</li>')
        parts.append('</ol>')

    if agent_citations:
        parts.append('<h4>Tool Invocations</h4>')
        parts.append('<ol>')
        for citation in agent_citations:
            label = _escape_html(
                str(citation.get('tool_name') or citation.get('function_name') or 'Tool')
            )
            parts.append(f'<li><b>{label}</b>')
            detail_items = _remove_empty_values({
                'function_name': citation.get('function_name'),
                'plugin_name': citation.get('plugin_name'),
                'success': citation.get('success'),
                'timestamp': citation.get('timestamp')
            })
            if detail_items:
                detail_str = '; '.join(
                    f'{_format_markdown_key(k)}: {_escape_html(str(v))}'
                    for k, v in detail_items.items()
                )
                parts.append(f'<br/><small>{detail_str}</small>')
            parts.append('</li>')
        parts.append('</ol>')

    if legacy_citations:
        parts.append('<h4>Legacy Citation Records</h4>')
        parts.append('<ol>')
        for citation in legacy_citations:
            parts.append(
                f'<li>{_escape_html(str(citation.get("label", "Legacy citation")))}</li>'
            )
        parts.append('</ol>')


def _append_html_code_block(parts: List[str], value: Any):
    """Append a code block in HTML format."""
    if isinstance(value, (dict, list)):
        code_text = json.dumps(value, indent=2, ensure_ascii=False, default=str)
    else:
        code_text = str(value)
    parts.append(f'<pre>{_escape_html(code_text)}</pre>')
