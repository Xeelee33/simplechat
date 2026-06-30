# functions_simplechat_operations.py
"""Shared SimpleChat-native operations for routes and Semantic Kernel plugins."""

import logging
import mimetypes
import os
import re
import tempfile
import uuid
from datetime import datetime, timezone
from typing import Any, Dict, Iterable, List, Optional, Tuple
from urllib.parse import quote

import requests
from azure.cosmos.exceptions import CosmosResourceNotFoundError
from flask import current_app, has_app_context, session

from collaboration_models import normalize_collaboration_user
from config import (
    CLIENTS,
    cosmos_activity_logs_container,
    cosmos_conversations_container,
    cosmos_groups_container,
    cosmos_messages_container,
    storage_account_personal_chat_container_name,
    TABULAR_EXTENSIONS,
)
from functions_activity_logging import (
    log_chat_activity,
    log_conversation_creation,
    log_document_upload,
    log_group_status_change,
    log_workflow_creation,
)
from functions_appinsights import log_event
from functions_authentication import (
    get_current_user_info,
    get_graph_endpoint,
    get_valid_access_token,
)
from functions_collaboration import (
    assert_user_can_participate_in_collaboration_conversation,
    create_collaboration_message_notifications,
    create_group_collaboration_conversation_record,
    create_personal_collaboration_conversation_record,
    get_collaboration_conversation,
    invite_personal_collaboration_participants,
    is_group_collaboration_conversation,
    persist_collaboration_message,
)
from functions_documents import allowed_file, create_document, process_document_upload_background, update_document
from functions_group import (
    assert_group_role,
    check_group_status_allows_operation,
    create_group,
    find_group_by_id,
    get_user_role_in_group,
    require_active_group,
)
from functions_notifications import create_notification
from functions_personal_workflows import save_personal_workflow
from functions_settings import get_settings, is_user_workflows_enabled_for_user
from utils_cache import invalidate_group_search_cache, invalidate_personal_search_cache


SIMPLECHAT_PLUGIN_TYPE = "simplechat"
SIMPLECHAT_DEFAULT_ENDPOINT = "simplechat://internal"
SIMPLECHAT_CAPABILITY_TO_FUNCTION = {
    "create_group": "create_group",
    "add_group_member": "add_user_to_group",
    "make_group_inactive": "make_group_inactive",
    "create_group_conversation": "create_group_conversation",
    "invite_group_conversation_members": "invite_group_conversation_members",
    "add_conversation_message": "add_conversation_message",
    "upload_markdown_document": "upload_markdown_document",
    "upload_word_document": "upload_word_document",
    "upload_powerpoint_document": "upload_powerpoint_document",
    "create_personal_conversation": "create_personal_conversation",
    "create_personal_workflow": "create_personal_workflow",
    "create_personal_collaboration_conversation": "create_personal_collaboration_conversation",
}
SIMPLECHAT_CAPABILITY_DEFINITIONS = [
    {
        "key": "create_group",
        "function_name": SIMPLECHAT_CAPABILITY_TO_FUNCTION["create_group"],
        "label": "Create Groups",
        "description": "Create a new group workspace as the current user.",
    },
    {
        "key": "add_group_member",
        "function_name": SIMPLECHAT_CAPABILITY_TO_FUNCTION["add_group_member"],
        "label": "Add Group Members",
        "description": "Add a user directly to a group as a member, admin, or document manager.",
    },
    {
        "key": "make_group_inactive",
        "function_name": SIMPLECHAT_CAPABILITY_TO_FUNCTION["make_group_inactive"],
        "label": "Make Groups Inactive",
        "description": "Mark a group inactive using the current user's Control Center admin permissions.",
    },
    {
        "key": "create_group_conversation",
        "function_name": SIMPLECHAT_CAPABILITY_TO_FUNCTION["create_group_conversation"],
        "label": "Create Group Multi-User Conversations",
        "description": "Create an invite-managed multi-user conversation in a group the current user can access, then add current group members as participants to grant access.",
    },
    {
        "key": "invite_group_conversation_members",
        "function_name": SIMPLECHAT_CAPABILITY_TO_FUNCTION["invite_group_conversation_members"],
        "label": "Invite Group Conversation Members",
        "description": "Invite current group members into an existing invite-managed group multi-user conversation the current user manages.",
    },
    {
        "key": "add_conversation_message",
        "function_name": SIMPLECHAT_CAPABILITY_TO_FUNCTION["add_conversation_message"],
        "label": "Add Conversation Messages",
        "description": "Add a user-authored message to a personal or collaborative conversation the current user can access.",
    },
    {
        "key": "upload_markdown_document",
        "function_name": SIMPLECHAT_CAPABILITY_TO_FUNCTION["upload_markdown_document"],
        "label": "Upload Markdown Documents",
        "description": "Create and upload a Markdown document into the current user's personal workspace or an allowed group workspace.",
    },
    {
        "key": "upload_word_document",
        "function_name": SIMPLECHAT_CAPABILITY_TO_FUNCTION["upload_word_document"],
        "label": "Upload Word Documents",
        "description": "Create and upload a Word document into the current user's personal workspace or an allowed group workspace.",
    },
    {
        "key": "upload_powerpoint_document",
        "function_name": SIMPLECHAT_CAPABILITY_TO_FUNCTION["upload_powerpoint_document"],
        "label": "Upload PowerPoint Documents",
        "description": "Create and upload a PowerPoint presentation into the current user's personal workspace or an allowed group workspace.",
    },
    {
        "key": "create_personal_conversation",
        "function_name": SIMPLECHAT_CAPABILITY_TO_FUNCTION["create_personal_conversation"],
        "label": "Create Personal Conversations",
        "description": "Create a standard one-user personal conversation.",
    },
    {
        "key": "create_personal_workflow",
        "function_name": SIMPLECHAT_CAPABILITY_TO_FUNCTION["create_personal_workflow"],
        "label": "Create Personal Workflows",
        "description": "Create a personal workflow for the current user using the existing workflow engine and permissions.",
    },
    {
        "key": "create_personal_collaboration_conversation",
        "function_name": SIMPLECHAT_CAPABILITY_TO_FUNCTION["create_personal_collaboration_conversation"],
        "label": "Create Personal Collaborative Conversations",
        "description": "Create a personal collaborative conversation and invite other users.",
    },
]
CONVERSATION_ACCESS_ERROR = "Conversation not found or access denied"


def get_default_simplechat_capabilities() -> Dict[str, bool]:
    return {definition["key"]: True for definition in SIMPLECHAT_CAPABILITY_DEFINITIONS}


def normalize_simplechat_capabilities(raw_capabilities: Any = None) -> Dict[str, bool]:
    normalized = get_default_simplechat_capabilities()

    if raw_capabilities is None:
        return normalized

    if isinstance(raw_capabilities, dict):
        for capability_key in normalized:
            if capability_key in raw_capabilities:
                normalized[capability_key] = bool(raw_capabilities[capability_key])
        if "upload_markdown_document" in raw_capabilities:
            upload_enabled = bool(raw_capabilities["upload_markdown_document"])
            for upload_capability_key in ("upload_word_document", "upload_powerpoint_document"):
                if upload_capability_key not in raw_capabilities:
                    normalized[upload_capability_key] = upload_enabled
        return normalized

    if isinstance(raw_capabilities, (list, tuple, set)):
        enabled_items = {str(item or "").strip() for item in raw_capabilities if str(item or "").strip()}
        return {
            definition["key"]: (
                definition["key"] in enabled_items or definition["function_name"] in enabled_items
            )
            for definition in SIMPLECHAT_CAPABILITY_DEFINITIONS
        }

    return normalized


def get_simplechat_enabled_function_names(raw_capabilities: Any = None) -> List[str]:
    normalized = normalize_simplechat_capabilities(raw_capabilities)
    return [
        definition["function_name"]
        for definition in SIMPLECHAT_CAPABILITY_DEFINITIONS
        if normalized.get(definition["key"], False)
    ]


def resolve_simplechat_action_capabilities(
    action_capability_map: Any,
    action_defaults: Any = None,
    action_id: Optional[str] = None,
    action_name: Optional[str] = None,
) -> Dict[str, bool]:
    resolved_defaults = normalize_simplechat_capabilities(action_defaults)

    if not isinstance(action_capability_map, dict):
        return resolved_defaults

    for candidate_key in (str(action_id or "").strip(), str(action_name or "").strip()):
        if candidate_key and candidate_key in action_capability_map:
            return normalize_simplechat_capabilities(action_capability_map.get(candidate_key))

    return resolved_defaults


def derive_conversation_title_from_message(content: str) -> str:
    normalized_content = re.sub(r"\s+", " ", str(content or "").strip())
    if not normalized_content:
        return "New Conversation"
    return f"{normalized_content[:30]}..." if len(normalized_content) > 30 else normalized_content


def create_personal_conversation_for_current_user(
    title: str = "New Conversation",
    notify_creation: bool = False,
) -> Dict[str, Any]:
    current_user = _require_current_user_info()
    normalized_title = str(title or "").strip() or "New Conversation"
    conversation_id = str(uuid.uuid4())
    conversation_item = {
        "id": conversation_id,
        "user_id": current_user["userId"],
        "last_updated": datetime.utcnow().isoformat(),
        "title": normalized_title,
        "context": [],
        "tags": [],
        "strict": False,
        "is_pinned": False,
        "is_hidden": False,
        "chat_type": "new",
        "has_unread_assistant_response": False,
        "last_unread_assistant_message_id": None,
        "last_unread_assistant_at": None,
    }
    cosmos_conversations_container.upsert_item(conversation_item)

    log_conversation_creation(
        user_id=current_user["userId"],
        conversation_id=conversation_id,
        title=normalized_title,
        workspace_type="personal",
    )

    conversation_item["added_to_activity_log"] = True
    cosmos_conversations_container.upsert_item(conversation_item)

    if notify_creation:
        _notify_personal_conversation_created(
            conversation_item=conversation_item,
            current_user=current_user,
        )

    return conversation_item


def create_personal_workflow_for_current_user(
    name: str,
    task_prompt: str,
    description: str = "",
    runner_type: str = "model",
    trigger_type: str = "manual",
    selected_agent_name: str = "",
    selected_agent_id: str = "",
    selected_agent_is_global: bool = False,
    model_endpoint_id: str = "",
    model_id: str = "",
    alert_priority: str = "none",
    is_enabled: bool = True,
    schedule_value: int = 1,
    schedule_unit: str = "hours",
    conversation_id: str = "",
) -> Dict[str, Any]:
    _require_user_workflows_enabled()
    current_user_info = _require_current_user_info()

    normalized_runner_type = str(runner_type or "model").strip().lower() or "model"
    normalized_trigger_type = str(trigger_type or "manual").strip().lower() or "manual"
    workflow_payload = {
        "name": str(name or "").strip(),
        "description": str(description or "").strip(),
        "task_prompt": str(task_prompt or "").strip(),
        "runner_type": normalized_runner_type,
        "trigger_type": normalized_trigger_type,
        "alert_priority": str(alert_priority or "none").strip().lower() or "none",
        "is_enabled": bool(is_enabled) if normalized_trigger_type == "interval" else True,
        "conversation_id": str(conversation_id or "").strip(),
    }

    if normalized_runner_type == "agent":
        workflow_payload["selected_agent"] = {
            "id": str(selected_agent_id or "").strip(),
            "name": str(selected_agent_name or "").strip(),
            "is_global": bool(selected_agent_is_global),
        }
    else:
        workflow_payload["model_endpoint_id"] = str(model_endpoint_id or "").strip()
        workflow_payload["model_id"] = str(model_id or "").strip()

    if normalized_trigger_type == "interval":
        workflow_payload["schedule"] = {
            "value": int(schedule_value),
            "unit": str(schedule_unit or "hours").strip().lower() or "hours",
        }

    workflow = save_personal_workflow(
        current_user_info["userId"],
        workflow_payload,
        actor_user_id=current_user_info["userId"],
    )
    log_workflow_creation(
        user_id=current_user_info["userId"],
        workflow_id=str(workflow.get("id") or "").strip(),
        workflow_name=str(workflow.get("name") or "").strip(),
        runner_type=workflow.get("runner_type"),
        trigger_type=workflow.get("trigger_type"),
    )
    return {
        "workflow": workflow,
        "message": f"Created workflow '{workflow.get('name', 'Workflow')}'.",
    }


def add_conversation_message_for_current_user(
    conversation_id: str,
    content: str,
    reply_to_message_id: str = "",
) -> Dict[str, Any]:
    current_user_info = _require_current_user_info()
    normalized_conversation_id = str(conversation_id or "").strip()
    normalized_content = str(content or "").strip()
    normalized_reply_to_message_id = str(reply_to_message_id or "").strip() or None

    if not normalized_conversation_id:
        raise ValueError("conversation_id is required")
    if not normalized_content:
        raise ValueError("content is required")

    try:
        conversation_item = cosmos_conversations_container.read_item(
            item=normalized_conversation_id,
            partition_key=normalized_conversation_id,
        )
    except CosmosResourceNotFoundError:
        conversation_item = None

    if conversation_item is not None:
        if str(conversation_item.get("user_id") or "").strip() != current_user_info["userId"]:
            raise LookupError(CONVERSATION_ACCESS_ERROR)

        message_doc, updated_conversation = _persist_personal_conversation_message(
            conversation_item=conversation_item,
            current_user_info=current_user_info,
            content=normalized_content,
            reply_to_message_id=normalized_reply_to_message_id,
        )
        return {
            "conversation": updated_conversation,
            "message": message_doc,
            "conversation_kind": "personal",
        }

    current_user = normalize_collaboration_user(current_user_info)
    if not current_user:
        raise PermissionError("User not authenticated")

    try:
        collaboration_conversation = get_collaboration_conversation(normalized_conversation_id)
    except CosmosResourceNotFoundError as exc:
        raise LookupError(CONVERSATION_ACCESS_ERROR) from exc

    try:
        assert_user_can_participate_in_collaboration_conversation(
            current_user["user_id"],
            collaboration_conversation,
        )
    except (LookupError, PermissionError) as exc:
        raise LookupError(CONVERSATION_ACCESS_ERROR) from exc
    message_doc, updated_conversation = persist_collaboration_message(
        collaboration_conversation,
        current_user,
        normalized_content,
        reply_to_message_id=normalized_reply_to_message_id,
    )
    create_collaboration_message_notifications(updated_conversation, message_doc)
    return {
        "conversation": updated_conversation,
        "message": message_doc,
        "conversation_kind": "collaboration",
    }


def upload_markdown_document_for_current_user(
    file_name: str,
    markdown_content: str,
    workspace_scope: str = "personal",
    group_id: str = "",
    default_group_id: str = "",
) -> Dict[str, Any]:
    current_user_info = _require_current_user_info()
    current_user_id = current_user_info["userId"]
    normalized_workspace_scope = _normalize_document_workspace_scope(workspace_scope)
    normalized_file_name = _normalize_markdown_file_name(file_name)
    raw_markdown_content = str(markdown_content or "")

    if not raw_markdown_content.strip():
        raise ValueError("markdown_content is required")
    if not allowed_file(normalized_file_name, allowed_extensions={"md"}):
        raise ValueError("Only Markdown files are supported")

    return _upload_generated_document_for_current_user(
        current_user_id=current_user_id,
        normalized_file_name=normalized_file_name,
        file_content_bytes=raw_markdown_content.encode("utf-8"),
        normalized_workspace_scope=normalized_workspace_scope,
        group_id=group_id,
        default_group_id=default_group_id,
        process_inline=False,
    )


def upload_generated_document_for_current_user(
    file_name: str,
    file_content: Any,
    workspace_scope: str = "personal",
    group_id: str = "",
    default_group_id: str = "",
    process_inline: bool = False,
) -> Dict[str, Any]:
    """Upload generated JSON/CSV-style content into the current user's workspace."""
    current_user_info = _require_current_user_info()
    current_user_id = current_user_info["userId"]
    normalized_workspace_scope = _normalize_document_workspace_scope(workspace_scope)
    normalized_file_name = _normalize_generated_document_file_name(file_name)

    if isinstance(file_content, bytes):
        file_content_bytes = file_content
    else:
        file_content_bytes = str(file_content or "").encode("utf-8")

    if not file_content_bytes.strip():
        raise ValueError("file_content is required")
    if not allowed_file(normalized_file_name):
        raise ValueError("Generated file type is not supported")

    return _upload_generated_document_for_current_user(
        current_user_id=current_user_id,
        normalized_file_name=normalized_file_name,
        file_content_bytes=file_content_bytes,
        normalized_workspace_scope=normalized_workspace_scope,
        group_id=group_id,
        default_group_id=default_group_id,
        process_inline=process_inline,
    )


def upload_word_document_for_current_user(
    file_name: str,
    title: str = "",
    markdown_content: str = "",
    workspace_scope: str = "personal",
    group_id: str = "",
    default_group_id: str = "",
) -> Dict[str, Any]:
    """Create a simple DOCX document from markdown-like text and upload it."""
    # Optional Office rendering dependency; keep lazy so non-export SimpleChat actions still import locally.
    from docx import Document as DocxDocument

    normalized_file_name = _normalize_generated_document_file_name(file_name or title or "generated_word_document.docx")
    if not normalized_file_name.lower().endswith(".docx"):
        normalized_file_name = f"{os.path.splitext(normalized_file_name)[0] or 'generated_word_document'}.docx"

    normalized_title = str(title or os.path.splitext(normalized_file_name)[0] or "Generated Document").strip()
    normalized_content = str(markdown_content or "").strip()
    if not normalized_content:
        raise ValueError("markdown_content is required")

    document = DocxDocument()
    if normalized_title:
        document.add_heading(normalized_title, level=1)
    _append_markdown_like_content_to_docx(document, normalized_content)

    buffer = tempfile.SpooledTemporaryFile(max_size=2 * 1024 * 1024)
    document.save(buffer)
    buffer.seek(0)
    return upload_generated_document_for_current_user(
        file_name=normalized_file_name,
        file_content=buffer.read(),
        workspace_scope=workspace_scope,
        group_id=group_id,
        default_group_id=default_group_id,
    )


def upload_powerpoint_document_for_current_user(
    file_name: str,
    title: str = "",
    markdown_content: str = "",
    workspace_scope: str = "personal",
    group_id: str = "",
    default_group_id: str = "",
) -> Dict[str, Any]:
    """Create a simple PPTX presentation from markdown-like text and upload it."""
    # Optional Office rendering dependency; keep lazy so non-export SimpleChat actions still import locally.
    from pptx import Presentation

    normalized_file_name = _normalize_generated_document_file_name(file_name or title or "generated_presentation.pptx")
    if not normalized_file_name.lower().endswith(".pptx"):
        normalized_file_name = f"{os.path.splitext(normalized_file_name)[0] or 'generated_presentation'}.pptx"

    normalized_title = str(title or os.path.splitext(normalized_file_name)[0] or "Generated Presentation").strip()
    normalized_content = str(markdown_content or "").strip()
    if not normalized_content:
        raise ValueError("markdown_content is required")

    presentation = Presentation()
    _populate_simple_presentation(presentation, normalized_title, normalized_content)

    buffer = tempfile.SpooledTemporaryFile(max_size=2 * 1024 * 1024)
    presentation.save(buffer)
    buffer.seek(0)
    return upload_generated_document_for_current_user(
        file_name=normalized_file_name,
        file_content=buffer.read(),
        workspace_scope=workspace_scope,
        group_id=group_id,
        default_group_id=default_group_id,
    )


def upload_generated_chat_artifact_for_current_user(
    conversation_id: str,
    file_name: str,
    file_content: Any,
) -> Dict[str, Any]:
    """Upload generated JSON/CSV-style content into the current user's current chat."""
    return upload_generated_analysis_artifact_for_current_user(
        conversation_id=conversation_id,
        file_name=file_name,
        file_content=file_content,
        capability="tabular",
    )


def upload_generated_analysis_artifact_for_current_user(
    conversation_id: str,
    file_name: str,
    file_content: Any,
    capability: str = "analysis",
    output_format: str = "",
    summary: str = "",
) -> Dict[str, Any]:
    """Upload generated analysis content into the current user's current chat."""
    current_user_info = _require_current_user_info()
    current_user_id = current_user_info["userId"]
    normalized_conversation_id = str(conversation_id or "").strip()
    normalized_file_name = _normalize_generated_document_file_name(file_name)
    normalized_capability = str(capability or "analysis").strip().lower() or "analysis"
    normalized_output_format = str(output_format or "").strip().lower() or os.path.splitext(normalized_file_name)[1].lower().lstrip(".")
    normalized_summary = str(summary or "").strip()

    if isinstance(file_content, bytes):
        file_content_bytes = file_content
    else:
        file_content_bytes = str(file_content or "").encode("utf-8")

    if not normalized_conversation_id:
        raise ValueError("conversation_id is required")
    if not file_content_bytes.strip():
        raise ValueError("file_content is required")
    if not allowed_file(normalized_file_name):
        raise ValueError("Generated file type is not supported")

    settings = get_settings()
    max_artifact_size_mb = settings.get("max_generated_chat_artifact_size_mb", 500)
    try:
        max_artifact_size_mb = max(1, int(max_artifact_size_mb))
    except (TypeError, ValueError):
        max_artifact_size_mb = 500

    max_artifact_size_bytes = max_artifact_size_mb * 1024 * 1024
    if len(file_content_bytes) > max_artifact_size_bytes:
        raise ValueError(
            f"Generated artifact exceeds the {max_artifact_size_mb} MB size limit"
        )

    return _upload_generated_chat_artifact_for_current_user(
        current_user_id=current_user_id,
        conversation_id=normalized_conversation_id,
        normalized_file_name=normalized_file_name,
        file_content_bytes=file_content_bytes,
        artifact_metadata={
            "capability": normalized_capability,
            "output_format": normalized_output_format,
            "summary": normalized_summary,
        },
    )


def upload_generated_analysis_artifact_for_user(
    current_user_id: str,
    conversation_id: str,
    file_name: str,
    file_content: Any,
    capability: str = "analysis",
    output_format: str = "",
    summary: str = "",
) -> Dict[str, Any]:
    """Upload generated analysis content for a known authorized user outside request context."""
    normalized_user_id = str(current_user_id or "").strip()
    normalized_conversation_id = str(conversation_id or "").strip()
    normalized_file_name = _normalize_generated_document_file_name(file_name)
    normalized_capability = str(capability or "analysis").strip().lower() or "analysis"
    normalized_output_format = str(output_format or "").strip().lower() or os.path.splitext(normalized_file_name)[1].lower().lstrip(".")
    normalized_summary = str(summary or "").strip()

    if isinstance(file_content, bytes):
        file_content_bytes = file_content
    else:
        file_content_bytes = str(file_content or "").encode("utf-8")

    if not normalized_user_id:
        raise ValueError("current_user_id is required")
    if not normalized_conversation_id:
        raise ValueError("conversation_id is required")
    if not file_content_bytes.strip():
        raise ValueError("file_content is required")
    if not allowed_file(normalized_file_name):
        raise ValueError("Generated file type is not supported")

    settings = get_settings()
    max_artifact_size_mb = settings.get("max_generated_chat_artifact_size_mb", 500)
    try:
        max_artifact_size_mb = max(1, int(max_artifact_size_mb))
    except (TypeError, ValueError):
        max_artifact_size_mb = 500

    max_artifact_size_bytes = max_artifact_size_mb * 1024 * 1024
    if len(file_content_bytes) > max_artifact_size_bytes:
        raise ValueError(
            f"Generated artifact exceeds the {max_artifact_size_mb} MB size limit"
        )

    return _upload_generated_chat_artifact_for_current_user(
        current_user_id=normalized_user_id,
        conversation_id=normalized_conversation_id,
        normalized_file_name=normalized_file_name,
        file_content_bytes=file_content_bytes,
        artifact_metadata={
            "capability": normalized_capability,
            "output_format": normalized_output_format,
            "summary": normalized_summary,
        },
    )


def delete_blob_backed_chat_message_files(messages: Iterable[Dict[str, Any]]) -> int:
    """Delete blob-backed chat files referenced by the provided message documents."""
    blob_service_client = CLIENTS.get("storage_account_office_docs_client")
    if not blob_service_client:
        return 0

    deleted_count = 0
    deleted_targets = set()

    for message in messages or []:
        if not isinstance(message, dict):
            continue

        if str(message.get("file_content_source") or "").strip().lower() != "blob":
            continue

        blob_container = str(message.get("blob_container") or "").strip()
        blob_path = str(message.get("blob_path") or "").strip()
        if not blob_container or not blob_path:
            continue

        target = (blob_container, blob_path)
        if target in deleted_targets:
            continue

        try:
            blob_client = blob_service_client.get_blob_client(
                container=blob_container,
                blob=blob_path,
            )
            if not blob_client.exists():
                deleted_targets.add(target)
                continue

            blob_client.delete_blob()
            deleted_targets.add(target)
            deleted_count += 1
        except Exception as exc:
            log_event(
                "[SimpleChat] Failed to delete blob-backed chat file",
                {
                    "blob_container": blob_container,
                    "blob_path": blob_path,
                    "error": str(exc),
                },
                debug_only=True,
            )

    return deleted_count


def download_blob_content(blob_container: str, blob_path: str) -> bytes:
    """Download a blob into memory for internal workspace promotion flows."""
    normalized_blob_container = str(blob_container or "").strip()
    normalized_blob_path = str(blob_path or "").strip()

    if not normalized_blob_container or not normalized_blob_path:
        raise ValueError("blob_container and blob_path are required")

    blob_service_client = CLIENTS.get("storage_account_office_docs_client")
    if not blob_service_client:
        raise RuntimeError("Blob storage client not available")

    blob_client = blob_service_client.get_blob_client(
        container=normalized_blob_container,
        blob=normalized_blob_path,
    )
    return blob_client.download_blob().readall()


def _normalize_chat_image_file_name(file_name: str, content_type: str = "image/png") -> str:
    normalized_file_name = str(file_name or "").replace("\\", "/").split("/")[-1].strip()
    normalized_content_type = str(content_type or "image/png").split(";", 1)[0].strip().lower() or "image/png"
    guessed_extension = mimetypes.guess_extension(normalized_content_type) or ".png"

    if not normalized_file_name:
        return f"chat-image{guessed_extension}"

    base_name, extension = os.path.splitext(normalized_file_name)
    if extension:
        return normalized_file_name

    normalized_base_name = base_name.strip() or normalized_file_name.strip() or "chat-image"
    return f"{normalized_base_name}{guessed_extension}"


def upload_chat_image_bytes_for_user(
    user_id: str,
    conversation_id: str,
    message_id: str,
    file_name: str,
    image_bytes: bytes,
    content_type: str = "image/png",
    image_source: str = "generated",
) -> Dict[str, Any]:
    """Upload chat image bytes to the conversation blob folder and return message fields."""
    normalized_user_id = str(user_id or "").strip()
    normalized_conversation_id = str(conversation_id or "").strip()
    normalized_message_id = str(message_id or "").strip()
    normalized_content_type = str(content_type or "image/png").split(";", 1)[0].strip() or "image/png"
    normalized_image_source = str(image_source or "generated").strip().lower() or "generated"

    if not normalized_user_id or not normalized_conversation_id or not normalized_message_id:
        raise ValueError("user_id, conversation_id, and message_id are required")
    if not isinstance(image_bytes, (bytes, bytearray)) or not image_bytes:
        raise ValueError("image_bytes are required")

    blob_service_client = CLIENTS.get("storage_account_office_docs_client")
    if not blob_service_client:
        raise RuntimeError("Blob storage client not available")

    file_content_bytes = bytes(image_bytes)
    normalized_file_name = _normalize_chat_image_file_name(file_name, normalized_content_type)
    blob_path = (
        f"{normalized_user_id}/{normalized_conversation_id}/images/"
        f"{normalized_message_id}/{normalized_file_name}"
    )
    blob_client = blob_service_client.get_blob_client(
        container=storage_account_personal_chat_container_name,
        blob=blob_path,
    )
    blob_client.upload_blob(
        file_content_bytes,
        overwrite=True,
        metadata={
            "conversation_id": normalized_conversation_id,
            "user_id": normalized_user_id,
            "message_id": normalized_message_id,
            "chat_image": "true",
            "image_source": normalized_image_source,
        },
    )

    log_event(
        "[SimpleChat] Chat image saved to blob storage",
        {
            "conversation_id": normalized_conversation_id,
            "message_id": normalized_message_id,
            "blob_container": storage_account_personal_chat_container_name,
            "blob_path": blob_path,
            "content_type": normalized_content_type,
            "image_source": normalized_image_source,
            "image_size": len(file_content_bytes),
        },
        debug_only=True,
    )

    return {
        "content": f"/api/image/{normalized_message_id}",
        "filename": normalized_file_name,
        "file_content_source": "blob",
        "blob_container": storage_account_personal_chat_container_name,
        "blob_path": blob_path,
        "mime_type": normalized_content_type,
        "image_size": len(file_content_bytes),
    }


def create_group_for_current_user(name: str, description: str = "") -> Dict[str, Any]:
    settings = _require_group_workspaces_enabled()
    _require_group_creation_enabled(settings)

    normalized_name = str(name or "").strip() or "Untitled Group"
    normalized_description = str(description or "").strip()
    current_user = _require_current_user_info()

    group_doc = create_group(normalized_name, normalized_description)
    _notify_group_created(group_doc=group_doc, actor_user=current_user)
    return group_doc


def make_group_inactive_for_current_user(
    group_id: str = "",
    reason: str = "",
    default_group_id: str = "",
) -> Dict[str, Any]:
    current_user_info = _require_current_user_info()
    admin_session_user = _require_control_center_admin_access()

    resolved_group_id = str(group_id or default_group_id or "").strip()
    if not resolved_group_id:
        resolved_group_id = require_active_group(current_user_info["userId"])

    group_doc = find_group_by_id(resolved_group_id)
    if not group_doc:
        raise LookupError("Group not found")

    old_status = str(group_doc.get("status") or "active").strip() or "active"
    if old_status == "inactive":
        return {
            "group": group_doc,
            "old_status": old_status,
            "new_status": old_status,
            "message": f"Group '{group_doc.get('name', 'Unknown')}' is already inactive.",
        }

    changed_at = datetime.utcnow().isoformat()
    changed_by_user_id = str(admin_session_user.get("oid") or current_user_info.get("userId") or "").strip() or "unknown"
    changed_by_email = str(
        admin_session_user.get("preferred_username")
        or current_user_info.get("email")
        or current_user_info.get("userPrincipalName")
        or ""
    ).strip() or "unknown"
    normalized_reason = str(reason or "").strip()

    group_doc["status"] = "inactive"
    group_doc["modifiedDate"] = changed_at
    group_doc.setdefault("statusHistory", []).append(
        {
            "old_status": old_status,
            "new_status": "inactive",
            "changed_by_user_id": changed_by_user_id,
            "changed_by_email": changed_by_email,
            "changed_at": changed_at,
            "reason": normalized_reason,
        }
    )
    updated_group_doc = cosmos_groups_container.upsert_item(group_doc)

    log_group_status_change(
        group_id=resolved_group_id,
        group_name=str(group_doc.get("name") or "Unknown").strip() or "Unknown",
        old_status=old_status,
        new_status="inactive",
        changed_by_user_id=changed_by_user_id,
        changed_by_email=changed_by_email,
        reason=normalized_reason or None,
    )
    log_event(
        "[SimpleChat] Group marked inactive",
        {
            "group_id": resolved_group_id,
            "group_name": group_doc.get("name"),
            "old_status": old_status,
            "new_status": "inactive",
            "changed_by_user_id": changed_by_user_id,
            "changed_by_email": changed_by_email,
            "reason": normalized_reason,
        },
    )

    return {
        "group": updated_group_doc,
        "old_status": old_status,
        "new_status": "inactive",
        "message": f"Marked group '{group_doc.get('name', 'Unknown')}' as inactive.",
    }


def create_group_collaboration_conversation_for_current_user(
    title: str = "",
    group_id: str = "",
    default_group_id: str = "",
) -> Tuple[Dict[str, Any], Dict[str, Any], Dict[str, Any]]:
    _require_collaboration_feature_enabled()
    current_user_info = _require_current_user_info()
    current_user = normalize_collaboration_user(current_user_info)
    if not current_user:
        raise PermissionError("User not authenticated")

    group_doc = _resolve_group_doc_for_current_user(
        current_user_info["userId"],
        group_id=group_id,
        default_group_id=default_group_id,
        allowed_roles=("Owner", "Admin", "DocumentManager", "User"),
        missing_group_message="group_id is required for group collaborative conversations",
    )
    allowed, reason = check_group_status_allows_operation(group_doc, "chat")
    if not allowed:
        raise PermissionError(reason)

    conversation_doc, _ = create_group_collaboration_conversation_record(
        title=str(title or "").strip(),
        creator_user=current_user,
        group_doc=group_doc,
    )
    _notify_group_conversation_created(
        conversation_doc=conversation_doc,
        group_doc=group_doc,
        creator_user=current_user,
    )
    return conversation_doc, current_user, group_doc


def invite_group_conversation_members_for_current_user(
    conversation_id: str,
    participants: Optional[Iterable[Dict[str, Any]]] = None,
    participant_identifiers: Any = None,
) -> Dict[str, Any]:
    _require_collaboration_feature_enabled()
    current_user_info = _require_current_user_info()
    current_user = normalize_collaboration_user(current_user_info)
    if not current_user:
        raise PermissionError("User not authenticated")

    normalized_conversation_id = str(conversation_id or "").strip()
    if not normalized_conversation_id:
        raise ValueError("conversation_id is required")

    try:
        conversation_doc = get_collaboration_conversation(normalized_conversation_id)
    except CosmosResourceNotFoundError as exc:
        raise LookupError(CONVERSATION_ACCESS_ERROR) from exc
    if not is_group_collaboration_conversation(conversation_doc):
        raise LookupError(CONVERSATION_ACCESS_ERROR)

    participants_to_add = _build_invited_participants(
        creator_user=current_user,
        participants=participants,
        participant_identifiers=participant_identifiers,
    )
    if not participants_to_add:
        raise ValueError("At least one participant identifier is required")

    try:
        updated_conversation_doc, invited_state_docs = invite_personal_collaboration_participants(
            normalized_conversation_id,
            current_user["user_id"],
            participants_to_add,
        )
    except (LookupError, PermissionError) as exc:
        raise LookupError(CONVERSATION_ACCESS_ERROR) from exc

    conversation_title = str((updated_conversation_doc or {}).get("title") or "Group Conversation").strip() or "Group Conversation"
    scope = (updated_conversation_doc or {}).get("scope") if isinstance((updated_conversation_doc or {}).get("scope"), dict) else {}
    group_id = str(scope.get("group_id") or "").strip()
    group_name = str(scope.get("group_name") or "Group Workspace").strip() or "Group Workspace"
    invited_participants = [
        {
            "user_id": state_doc.get("user_id"),
            "display_name": state_doc.get("user_display_name"),
            "email": state_doc.get("user_email"),
            "membership_status": state_doc.get("membership_status"),
        }
        for state_doc in invited_state_docs
    ]

    if invited_participants:
        message = (
            f"Invited {len(invited_participants)} current group member(s) to "
            f"'{conversation_title}' in '{group_name}'."
        )
    else:
        message = (
            f"No new group members were invited to '{conversation_title}' in '{group_name}'."
        )

    return {
        "conversation": updated_conversation_doc,
        "group": {
            "id": group_id,
            "name": group_name,
        },
        "invited_participants": invited_participants,
        "message": message,
    }


def create_personal_collaboration_conversation_for_current_user(
    title: str = "",
    participants: Optional[Iterable[Dict[str, Any]]] = None,
    participant_identifiers: Any = None,
) -> Tuple[Dict[str, Any], List[Dict[str, Any]], Dict[str, Any]]:
    _require_collaboration_feature_enabled()
    current_user_info = _require_current_user_info()
    creator_user = normalize_collaboration_user(current_user_info)
    if not creator_user:
        raise PermissionError("User not authenticated")

    invited_participants = _build_invited_participants(
        creator_user=creator_user,
        participants=participants,
        participant_identifiers=participant_identifiers,
    )
    conversation_doc, user_states = create_personal_collaboration_conversation_record(
        title=str(title or "").strip(),
        creator_user=creator_user,
        invited_participants=invited_participants,
    )
    _notify_personal_collaboration_conversation_created(
        conversation_doc=conversation_doc,
        creator_user=creator_user,
        invited_participants=invited_participants,
    )
    return conversation_doc, user_states, creator_user


def add_group_member_for_current_user(
    group_id: str = "",
    user_id: str = "",
    user_identifier: str = "",
    email: str = "",
    display_name: str = "",
    role: str = "user",
    default_group_id: str = "",
) -> Dict[str, Any]:
    current_user = _require_current_user_info()
    group_doc = _resolve_group_doc_for_current_user(
        current_user["userId"],
        group_id=group_id,
        default_group_id=default_group_id,
        allowed_roles=("Owner", "Admin"),
        missing_group_message="group_id is required when adding a user to a group",
    )
    actor_role = get_user_role_in_group(group_doc, current_user["userId"])
    if actor_role not in ["Owner", "Admin"]:
        raise PermissionError("Only the owner or admin can add members")

    member_role = str(role or "user").strip().lower()
    valid_roles = ["admin", "document_manager", "user"]
    if member_role not in valid_roles:
        raise ValueError(f"Invalid role. Must be: {', '.join(valid_roles)}")

    resolved_user = resolve_directory_user(
        user_id=user_id,
        user_identifier=user_identifier,
        email=email,
        display_name=display_name,
    )
    target_user_id = resolved_user["id"]
    if get_user_role_in_group(group_doc, target_user_id):
        raise ValueError("User is already a member")

    new_member_doc = {
        "userId": target_user_id,
        "email": resolved_user.get("email", ""),
        "displayName": resolved_user.get("displayName") or resolved_user.get("email") or target_user_id,
    }
    group_doc.setdefault("users", []).append(new_member_doc)

    if member_role == "admin":
        if target_user_id not in group_doc.get("admins", []):
            group_doc.setdefault("admins", []).append(target_user_id)
    elif member_role == "document_manager":
        if target_user_id not in group_doc.get("documentManagers", []):
            group_doc.setdefault("documentManagers", []).append(target_user_id)

    group_doc["modifiedDate"] = datetime.utcnow().isoformat()
    updated_group_doc = cosmos_groups_container.upsert_item(group_doc)

    _log_group_member_addition(
        actor_user=current_user,
        actor_role=actor_role,
        group_doc=group_doc,
        member_doc=new_member_doc,
        member_role=member_role,
    )
    _notify_group_member_addition(
        group_doc=group_doc,
        member_doc=new_member_doc,
        member_role=member_role,
        added_by_email=current_user.get("email", "unknown"),
        actor_user=current_user,
    )

    return {
        "success": True,
        "message": "Member added",
        "group_id": group_doc.get("id"),
        "group_name": group_doc.get("name", "Unknown"),
        "member": new_member_doc,
        "member_role": member_role,
        "group": updated_group_doc,
    }


def resolve_directory_user(
    user_id: str = "",
    user_identifier: str = "",
    email: str = "",
    display_name: str = "",
) -> Dict[str, str]:
    normalized_user_id = str(user_id or "").strip()
    normalized_identifier = str(user_identifier or "").strip()
    normalized_email = str(email or "").strip()
    normalized_display_name = str(display_name or "").strip()

    if normalized_user_id and (normalized_email or normalized_display_name):
        return {
            "id": normalized_user_id,
            "displayName": normalized_display_name or normalized_email or normalized_user_id,
            "email": normalized_email,
        }

    if normalized_user_id:
        try:
            direct_match = _get_directory_user_by_id(normalized_user_id)
        except PermissionError:
            direct_match = None
        if direct_match:
            return direct_match
        if not (normalized_identifier or normalized_email or normalized_display_name):
            return {
                "id": normalized_user_id,
                "displayName": normalized_user_id,
                "email": "",
            }

    if normalized_email or "@" in normalized_identifier:
        lookup_value = normalized_email or normalized_identifier
        exact_matches = _find_directory_users_by_email(lookup_value)
        if len(exact_matches) == 1:
            return exact_matches[0]
        if len(exact_matches) > 1:
            raise ValueError(f"Multiple directory users matched '{lookup_value}'")

    lookup_query = normalized_identifier or normalized_display_name or normalized_email or normalized_user_id
    if not lookup_query:
        raise ValueError("Missing userId or user identifier")

    search_results = search_directory_users(lookup_query, limit=10)
    if not search_results:
        raise LookupError(f"User '{lookup_query}' was not found in the directory")

    exact_matches = []
    lowered_lookup = lookup_query.lower()
    for candidate in search_results:
        candidate_email = str(candidate.get("email") or "").strip().lower()
        candidate_name = str(candidate.get("displayName") or "").strip().lower()
        candidate_id = str(candidate.get("id") or "").strip().lower()
        if lowered_lookup in {candidate_email, candidate_name, candidate_id}:
            exact_matches.append(candidate)

    if len(exact_matches) == 1:
        return exact_matches[0]
    if len(exact_matches) > 1:
        raise ValueError(f"Multiple directory users matched '{lookup_query}'")
    if len(search_results) == 1:
        return search_results[0]

    raise ValueError(
        f"Multiple directory users matched '{lookup_query}'. Provide a more specific email or user ID."
    )


def search_directory_users(query: str, limit: int = 10) -> List[Dict[str, str]]:
    normalized_query = str(query or "").strip()
    if not normalized_query:
        return []

    escaped_query = _escape_odata_value(normalized_query)
    payload = _graph_get_json(
        "/users",
        params={
            "$filter": (
                f"startswith(displayName, '{escaped_query}') "
                f"or startswith(mail, '{escaped_query}') "
                f"or startswith(userPrincipalName, '{escaped_query}')"
            ),
            "$top": max(1, min(int(limit or 10), 25)),
            "$select": "id,displayName,mail,userPrincipalName",
        },
    )
    return _normalize_directory_users(payload.get("value", []))


def _build_invited_participants(
    creator_user: Dict[str, str],
    participants: Optional[Iterable[Dict[str, Any]]] = None,
    participant_identifiers: Any = None,
) -> List[Dict[str, str]]:
    invited_participants: List[Dict[str, str]] = []
    seen_user_ids = {creator_user.get("user_id")}

    for raw_participant in participants or []:
        normalized_participant = normalize_collaboration_user(raw_participant)
        if not normalized_participant:
            continue
        participant_user_id = normalized_participant.get("user_id")
        if participant_user_id in seen_user_ids:
            continue
        seen_user_ids.add(participant_user_id)
        invited_participants.append(normalized_participant)

    for raw_identifier in _split_participant_identifiers(participant_identifiers):
        resolved_user = resolve_directory_user(user_identifier=raw_identifier)
        normalized_participant = normalize_collaboration_user(resolved_user)
        if not normalized_participant:
            continue
        participant_user_id = normalized_participant.get("user_id")
        if participant_user_id in seen_user_ids:
            continue
        seen_user_ids.add(participant_user_id)
        invited_participants.append(normalized_participant)

    return invited_participants


def _split_participant_identifiers(raw_identifiers: Any) -> List[str]:
    if raw_identifiers is None:
        return []
    if isinstance(raw_identifiers, str):
        values = re.split(r"[,;\n]+", raw_identifiers)
    elif isinstance(raw_identifiers, (list, tuple, set)):
        values = []
        for item in raw_identifiers:
            if isinstance(item, str):
                values.extend(re.split(r"[,;\n]+", item))
    else:
        values = [str(raw_identifiers)]

    return [str(value or "").strip() for value in values if str(value or "").strip()]


def _persist_personal_conversation_message(
    conversation_item: Dict[str, Any],
    current_user_info: Dict[str, str],
    content: str,
    reply_to_message_id: Optional[str] = None,
) -> Tuple[Dict[str, Any], Dict[str, Any]]:
    conversation_id = str(conversation_item.get("id") or "").strip()
    if not conversation_id:
        raise ValueError("Conversation is missing an id")

    timestamp = datetime.now(timezone.utc).isoformat()
    current_thread_id = str(uuid.uuid4())
    previous_thread_id = _get_latest_personal_thread_id(conversation_id)
    normalized_chat_type = str(conversation_item.get("chat_type") or "personal_single_user").strip() or "personal_single_user"

    message_doc = {
        "id": f"{conversation_id}_user_{uuid.uuid4().hex}",
        "conversation_id": conversation_id,
        "role": "user",
        "content": str(content or "").strip(),
        "reply_to_message_id": reply_to_message_id,
        "timestamp": timestamp,
        "model_deployment_name": None,
        "metadata": {
            "user_info": {
                "user_id": current_user_info.get("userId"),
                "username": current_user_info.get("userPrincipalName"),
                "display_name": current_user_info.get("displayName"),
                "email": current_user_info.get("email"),
                "timestamp": timestamp,
            },
            "button_states": {
                "image_generation": False,
                "document_search": False,
                "web_search": False,
            },
            "workspace_search": {
                "search_enabled": False,
            },
            "chat_context": {
                "conversation_id": conversation_id,
                "chat_type": normalized_chat_type,
            },
            "thread_info": {
                "thread_id": current_thread_id,
                "previous_thread_id": previous_thread_id,
                "active_thread": True,
                "thread_attempt": 1,
            },
        },
    }

    cosmos_messages_container.upsert_item(message_doc)

    conversation_item["chat_type"] = normalized_chat_type
    if str(conversation_item.get("title") or "").strip() in {"", "New Conversation"}:
        conversation_item["title"] = _derive_personal_conversation_title(message_doc["content"])
    conversation_item["last_updated"] = timestamp
    cosmos_conversations_container.upsert_item(conversation_item)

    log_chat_activity(
        user_id=current_user_info["userId"],
        conversation_id=conversation_id,
        message_type="user_message",
        message_length=len(message_doc["content"]),
        has_document_search=False,
        has_image_generation=False,
        chat_context=normalized_chat_type,
        workspace_type="personal",
    )

    return message_doc, conversation_item


def _get_latest_personal_thread_id(conversation_id: str) -> Optional[str]:
    query = (
        "SELECT TOP 1 c.metadata.thread_info.thread_id AS thread_id "
        "FROM c WHERE c.conversation_id = @conversation_id ORDER BY c.timestamp DESC"
    )
    items = list(cosmos_messages_container.query_items(
        query=query,
        parameters=[{"name": "@conversation_id", "value": conversation_id}],
        partition_key=conversation_id,
    ))
    if not items:
        return None
    return str(items[0].get("thread_id") or "").strip() or None


def _normalize_document_workspace_scope(workspace_scope: str = "personal") -> str:
    normalized_workspace_scope = str(workspace_scope or "personal").strip().lower()
    if normalized_workspace_scope not in {"personal", "group"}:
        raise ValueError("workspace_scope must be 'personal' or 'group'")
    return normalized_workspace_scope


def _normalize_markdown_file_name(file_name: str) -> str:
    normalized_file_name = str(file_name or "").replace("\\", "/").split("/")[-1].strip()
    if not normalized_file_name:
        normalized_file_name = "generated_markdown_document"

    base_name, extension = os.path.splitext(normalized_file_name)
    if extension.lower() == ".md" and base_name.strip():
        return normalized_file_name

    normalized_base_name = base_name.strip() or normalized_file_name.strip() or "generated_markdown_document"
    return f"{normalized_base_name}.md"


def _normalize_generated_document_file_name(file_name: str) -> str:
    normalized_file_name = str(file_name or "").replace("\\", "/").split("/")[-1].strip()
    if not normalized_file_name:
        normalized_file_name = "generated_tabular_output.json"

    base_name, extension = os.path.splitext(normalized_file_name)
    normalized_extension = extension.lower().strip()
    if normalized_extension and normalized_extension != "." and base_name.strip():
        return normalized_file_name

    normalized_base_name = base_name.strip() or normalized_file_name.strip() or "generated_tabular_output"
    return f"{normalized_base_name}.json"


def _clean_markdown_like_text(value: str) -> str:
    cleaned_text = str(value or "").strip()
    cleaned_text = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r"\1 (\2)", cleaned_text)
    cleaned_text = cleaned_text.replace("**", "").replace("__", "").replace("`", "")
    cleaned_text = cleaned_text.replace("*_", "").replace("_*", "")
    return cleaned_text.strip()


def _append_markdown_like_content_to_docx(document, markdown_content: str) -> None:
    in_code_block = False
    for raw_line in str(markdown_content or "").replace("\r\n", "\n").replace("\r", "\n").split("\n"):
        line = raw_line.rstrip()
        stripped_line = line.strip()
        if stripped_line.startswith("```"):
            in_code_block = not in_code_block
            continue
        if not stripped_line:
            document.add_paragraph("")
            continue
        if in_code_block:
            document.add_paragraph(line)
            continue

        heading_match = re.match(r"^(#{1,6})\s+(.+)$", stripped_line)
        if heading_match:
            heading_level = min(len(heading_match.group(1)) + 1, 4)
            document.add_heading(_clean_markdown_like_text(heading_match.group(2)), level=heading_level)
            continue

        bullet_match = re.match(r"^[-*+]\s+(.+)$", stripped_line)
        if bullet_match:
            document.add_paragraph(_clean_markdown_like_text(bullet_match.group(1)), style="List Bullet")
            continue

        number_match = re.match(r"^\d+[.)]\s+(.+)$", stripped_line)
        if number_match:
            document.add_paragraph(_clean_markdown_like_text(number_match.group(1)), style="List Number")
            continue

        document.add_paragraph(_clean_markdown_like_text(stripped_line))


def _split_markdown_like_content_for_slides(markdown_content: str, fallback_title: str) -> List[Dict[str, Any]]:
    sections = []
    current_section = {
        "title": fallback_title or "Overview",
        "bullets": [],
    }

    for raw_line in str(markdown_content or "").replace("\r\n", "\n").replace("\r", "\n").split("\n"):
        stripped_line = raw_line.strip()
        if not stripped_line or stripped_line.startswith("```"):
            continue

        heading_match = re.match(r"^(#{1,3})\s+(.+)$", stripped_line)
        if heading_match:
            if current_section["bullets"]:
                sections.append(current_section)
            current_section = {
                "title": _clean_markdown_like_text(heading_match.group(2)) or fallback_title or "Overview",
                "bullets": [],
            }
            continue

        bullet_match = re.match(r"^(?:[-*+]|\d+[.)])\s+(.+)$", stripped_line)
        bullet_text = bullet_match.group(1) if bullet_match else stripped_line
        cleaned_bullet = _clean_markdown_like_text(bullet_text)
        if cleaned_bullet:
            current_section["bullets"].append(cleaned_bullet)

    if current_section["bullets"] or not sections:
        sections.append(current_section)

    split_sections = []
    for section in sections:
        bullets = section.get("bullets") or ["No content recorded."]
        for start_index in range(0, len(bullets), 8):
            split_sections.append({
                "title": str(section.get("title") or fallback_title or "Overview")[:120],
                "bullets": [str(item or "")[:220] for item in bullets[start_index:start_index + 8]],
            })
            if len(split_sections) >= 20:
                return split_sections
    return split_sections


def _get_slide_body_shape(slide):
    # Optional PowerPoint rendering dependency used only while creating generated presentations.
    from pptx.util import Inches as PptxInches

    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.idx == 1:
            return placeholder
    return slide.shapes.add_textbox(PptxInches(0.8), PptxInches(1.6), PptxInches(11.7), PptxInches(5.2))


def _add_simple_content_slide(presentation, title: str, bullets: List[str]) -> None:
    slide_layout = presentation.slide_layouts[1] if len(presentation.slide_layouts) > 1 else presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = str(title or "Overview")[:120]

    body_shape = _get_slide_body_shape(slide)
    text_frame = body_shape.text_frame
    text_frame.clear()
    for index, bullet in enumerate(bullets or ["No content recorded."]):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = str(bullet or "")[:220]
        paragraph.level = 0


def _populate_simple_presentation(presentation, title: str, markdown_content: str) -> None:
    normalized_title = str(title or "Generated Presentation").strip() or "Generated Presentation"
    title_slide_layout = presentation.slide_layouts[0]
    title_slide = presentation.slides.add_slide(title_slide_layout)
    if title_slide.shapes.title:
        title_slide.shapes.title.text = normalized_title[:120]
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "Generated by SimpleChat"

    for section in _split_markdown_like_content_for_slides(markdown_content, normalized_title):
        _add_simple_content_slide(
            presentation,
            section.get("title") or normalized_title,
            section.get("bullets") or [],
        )


def _write_temp_markdown_file(markdown_content: str) -> str:
    sc_temp_files_dir = "/sc-temp-files" if os.path.exists("/sc-temp-files") else None
    with tempfile.NamedTemporaryFile(delete=False, suffix=".md", dir=sc_temp_files_dir) as temp_file:
        temp_file.write(str(markdown_content or "").encode("utf-8"))
        return temp_file.name


def _write_temp_generated_file(file_content_bytes: bytes, suffix: str) -> str:
    sc_temp_files_dir = "/sc-temp-files" if os.path.exists("/sc-temp-files") else None
    normalized_suffix = suffix if suffix.startswith('.') else f'.{suffix}' if suffix else '.json'
    with tempfile.NamedTemporaryFile(delete=False, suffix=normalized_suffix, dir=sc_temp_files_dir) as temp_file:
        temp_file.write(file_content_bytes)
        return temp_file.name


def _queue_document_upload_background_task(
    document_id: str,
    user_id: str,
    temp_file_path: str,
    original_filename: str,
    group_id: Optional[str] = None,
    public_workspace_id: Optional[str] = None,
) -> None:
    task_kwargs = {
        "document_id": document_id,
        "user_id": user_id,
        "temp_file_path": temp_file_path,
        "original_filename": original_filename,
    }
    if group_id:
        task_kwargs["group_id"] = group_id
    if public_workspace_id:
        task_kwargs["public_workspace_id"] = public_workspace_id

    if not has_app_context():
        raise RuntimeError("SimpleChat document uploads require an active app context")

    executor = current_app.extensions.get("executor")
    if executor and hasattr(executor, "submit_stored"):
        executor.submit_stored(
            document_id,
            process_document_upload_background,
            **task_kwargs,
        )
        return

    if executor and hasattr(executor, "submit"):
        executor.submit(process_document_upload_background, **task_kwargs)
        return

    process_document_upload_background(**task_kwargs)


def queue_generated_document_processing(
    document_id: str,
    owner_user_id: str,
    normalized_file_name: str,
    file_content_bytes: Any,
    group_id: Optional[str] = None,
    public_workspace_id: Optional[str] = None,
    process_inline: bool = False,
) -> None:
    """Queue processing for an existing generated document shell."""
    normalized_document_id = str(document_id or "").strip()
    normalized_owner_user_id = str(owner_user_id or "").strip()
    normalized_name = str(normalized_file_name or "").strip()

    if not normalized_document_id:
        raise ValueError("document_id is required")
    if not normalized_owner_user_id:
        raise ValueError("owner_user_id is required")
    if not normalized_name:
        raise ValueError("normalized_file_name is required")

    if isinstance(file_content_bytes, bytes):
        normalized_file_content_bytes = file_content_bytes
    else:
        normalized_file_content_bytes = str(file_content_bytes or "").encode("utf-8")

    if not normalized_file_content_bytes.strip():
        raise ValueError("file_content_bytes is required")

    file_extension = os.path.splitext(normalized_name)[1].lower() or ".json"
    temp_file_path = _write_temp_generated_file(normalized_file_content_bytes, file_extension)

    try:
        if process_inline:
            process_document_upload_background(
                document_id=normalized_document_id,
                user_id=normalized_owner_user_id,
                temp_file_path=temp_file_path,
                original_filename=normalized_name,
                group_id=group_id,
                public_workspace_id=public_workspace_id,
            )
            return

        _queue_document_upload_background_task(
            document_id=normalized_document_id,
            user_id=normalized_owner_user_id,
            temp_file_path=temp_file_path,
            original_filename=normalized_name,
            group_id=group_id,
            public_workspace_id=public_workspace_id,
        )
    except Exception:
        if os.path.exists(temp_file_path):
            os.remove(temp_file_path)
        raise


def _upload_generated_document_for_current_user(
    current_user_id: str,
    normalized_file_name: str,
    file_content_bytes: bytes,
    normalized_workspace_scope: str,
    group_id: str = "",
    default_group_id: str = "",
    process_inline: bool = False,
) -> Dict[str, Any]:
    document_id = str(uuid.uuid4())
    file_extension = os.path.splitext(normalized_file_name)[1].lower() or '.json'
    temp_file_path = _write_temp_generated_file(file_content_bytes, file_extension)
    resolved_group_id = None
    initial_status = "Processing file..." if process_inline else "Queued for processing"

    try:
        if normalized_workspace_scope == "group":
            resolved_group_id = _resolve_group_upload_target_for_current_user(
                current_user_id,
                group_id=group_id,
                default_group_id=default_group_id,
            )
            create_document(
                file_name=normalized_file_name,
                group_id=resolved_group_id,
                user_id=current_user_id,
                document_id=document_id,
                num_file_chunks=0,
                status=initial_status,
            )
            update_document(
                document_id=document_id,
                user_id=current_user_id,
                group_id=resolved_group_id,
                percentage_complete=0,
            )
        else:
            create_document(
                file_name=normalized_file_name,
                user_id=current_user_id,
                document_id=document_id,
                num_file_chunks=0,
                status=initial_status,
            )
            update_document(
                document_id=document_id,
                user_id=current_user_id,
                percentage_complete=0,
            )

        if process_inline:
            process_document_upload_background(
                document_id=document_id,
                user_id=current_user_id,
                temp_file_path=temp_file_path,
                original_filename=normalized_file_name,
                group_id=resolved_group_id,
            )
        else:
            _queue_document_upload_background_task(
                document_id=document_id,
                user_id=current_user_id,
                temp_file_path=temp_file_path,
                original_filename=normalized_file_name,
                group_id=resolved_group_id,
            )

        if normalized_workspace_scope == "group":
            invalidate_group_search_cache(resolved_group_id)
            log_document_upload(
                user_id=current_user_id,
                container_type="group",
                document_id=document_id,
                file_size=len(file_content_bytes),
                file_type=file_extension,
            )
        else:
            invalidate_personal_search_cache(current_user_id)
            log_document_upload(
                user_id=current_user_id,
                container_type="personal",
                document_id=document_id,
                file_size=len(file_content_bytes),
                file_type=file_extension,
            )

        return {
            "document": {
                "id": document_id,
                "file_name": normalized_file_name,
                "status": "Processing complete" if process_inline else "Queued for processing",
            },
            "workspace_scope": normalized_workspace_scope,
            "group_id": resolved_group_id,
        }
    except Exception:
        if os.path.exists(temp_file_path):
            os.remove(temp_file_path)
        raise


def _upload_generated_chat_artifact_for_current_user(
    current_user_id: str,
    conversation_id: str,
    normalized_file_name: str,
    file_content_bytes: bytes,
    artifact_metadata: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    try:
        conversation_item = cosmos_conversations_container.read_item(
            item=conversation_id,
            partition_key=conversation_id,
        )
    except CosmosResourceNotFoundError as exc:
        raise LookupError(f"Conversation {conversation_id} not found") from exc

    if str(conversation_item.get("user_id") or "").strip() != current_user_id:
        raise PermissionError("Forbidden")

    blob_service_client = CLIENTS.get("storage_account_office_docs_client")
    if not blob_service_client:
        raise RuntimeError("Blob storage client not available")

    artifact_message_id = f"{conversation_id}_generated_file_{uuid.uuid4().hex}"
    blob_path = (
        f"{current_user_id}/{conversation_id}/generated/"
        f"{artifact_message_id}/{normalized_file_name}"
    )
    blob_client = blob_service_client.get_blob_client(
        container=storage_account_personal_chat_container_name,
        blob=blob_path,
    )
    blob_client.upload_blob(
        file_content_bytes,
        overwrite=True,
        metadata={
            "conversation_id": conversation_id,
            "user_id": current_user_id,
            "generated_artifact": "true",
        },
    )

    timestamp = datetime.now(timezone.utc).isoformat()
    current_thread_id = str(uuid.uuid4())
    previous_thread_id = _get_latest_personal_thread_id(conversation_id)
    file_extension = os.path.splitext(normalized_file_name)[1].lower().lstrip(".")
    artifact_metadata = artifact_metadata if isinstance(artifact_metadata, dict) else {}
    artifact_capability = str(artifact_metadata.get("capability") or "analysis").strip().lower() or "analysis"
    artifact_output_format = str(artifact_metadata.get("output_format") or file_extension).strip().lower() or file_extension
    artifact_summary = str(artifact_metadata.get("summary") or "").strip()

    message_doc = {
        "id": artifact_message_id,
        "conversation_id": conversation_id,
        "role": "file",
        "filename": normalized_file_name,
        "is_table": file_extension in TABULAR_EXTENSIONS,
        "file_content_source": "blob",
        "blob_container": storage_account_personal_chat_container_name,
        "blob_path": blob_path,
        "timestamp": timestamp,
        "model_deployment_name": None,
        "metadata": {
            "is_generated_chat_artifact": True,
            "generated_artifact_storage_scope": "chat",
            "generated_artifact_capability": artifact_capability,
            "generated_artifact_output_format": artifact_output_format,
            "generated_artifact_summary": artifact_summary,
            "thread_info": {
                "thread_id": current_thread_id,
                "previous_thread_id": previous_thread_id,
                "active_thread": False,
                "thread_attempt": 1,
            },
        },
    }
    cosmos_messages_container.upsert_item(message_doc)

    log_event(
        "[SimpleChat] Generated chat artifact saved",
        {
            "conversation_id": conversation_id,
            "message_id": artifact_message_id,
            "file_name": normalized_file_name,
            "blob_path": blob_path,
            "storage_scope": "chat",
            "capability": artifact_capability,
            "output_format": artifact_output_format,
        },
        debug_only=True,
    )

    return {
        "message": {
            "id": artifact_message_id,
            "file_name": normalized_file_name,
            "blob_container": storage_account_personal_chat_container_name,
            "blob_path": blob_path,
            "capability": artifact_capability,
            "output_format": artifact_output_format,
        },
        "conversation_id": conversation_id,
    }


def _resolve_group_upload_target_for_current_user(
    current_user_id: str,
    group_id: str = "",
    default_group_id: str = "",
) -> str:
    normalized_group_id = str(group_id or default_group_id or "").strip()
    if not normalized_group_id:
        normalized_group_id = require_active_group(current_user_id)

    group_doc = find_group_by_id(normalized_group_id)
    if not group_doc:
        raise LookupError("Group not found")

    allowed, reason = check_group_status_allows_operation(group_doc, "upload")
    if not allowed:
        raise PermissionError(reason)

    assert_group_role(
        current_user_id,
        normalized_group_id,
        allowed_roles=("Owner", "Admin", "DocumentManager"),
    )
    return normalized_group_id


def _derive_personal_conversation_title(content: str) -> str:
    return derive_conversation_title_from_message(content)


def _resolve_group_doc_for_current_user(
    current_user_id: str,
    group_id: str = "",
    default_group_id: str = "",
    allowed_roles: Tuple[str, ...] = ("Owner", "Admin", "DocumentManager", "User"),
    missing_group_message: str = "group_id is required",
) -> Dict[str, Any]:
    _require_group_workspaces_enabled()
    resolved_group_id = str(group_id or default_group_id or "").strip()
    if not resolved_group_id:
        resolved_group_id = require_active_group(
            current_user_id,
            allowed_roles=allowed_roles,
        )

    if not resolved_group_id:
        raise ValueError(missing_group_message)

    group_doc = find_group_by_id(resolved_group_id)
    if not group_doc:
        raise LookupError("Group not found")

    assert_group_role(current_user_id, resolved_group_id, allowed_roles=allowed_roles)
    return group_doc


def _require_current_user_info() -> Dict[str, str]:
    current_user = get_current_user_info()
    if not current_user or not current_user.get("userId"):
        raise PermissionError("User not authenticated")
    return current_user


def _require_group_workspaces_enabled() -> Dict[str, Any]:
    settings = get_settings() or {}
    if not settings.get("enable_group_workspaces", False):
        raise PermissionError("Group workspaces are disabled by configuration")
    return settings


def _require_group_creation_enabled(settings: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    settings = settings or get_settings() or {}
    if not settings.get("enable_group_creation", False):
        raise PermissionError("Group creation is disabled by configuration")

    if settings.get("require_member_of_create_group", False):
        user_roles = (session.get("user") or {}).get("roles") or []
        if "CreateGroups" not in user_roles:
            raise PermissionError("Insufficient permissions (CreateGroups role required)")
    return settings


def _require_user_workflows_enabled() -> Dict[str, Any]:
    settings = get_settings() or {}
    user_roles = (session.get("user") or {}).get("roles") or []
    if not is_user_workflows_enabled_for_user(settings, user_roles=user_roles):
        if settings.get("allow_user_workflows", False):
            raise PermissionError("Insufficient permissions (WorkflowUser role required)")
        raise PermissionError("Personal workflows are disabled by configuration")
    return settings


def _require_control_center_admin_access(settings: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    settings = settings or get_settings() or {}
    session_user = (session.get("user") or {})
    user_roles = session_user.get("roles") or []
    require_member_of_control_center_admin = settings.get("require_member_of_control_center_admin", False)

    has_control_center_admin_role = "ControlCenterAdmin" in user_roles
    has_regular_admin_role = "Admin" in user_roles

    if require_member_of_control_center_admin:
        if not has_control_center_admin_role:
            raise PermissionError("Insufficient permissions (ControlCenterAdmin role required)")
        return session_user

    if not has_regular_admin_role:
        raise PermissionError("Insufficient permissions (Admin role required)")
    return session_user


def _require_collaboration_feature_enabled() -> Dict[str, Any]:
    settings = get_settings() or {}
    if not settings.get("enable_collaborative_conversations", False):
        raise PermissionError("Collaborative conversations are disabled by configuration")
    return settings


def _graph_get_json(path: str, params: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    token = get_valid_access_token()
    if not token:
        raise PermissionError("Could not acquire access token")

    response = requests.get(
        get_graph_endpoint(path),
        headers={
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json",
        },
        params=params,
        timeout=20,
    )
    response.raise_for_status()
    return response.json()


def _get_directory_user_by_id(user_id: str) -> Optional[Dict[str, str]]:
    normalized_user_id = str(user_id or "").strip()
    if not normalized_user_id:
        return None

    token = get_valid_access_token()
    if not token:
        raise PermissionError("Could not acquire access token")

    response = requests.get(
        get_graph_endpoint(f"/users/{quote(normalized_user_id)}"),
        headers={
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json",
        },
        params={"$select": "id,displayName,mail,userPrincipalName"},
        timeout=20,
    )
    if response.status_code == 404:
        return None
    response.raise_for_status()
    return _normalize_directory_user(response.json())


def _find_directory_users_by_email(email: str) -> List[Dict[str, str]]:
    normalized_email = str(email or "").strip()
    if not normalized_email:
        return []

    escaped_email = _escape_odata_value(normalized_email)
    payload = _graph_get_json(
        "/users",
        params={
            "$filter": f"mail eq '{escaped_email}' or userPrincipalName eq '{escaped_email}'",
            "$top": 5,
            "$select": "id,displayName,mail,userPrincipalName",
        },
    )
    return _normalize_directory_users(payload.get("value", []))


def _normalize_directory_users(raw_users: Iterable[Dict[str, Any]]) -> List[Dict[str, str]]:
    normalized_users = []
    for raw_user in raw_users or []:
        normalized_user = _normalize_directory_user(raw_user)
        if normalized_user:
            normalized_users.append(normalized_user)
    return normalized_users


def _normalize_directory_user(raw_user: Dict[str, Any]) -> Optional[Dict[str, str]]:
    if not isinstance(raw_user, dict):
        return None

    user_id = str(raw_user.get("id") or "").strip()
    if not user_id:
        return None

    email = str(raw_user.get("mail") or raw_user.get("userPrincipalName") or "").strip()
    display_name = str(raw_user.get("displayName") or email or user_id).strip()
    return {
        "id": user_id,
        "displayName": display_name,
        "email": email,
    }


def _escape_odata_value(value: str) -> str:
    return str(value or "").replace("'", "''").strip()


def _log_group_member_addition(
    actor_user: Dict[str, str],
    actor_role: str,
    group_doc: Dict[str, Any],
    member_doc: Dict[str, str],
    member_role: str,
) -> None:
    try:
        activity_record = {
            "id": str(uuid.uuid4()),
            "activity_type": "add_member_directly",
            "timestamp": datetime.utcnow().isoformat(),
            "added_by_user_id": actor_user.get("userId"),
            "added_by_email": actor_user.get("email", "unknown"),
            "added_by_role": actor_role,
            "group_id": group_doc.get("id"),
            "group_name": group_doc.get("name", "Unknown"),
            "member_user_id": member_doc.get("userId", ""),
            "member_email": member_doc.get("email", ""),
            "member_name": member_doc.get("displayName", ""),
            "member_role": member_role,
            "description": (
                f"{actor_role} {actor_user.get('email', 'unknown')} added member "
                f"{member_doc.get('displayName', '')} ({member_doc.get('email', '')}) to group "
                f"{group_doc.get('name', group_doc.get('id', 'Unknown'))} as {member_role}"
            ),
        }
        cosmos_activity_logs_container.create_item(body=activity_record)
    except Exception as exc:
        log_event(
            f"[SimpleChat] Failed to log group member addition: {exc}",
            level=logging.WARNING,
            exceptionTraceback=True,
        )


def _notify_group_member_addition(
    group_doc: Dict[str, Any],
    member_doc: Dict[str, str],
    member_role: str,
    added_by_email: str,
    actor_user: Optional[Dict[str, str]] = None,
) -> None:
    role_display = {
        "admin": "Admin",
        "document_manager": "Document Manager",
        "user": "Member",
    }.get(member_role, "Member")

    try:
        create_notification(
            user_id=member_doc.get("userId", ""),
            notification_type="group_member_added",
            title="Added to Group",
            message=(
                f"You have been added to the group '{group_doc.get('name', 'Unknown')}' "
                f"as {role_display} by {added_by_email}."
            ),
            link_url=f"/manage_group/{group_doc.get('id', '')}",
            link_context={
                "workspace_type": "group",
                "group_id": group_doc.get("id", ""),
            },
            metadata={
                "group_id": group_doc.get("id", ""),
                "group_name": group_doc.get("name", "Unknown"),
                "added_by": added_by_email,
                "role": member_role,
                "audience": "member",
            },
        )
    except Exception as exc:
        log_event(
            f"[SimpleChat] Failed to notify group member addition: {exc}",
            level=logging.WARNING,
            exceptionTraceback=True,
        )

    actor_user_id = str((actor_user or {}).get("userId") or "").strip()
    if not actor_user_id or actor_user_id == str(member_doc.get("userId") or "").strip():
        return

    try:
        create_notification(
            user_id=actor_user_id,
            notification_type="group_member_added",
            title="Group member added",
            message=(
                f"Added {member_doc.get('displayName', 'a new member')} to '{group_doc.get('name', 'Unknown')}' "
                f"as {role_display}."
            ),
            link_url=f"/manage_group/{group_doc.get('id', '')}",
            link_context={
                "workspace_type": "group",
                "group_id": group_doc.get("id", ""),
            },
            metadata={
                "group_id": group_doc.get("id", ""),
                "group_name": group_doc.get("name", "Unknown"),
                "member_user_id": member_doc.get("userId", ""),
                "member_email": member_doc.get("email", ""),
                "member_display_name": member_doc.get("displayName", ""),
                "role": member_role,
                "audience": "actor",
            },
        )
    except Exception as exc:
        log_event(
            f"[SimpleChat] Failed to notify actor about group member addition: {exc}",
            level=logging.WARNING,
            exceptionTraceback=True,
        )


def _create_personal_notification(
    user_id: str,
    notification_type: str,
    title: str,
    message: str,
    link_url: str = "",
    link_context: Optional[Dict[str, Any]] = None,
    metadata: Optional[Dict[str, Any]] = None,
) -> Optional[Dict[str, Any]]:
    normalized_user_id = str(user_id or "").strip()
    if not normalized_user_id:
        return None

    try:
        return create_notification(
            user_id=normalized_user_id,
            notification_type=notification_type,
            title=title,
            message=message,
            link_url=link_url,
            link_context=link_context or {},
            metadata=metadata or {},
        )
    except Exception as exc:
        log_event(
            f"[SimpleChat] Failed to create notification '{notification_type}': {exc}",
            level=logging.WARNING,
            exceptionTraceback=True,
        )
        return None


def _build_group_link_context(group_doc: Dict[str, Any]) -> Dict[str, Any]:
    return {
        "workspace_type": "group",
        "group_id": str((group_doc or {}).get("id") or "").strip(),
    }


def _build_conversation_link_context(conversation_doc: Dict[str, Any]) -> Dict[str, Any]:
    conversation_doc = conversation_doc if isinstance(conversation_doc, dict) else {}
    scope = conversation_doc.get("scope") if isinstance(conversation_doc.get("scope"), dict) else {}
    group_id = str(scope.get("group_id") or conversation_doc.get("group_id") or "").strip()
    chat_type = str(conversation_doc.get("chat_type") or "").strip().lower()

    link_context = {
        "conversation_id": str(conversation_doc.get("id") or "").strip(),
        "workspace_type": "group" if group_id or chat_type.startswith("group") else "personal",
    }
    if group_id:
        link_context["group_id"] = group_id
    if conversation_doc.get("conversation_kind"):
        link_context["conversation_kind"] = conversation_doc.get("conversation_kind")
    return link_context


def _build_conversation_link_url(conversation_doc: Dict[str, Any]) -> str:
    conversation_id = str((conversation_doc or {}).get("id") or "").strip()
    if not conversation_id:
        return ""
    return f"/chats?conversationId={conversation_id}"


def _get_group_notification_recipient_ids(group_doc: Dict[str, Any]) -> List[str]:
    recipient_ids = set()
    owner_user_id = str(((group_doc or {}).get("owner") or {}).get("id") or "").strip()
    if owner_user_id:
        recipient_ids.add(owner_user_id)

    for member in list((group_doc or {}).get("users", []) or []):
        member_user_id = str(member.get("userId") or "").strip()
        if member_user_id:
            recipient_ids.add(member_user_id)

    return sorted(recipient_ids)


def _notify_group_created(group_doc: Dict[str, Any], actor_user: Dict[str, str]) -> None:
    group_id = str((group_doc or {}).get("id") or "").strip()
    group_name = str((group_doc or {}).get("name") or "Untitled Group").strip() or "Untitled Group"
    actor_user_id = str((actor_user or {}).get("userId") or "").strip()
    if not group_id or not actor_user_id:
        return

    _create_personal_notification(
        user_id=actor_user_id,
        notification_type="group_created",
        title=f"Group created: {group_name}",
        message=f"You created the group '{group_name}'.",
        link_url=f"/manage_group/{group_id}",
        link_context=_build_group_link_context(group_doc),
        metadata={
            "group_id": group_id,
            "group_name": group_name,
        },
    )


def _notify_personal_conversation_created(
    conversation_item: Dict[str, Any],
    current_user: Dict[str, str],
) -> None:
    conversation_title = str((conversation_item or {}).get("title") or "New Conversation").strip() or "New Conversation"
    _create_personal_notification(
        user_id=str((current_user or {}).get("userId") or "").strip(),
        notification_type="conversation_created",
        title=f"Conversation created: {conversation_title}",
        message=f"Created a new personal conversation named '{conversation_title}'.",
        link_url=_build_conversation_link_url(conversation_item),
        link_context=_build_conversation_link_context(conversation_item),
        metadata={
            "conversation_id": str((conversation_item or {}).get("id") or "").strip(),
            "conversation_title": conversation_title,
            "chat_type": str((conversation_item or {}).get("chat_type") or "").strip(),
            "audience": "actor",
        },
    )


def _notify_group_conversation_created(
    conversation_doc: Dict[str, Any],
    group_doc: Dict[str, Any],
    creator_user: Dict[str, str],
) -> None:
    conversation_title = str((conversation_doc or {}).get("title") or "New group conversation").strip() or "New group conversation"
    group_name = str((group_doc or {}).get("name") or "Group Workspace").strip() or "Group Workspace"
    creator_display_name = str((creator_user or {}).get("display_name") or (creator_user or {}).get("displayName") or "A teammate").strip() or "A teammate"
    link_url = _build_conversation_link_url(conversation_doc)
    link_context = _build_conversation_link_context(conversation_doc)
    metadata = {
        "group_id": str((group_doc or {}).get("id") or "").strip(),
        "group_name": group_name,
        "conversation_id": str((conversation_doc or {}).get("id") or "").strip(),
        "conversation_title": conversation_title,
        "chat_type": str((conversation_doc or {}).get("chat_type") or "").strip(),
    }

    for recipient_user_id in _get_group_notification_recipient_ids(group_doc):
        audience = "actor" if recipient_user_id == str((creator_user or {}).get("user_id") or "").strip() else "member"
        if audience == "actor":
            title = f"Group conversation created: {conversation_title}"
            message = f"You created '{conversation_title}' in '{group_name}'."
        else:
            title = f"New group conversation in {group_name}"
            message = f"{creator_display_name} created '{conversation_title}' in '{group_name}'."

        _create_personal_notification(
            user_id=recipient_user_id,
            notification_type="conversation_created",
            title=title,
            message=message,
            link_url=link_url,
            link_context=link_context,
            metadata={
                **metadata,
                "audience": audience,
            },
        )


def _notify_personal_collaboration_conversation_created(
    conversation_doc: Dict[str, Any],
    creator_user: Dict[str, str],
    invited_participants: Optional[Iterable[Dict[str, Any]]] = None,
) -> None:
    conversation_title = str((conversation_doc or {}).get("title") or "Collaborative conversation").strip() or "Collaborative conversation"
    creator_user_id = str((creator_user or {}).get("user_id") or "").strip()
    creator_display_name = str((creator_user or {}).get("display_name") or "You").strip() or "You"
    link_url = _build_conversation_link_url(conversation_doc)
    link_context = _build_conversation_link_context(conversation_doc)
    base_metadata = {
        "conversation_id": str((conversation_doc or {}).get("id") or "").strip(),
        "conversation_title": conversation_title,
        "chat_type": str((conversation_doc or {}).get("chat_type") or "").strip(),
        "participant_count": len(list(invited_participants or [])) + 1,
    }

    _create_personal_notification(
        user_id=creator_user_id,
        notification_type="conversation_created",
        title=f"Collaborative conversation created: {conversation_title}",
        message=(
            f"Created '{conversation_title}'"
            f" with {len(list(invited_participants or []))} invited participant(s)."
        ),
        link_url=link_url,
        link_context=link_context,
        metadata={
            **base_metadata,
            "audience": "actor",
        },
    )

    for participant in invited_participants or []:
        participant_user_id = str((participant or {}).get("user_id") or "").strip()
        if not participant_user_id or participant_user_id == creator_user_id:
            continue

        _create_personal_notification(
            user_id=participant_user_id,
            notification_type="conversation_created",
            title=f"Added to collaborative conversation: {conversation_title}",
            message=f"{creator_display_name} added you to '{conversation_title}'.",
            link_url=link_url,
            link_context=link_context,
            metadata={
                **base_metadata,
                "created_by_user_id": creator_user_id,
                "created_by_display_name": creator_display_name,
                "audience": "participant",
            },
        )